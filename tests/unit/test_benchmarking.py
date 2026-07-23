import hashlib
import json
import tempfile
import unittest
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from rfs.evaluation.benchmarking import (
    _score_planning_contract,
    list_benchmark_cases,
    run_fast_benchmark_case,
    run_fast_benchmark_suite,
    score_benchmark_case,
    validate_benchmark_case,
)
from rfs.evaluation.pdf_extraction_benchmark import run_pdf_extraction_stress_suite
from rfs.evaluation.scanned_pdf import rasterize_pdf_as_scan
from rfs.paper_to_image.analyzer import parse_paper


ROOT = Path(__file__).resolve().parents[2]


class BenchmarkingTests(unittest.TestCase):
    def test_rasterized_scan_fixture_preserves_pages_and_removes_native_text(self):
        import fitz

        with tempfile.TemporaryDirectory() as temp:
            source = Path(temp) / "source.pdf"
            target = Path(temp) / "scanned.pdf"
            second_target = Path(temp) / "scanned_again.pdf"
            document = fitz.open()
            for text in ("Abstract native text", "Method native text"):
                page = document.new_page(width=600, height=800)
                page.insert_text((50, 80), text, fontsize=14)
            document.save(source)
            document.close()

            rasterize_pdf_as_scan(source, target, dpi=110)
            rasterize_pdf_as_scan(source, second_target, dpi=110)
            parsed = parse_paper(target, ocr_engine="off")
            first_hash = hashlib.sha256(target.read_bytes()).hexdigest()
            second_hash = hashlib.sha256(second_target.read_bytes()).hexdigest()

        self.assertEqual(parsed["page_count"], 2)
        self.assertEqual(parsed["extraction_report"]["pdf_type"], "scanned")
        self.assertTrue(all(page["char_count"] == 0 for page in parsed["pages"]))
        self.assertEqual(first_hash, second_hash)
    def test_planning_contract_treats_relation_labels_as_intermediate_artifacts(self):
        expected = {
            "entities": [
                {"id": "generator", "label": "Generate"},
                {"id": "draft", "label": "Initial Output"},
                {"id": "critic", "label": "Feedback"},
            ],
            "relations": [
                {"source": "generator", "target": "draft"},
                {"source": "draft", "target": "critic"},
            ],
        }
        specification = {
            "modules": [
                {"id": "init", "name": "Generate"},
                {"id": "feedback", "name": "Feedback"},
            ],
            "relations": [
                {"source": "init", "target": "feedback", "label": "Initial Output"},
            ],
        }

        result = _score_planning_contract(expected, specification)

        self.assertEqual(result["entity_recall"], 1.0)
        self.assertEqual(result["relation_recall"], 1.0)

    def test_planning_contract_prefers_real_entity_with_variable_suffix_over_relation_label(self):
        expected = {
            "entities": [
                {"id": "refiner", "label": "Refine"},
                {"id": "result", "label": "Refined Output"},
            ],
            "relations": [{"source": "refiner", "target": "result"}],
        }
        specification = {
            "modules": [{"id": "refine", "name": "Refine"}],
            "outputs": [{"id": "final", "name": "Refined Output (yn)"}],
            "relations": [{"source": "refine", "target": "final", "label": "Refined Output"}],
        }

        result = _score_planning_contract(expected, specification)

        self.assertEqual(result["entity_mapping"]["result"], "final")
        self.assertEqual(result["relation_recall"], 1.0)

    def test_planning_contract_uses_relations_to_disambiguate_equal_input_labels(self):
        expected = {
            "entities": [
                {"id": "task_input", "label": "Input"},
                {"id": "generator", "label": "Generate"},
            ],
            "relations": [{"source": "task_input", "target": "generator"}],
        }
        specification = {
            "inputs": [
                {"id": "iteration_input", "name": "Input (y0)"},
                {"id": "task_input", "name": "Input"},
            ],
            "modules": [
                {"id": "model", "name": "Model"},
                {"id": "generate", "name": "Generate"},
            ],
            "relations": [
                {"source": "iteration_input", "target": "model"},
                {"source": "task_input", "target": "generate"},
            ],
        }

        result = _score_planning_contract(expected, specification)

        self.assertEqual(result["entity_mapping"]["task_input"], "task_input")
        self.assertEqual(result["relation_recall"], 1.0)

    def test_bundled_cases_validate_and_list(self):
        p2i = ROOT / "benchmarks" / "paper-to-image" / "cases" / "001_linear_pipeline"
        i2p = ROOT / "benchmarks" / "image-to-ppt" / "cases" / "001_three_stage_layout"

        self.assertTrue(validate_benchmark_case(p2i)["ok"])
        self.assertTrue(validate_benchmark_case(i2p)["ok"])
        listing = list_benchmark_cases(ROOT / "benchmarks")
        self.assertGreaterEqual(len(listing["cases"]), 7)

    def test_real_paper_case_is_valid_with_source_metadata(self):
        case_dir = ROOT / "benchmarks" / "paper-to-image" / "cases" / "101_vit_linear"
        validation = validate_benchmark_case(case_dir)

        self.assertTrue(validation["ok"])
        self.assertTrue((case_dir / "source.json").exists())

    def test_fast_benchmark_command_runs_without_image_generation(self):
        case_dir = ROOT / "benchmarks" / "paper-to-image" / "cases" / "001_linear_pipeline"
        with tempfile.TemporaryDirectory() as tmp:
            result = run_fast_benchmark_case(case_dir, tmp, planner_mode="heuristic", ocr_engine="off")

            self.assertTrue(result["ok"])
            self.assertTrue((Path(tmp) / "fast_benchmark_result.json").exists())
            self.assertFalse((Path(tmp) / "selected_image.png").exists())

    def test_fast_suite_aggregates_reliability_and_timing_metrics(self):
        with tempfile.TemporaryDirectory() as tmp:
            result = run_fast_benchmark_suite(
                ROOT / "benchmarks",
                tmp,
                case_ids=["001_linear_pipeline"],
                planner_mode="heuristic",
                ocr_engine="off",
            )

            self.assertTrue(result["ok"])
            self.assertEqual(result["aggregate"]["case_count"], 1)
            self.assertIn("mean_total_seconds", result["aggregate"])
            self.assertIn("provider_success_rate", result["aggregate"])
            self.assertIn("mean_evidence_page_coverage_ratio", result["aggregate"])
            self.assertIn("ocr_scheduled_page_total", result["aggregate"])
            self.assertIn("ocr_attempted_page_total", result["aggregate"])
            self.assertIn("ocr_incomplete_run_count", result["aggregate"])
            self.assertIn("max_detected_column_count", result["aggregate"])
            self.assertIn("mean_section_count", result["aggregate"])
            self.assertIn("typographic_heading_total", result["aggregate"])
            self.assertIn("merged_heading_line_total", result["aggregate"])
            self.assertIn("figure_caption_total", result["aggregate"])
            self.assertIn("max_ocr_worker_count", result["aggregate"])
            self.assertIn("repeated_margin_noise_removed_total", result["aggregate"])
            self.assertIn("native_hyphenation_repair_total", result["aggregate"])
            self.assertIn("ocr_spacing_repair_total", result["aggregate"])
            self.assertTrue((Path(tmp) / "fast_suite_report.json").exists())

    def test_pdf_stress_suite_reports_ocr_stage_timings(self):
        with tempfile.TemporaryDirectory() as temp:
            result = run_pdf_extraction_stress_suite(Path(temp) / "pdf", ocr_engine="off")

        self.assertIn("ocr_stage_seconds", result["aggregate"])
        self.assertIn("recognition_seconds", result["aggregate"]["ocr_stage_seconds"])
        adapter_case = next(item for item in result["cases"] if item["case_id"] == "mixed_scan_fixture_adapter")
        self.assertIn("ocr_page_durations", adapter_case)
        self.assertIn("render_seconds", adapter_case["ocr_stage_seconds"])

    def test_pdf_extraction_stress_suite_runs_deterministically_without_runtime_ocr(self):
        with tempfile.TemporaryDirectory() as tmp:
            result = run_pdf_extraction_stress_suite(tmp, ocr_engine="off")

            self.assertTrue(result["ok"])
            self.assertEqual(result["aggregate"]["case_count"], 8)
            self.assertEqual(result["aggregate"]["passed_case_count"], 8)
            self.assertTrue((Path(tmp) / "fixtures" / "native_two_column.pdf").exists())
            self.assertTrue((Path(tmp) / "fixtures" / "repeated_margin_noise.pdf").exists())
            self.assertTrue((Path(tmp) / "fixtures" / "native_chinese.pdf").exists())
            self.assertTrue((Path(tmp) / "fixtures" / "hyphenated_native.pdf").exists())
            self.assertTrue((Path(tmp) / "fixtures" / "rotated_repeated_margin.pdf").exists())
            self.assertTrue((Path(tmp) / "mixed_scan_fixture_adapter" / "document_model.json").exists())
            self.assertTrue((Path(tmp) / "pdf_extraction_stress_report.json").exists())

    def test_paper_to_image_score_uses_hard_scientific_gates(self):
        case_dir = ROOT / "benchmarks" / "paper-to-image" / "cases" / "001_linear_pipeline"
        with tempfile.TemporaryDirectory() as tmp:
            run = Path(tmp)
            review = {
                "scientific": {"score": 1.0, "missing_modules": [], "missing_relations": [], "reversed_relations": [], "invented_items": []},
                "ocr": {"score": 1.0, "missing_labels": [], "misspelled_labels": [], "forbidden_labels_found": []},
                "aesthetic": {"score": 0.9, "readability": 0.9},
                "clarity": {"score": 0.9},
                "information": {"score": 0.95},
                "stability": {"production_pass_rate": 1.0},
                "hard_errors": [],
            }
            (run / "candidate_review.json").write_text(json.dumps({
                "selected_candidate_id": "candidate_01",
                "candidates": [{"candidate_id": "candidate_01", "score": 0.95, "review": review}],
            }), encoding="utf-8")

            result = score_benchmark_case(case_dir, run)

            self.assertTrue(result["passed"])
            self.assertEqual(result["metrics"]["relation_recall"], 1.0)

    def test_paper_to_image_score_checks_planning_contract(self):
        case_dir = ROOT / "benchmarks" / "paper-to-image" / "cases" / "001_linear_pipeline"
        with tempfile.TemporaryDirectory() as tmp:
            run = Path(tmp)
            review = {
                "scientific": {"score": 1.0, "missing_modules": [], "missing_relations": [], "reversed_relations": [], "invented_items": []},
                "ocr": {"score": 1.0, "missing_labels": [], "misspelled_labels": [], "forbidden_labels_found": []},
                "aesthetic": {"score": 0.9, "readability": 0.9},
                "clarity": {"score": 0.9},
                "information": {"score": 0.95},
                "stability": {"production_pass_rate": 1.0},
                "hard_errors": [],
            }
            (run / "candidate_review.json").write_text(json.dumps({
                "selected_candidate_id": "candidate_01",
                "candidates": [{"candidate_id": "candidate_01", "score": 0.95, "review": review}],
            }), encoding="utf-8")
            (run / "figure_specification.json").write_text(json.dumps({
                "modules": [{"id": "only_module", "name": "Feature Encoder"}],
                "relations": [],
            }), encoding="utf-8")

            result = score_benchmark_case(case_dir, run)

            self.assertFalse(result["passed"])
            self.assertLess(result["metrics"]["plan_entity_recall"], 1.0)
            self.assertEqual(result["metrics"]["plan_relation_recall"], 0.0)
            self.assertIn("paper planning missed required benchmark entities", result["hard_failures"])
            self.assertIn("paper planning missed required benchmark relations", result["hard_failures"])

    def test_full_slide_picture_is_detected_as_editability_cheating(self):
        case_dir = ROOT / "benchmarks" / "image-to-ppt" / "cases" / "001_three_stage_layout"
        with tempfile.TemporaryDirectory() as tmp:
            run = Path(tmp)
            reference = case_dir / "reference.ppm"
            with Image.open(reference) as image:
                image.save(run / "rebuild_preview.png")
                image.save(run / "reference.png")
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(str(run / "reference.png"), 0, 0, width=prs.slide_width, height=prs.slide_height)
            prs.save(run / "editable_composition.pptx")
            (run / "composition_quality_report.json").write_text(json.dumps({"slots": [], "cards": [], "arrows": []}), encoding="utf-8")
            (run / "rebuild_visual_quality_report.json").write_text(json.dumps({"blocking_issue_count": 0}), encoding="utf-8")

            result = score_benchmark_case(case_dir, run)

            self.assertFalse(result["passed"])
            self.assertEqual(result["metrics"]["full_slide_image_count"], 1)
            self.assertIn("full-slide reference image detected", result["hard_failures"])


if __name__ == "__main__":
    unittest.main()
