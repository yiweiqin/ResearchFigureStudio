from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
REFERENCE = Path(r"C:\Users\zhang\Documents\xwechat_files\wxid_h824xk1qpfoh22_adc2\temp\RWTemp\2026-06\860847045f7ccc2e7052eb182bab099b\c60bea4f7eb67de472ee6cf6bcbd09d6.png")
OUT = ROOT / "output" / "c60_editable_rebuild"

REF_W = 1806
REF_H = 920
SLIDE_W = 15.60
SLIDE_H = SLIDE_W * REF_H / REF_W

FONT_REGULAR = "Arial"
FONT_UI = "Microsoft YaHei UI"
TEXT_STYLES = {
    "panel_title": {"font": FONT_REGULAR, "size": 11.2, "bold": True},
    "panel_title_small": {"font": FONT_REGULAR, "size": 10.3, "bold": True},
    "module_label": {"font": FONT_REGULAR, "size": 7.7, "bold": True},
    "module_output": {"font": FONT_REGULAR, "size": 7.3, "bold": False},
    "body": {"font": FONT_REGULAR, "size": 8.7, "bold": False},
    "body_small": {"font": FONT_REGULAR, "size": 7.7, "bold": False},
    "legend": {"font": FONT_REGULAR, "size": 8.8, "bold": True},
    "trait": {"font": FONT_REGULAR, "size": 6.7, "bold": False},
    "badge": {"font": FONT_REGULAR, "size": 10.5, "bold": True},
}


def x(px: float) -> float:
    return px / REF_W * SLIDE_W


def y(px: float) -> float:
    return px / REF_H * SLIDE_H


def w(px: float) -> float:
    return px / REF_W * SLIDE_W


def h(px: float) -> float:
    return px / REF_H * SLIDE_H


def rgb(hex_color: str) -> RGBColor:
    value = hex_color.strip().lstrip("#")
    return RGBColor(int(value[:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def set_line(shape, color: str, width: float = 1.4, dash: bool = False) -> None:
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width)
    if dash:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE

        shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH


def set_fill(shape, color: str, transparency: int = 0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.fill.transparency = transparency


def no_fill(shape) -> None:
    shape.fill.background()


def add_text(slide, text: str, left: float, top: float, width: float, height: float, *, size: float | None = None, color: str = "#111111", bold: bool | None = None, align=PP_ALIGN.CENTER, font: str | None = None, valign=MSO_ANCHOR.MIDDLE, style: str | None = None):
    spec = TEXT_STYLES.get(style or "", {})
    font = font or str(spec.get("font") or FONT_REGULAR)
    size = float(size if size is not None else spec.get("size", 11))
    bold = bool(bold if bold is not None else spec.get("bold", False))
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Inches(0.01)
    tf.margin_right = Inches(0.01)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    for idx, raw_line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = raw_line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font
        run.font.color.rgb = rgb(color)
    return box


def add_panel(slide, left_px, top_px, width_px, height_px, title: str, color="#0B3A91", title_size=12.5):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x(left_px)), Inches(y(top_px)), Inches(w(width_px)), Inches(h(height_px)))
    set_fill(panel, "#FFFFFF", 0)
    set_line(panel, color, 1.55)
    panel.shadow.inherit = False
    add_text(slide, title, x(left_px + 10), y(top_px + 8), w(width_px - 20), y(58), size=title_size, color=color, style="panel_title")
    return panel


def add_card(slide, left_px, top_px, width_px, height_px, border: str, fill: str = "#FFFFFF", radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    card = slide.shapes.add_shape(shape_type, Inches(x(left_px)), Inches(y(top_px)), Inches(w(width_px)), Inches(h(height_px)))
    set_fill(card, fill, 0)
    set_line(card, border, 1.2)
    card.shadow.inherit = False
    return card


def apply_arrow(connector, size: str = "sm") -> None:
    ln = connector.line._get_or_add_ln()
    ln.append(parse_xml(f'<a:tailEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" type="triangle" w="{size}" len="{size}"/>'))


def add_arrow(slide, x1, y1, x2, y2, color="#000000", width_pt=1.3):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x(x1)), Inches(y(y1)), Inches(x(x2)), Inches(y(y2)))
    set_line(conn, color, width_pt)
    apply_arrow(conn)
    return conn


def add_line(slide, x1, y1, x2, y2, color="#000000", width_pt=1.2, dash=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x(x1)), Inches(y(y1)), Inches(x(x2)), Inches(y(y2)))
    set_line(conn, color, width_pt, dash=dash)
    return conn


def add_polyline(slide, points, color="#000000", width_pt=1.2, arrow=False):
    lines = []
    for idx, (p1, p2) in enumerate(zip(points[:-1], points[1:])):
        conn = add_line(slide, p1[0], p1[1], p2[0], p2[1], color=color, width_pt=width_pt)
        if arrow and idx == len(points) - 2:
            apply_arrow(conn)
        lines.append(conn)
    return lines


def icon_video(slide, left_px, top_px, width_px=76, height_px=62, color="#173B5D"):
    box = add_card(slide, left_px, top_px, width_px, height_px, "#0B2844", "#2B4C6D", radius=False)
    for i in range(5):
        for side_x in (left_px + 6, left_px + width_px - 14):
            sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x(side_x)), Inches(y(top_px + 8 + i * 10)), Inches(w(8)), Inches(h(5)))
            set_fill(sq, "#C8D8EA")
            set_line(sq, "#C8D8EA", 0.2)
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(x(left_px + width_px * 0.43)), Inches(y(top_px + height_px * 0.28)), Inches(w(width_px * 0.25)), Inches(h(height_px * 0.42)))
    tri.rotation = 90
    set_fill(tri, "#FFFFFF")
    set_line(tri, "#FFFFFF", 0.5)
    return box


def icon_wave(slide, left_px, top_px, width_px=82, height_px=42, color="#0A459C"):
    cx = left_px + width_px / 2
    for i in range(15):
        px = left_px + 4 + i * (width_px - 8) / 14
        amp = (0.25 + 0.75 * abs(7 - i) / 7) * height_px * 0.42
        add_line(slide, px, top_px + height_px / 2 - amp / 2, px, top_px + height_px / 2 + amp / 2, color=color, width_pt=1.25)
    add_line(slide, left_px, top_px + height_px / 2, left_px + 8, top_px + height_px / 2, color=color, width_pt=1)
    add_line(slide, cx + width_px / 2 - 8, top_px + height_px / 2, cx + width_px / 2, top_px + height_px / 2, color=color, width_pt=1)


def icon_doc(slide, left_px, top_px, width_px=44, height_px=58, color="#111111"):
    doc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x(left_px)), Inches(y(top_px)), Inches(w(width_px)), Inches(h(height_px)))
    set_fill(doc, "#FFFFFF")
    set_line(doc, color, 1.2)
    fold = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(x(left_px + width_px - 15)), Inches(y(top_px)), Inches(w(15)), Inches(h(15)))
    set_fill(fold, "#E9EEF3")
    set_line(fold, color, 0.8)
    for i in range(3):
        add_line(slide, left_px + 8, top_px + 20 + i * 11, left_px + width_px - 8, top_px + 20 + i * 11, color="#1A1A1A", width_pt=0.9)


def icon_mic(slide, left_px, top_px, color="#173B5D"):
    oval = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x(left_px + 16)), Inches(y(top_px)), Inches(w(22)), Inches(h(48)))
    set_fill(oval, color)
    set_line(oval, color, 1.0)
    add_line(slide, left_px + 7, top_px + 26, left_px + 7, top_px + 40, color=color, width_pt=2)
    add_line(slide, left_px + 47, top_px + 26, left_px + 47, top_px + 40, color=color, width_pt=2)
    add_line(slide, left_px + 8, top_px + 40, left_px + 27, top_px + 55, color=color, width_pt=2)
    add_line(slide, left_px + 46, top_px + 40, left_px + 27, top_px + 55, color=color, width_pt=2)
    add_line(slide, left_px + 27, top_px + 55, left_px + 27, top_px + 72, color=color, width_pt=2)
    add_line(slide, left_px + 12, top_px + 72, left_px + 42, top_px + 72, color=color, width_pt=2)


def icon_face(slide, left_px, top_px, size_px=62):
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x(left_px)), Inches(y(top_px)), Inches(w(size_px)), Inches(h(size_px)))
    no_fill(frame)
    set_line(frame, "#6B57C8", 1.0)
    head = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x(left_px + size_px * 0.28)), Inches(y(top_px + size_px * 0.18)), Inches(w(size_px * 0.44)), Inches(h(size_px * 0.48)))
    set_fill(head, "#D7B08E")
    set_line(head, "#6B4C35", 0.8)
    hair = slide.shapes.add_shape(MSO_SHAPE.ARC, Inches(x(left_px + size_px * 0.23)), Inches(y(top_px + size_px * 0.13)), Inches(w(size_px * 0.54)), Inches(h(size_px * 0.32)))
    no_fill(hair)
    set_line(hair, "#202020", 2.5)
    for ex in (left_px + size_px * 0.41, left_px + size_px * 0.58):
        eye = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x(ex)), Inches(y(top_px + size_px * 0.40)), Inches(w(3)), Inches(h(3)))
        set_fill(eye, "#111111")
        set_line(eye, "#111111", 0.2)
    body = slide.shapes.add_shape(MSO_SHAPE.ARC, Inches(x(left_px + size_px * 0.25)), Inches(y(top_px + size_px * 0.64)), Inches(w(size_px * 0.50)), Inches(h(size_px * 0.25)))
    no_fill(body)
    set_line(body, "#284A7C", 2.2)


def icon_person(slide, left_px, top_px, scale=1.0, shirt="#546E7A"):
    head = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x(left_px + 18 * scale)), Inches(y(top_px)), Inches(w(18 * scale)), Inches(h(18 * scale)))
    set_fill(head, "#D0A079")
    set_line(head, "#5A4030", 0.6)
    body = slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID, Inches(x(left_px + 12 * scale)), Inches(y(top_px + 20 * scale)), Inches(w(30 * scale)), Inches(h(48 * scale)))
    set_fill(body, shirt)
    set_line(body, "#333333", 0.6)
    add_line(slide, left_px + 17 * scale, top_px + 68 * scale, left_px + 10 * scale, top_px + 95 * scale, color="#222222", width_pt=1.4)
    add_line(slide, left_px + 35 * scale, top_px + 68 * scale, left_px + 43 * scale, top_px + 95 * scale, color="#222222", width_pt=1.4)


def icon_pose(slide, left_px, top_px, scale=1.0, color="#111111"):
    pts = {
        "head": (left_px + 34 * scale, top_px + 8 * scale),
        "neck": (left_px + 34 * scale, top_px + 25 * scale),
        "hip": (left_px + 34 * scale, top_px + 58 * scale),
        "lh": (left_px + 10 * scale, top_px + 42 * scale),
        "rh": (left_px + 58 * scale, top_px + 42 * scale),
        "lf": (left_px + 17 * scale, top_px + 88 * scale),
        "rf": (left_px + 52 * scale, top_px + 88 * scale),
    }
    head = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x(pts["head"][0] - 5 * scale)), Inches(y(pts["head"][1] - 5 * scale)), Inches(w(10 * scale)), Inches(h(10 * scale)))
    set_fill(head, color)
    set_line(head, color, 0.5)
    for a, b in [("neck", "hip"), ("neck", "lh"), ("neck", "rh"), ("hip", "lf"), ("hip", "rf")]:
        add_line(slide, pts[a][0], pts[a][1], pts[b][0], pts[b][1], color=color, width_pt=1.5)
    for key, c in [("lh", "#E17721"), ("rh", "#7A55C7"), ("lf", "#E17721"), ("rf", "#0A7F9C")]:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x(pts[key][0] - 3 * scale)), Inches(y(pts[key][1] - 3 * scale)), Inches(w(6 * scale)), Inches(h(6 * scale)))
        set_fill(dot, c)
        set_line(dot, c, 0.3)


def icon_frames(slide, left_px, top_px):
    for i in range(4):
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x(left_px + i * 7)), Inches(y(top_px + i * 5)), Inches(w(44)), Inches(h(54)))
        set_fill(rect, "#F6F6F6")
        set_line(rect, "#333333", 0.8)
    icon_person(slide, left_px + 28, top_px + 26, scale=0.33)


def icon_checklist(slide, left_px, top_px):
    board = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x(left_px)), Inches(y(top_px)), Inches(w(70)), Inches(h(96)))
    set_fill(board, "#EAF2FB")
    set_line(board, "#0B3A91", 1.2)
    clip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x(left_px + 22)), Inches(y(top_px - 8)), Inches(w(28)), Inches(h(16)))
    set_fill(clip, "#2F5A88")
    set_line(clip, "#2F5A88", 0.6)
    for i in range(4):
        cy = top_px + 24 + i * 16
        add_line(slide, left_px + 16, cy, left_px + 20, cy + 5, color="#173B5D", width_pt=1.5)
        add_line(slide, left_px + 20, cy + 5, left_px + 28, cy - 4, color="#173B5D", width_pt=1.5)
        add_line(slide, left_px + 36, cy, left_px + 58, cy, color="#111111", width_pt=1.0)


def icon_database_people(slide, left_px, top_px):
    top = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x(left_px)), Inches(y(top_px)), Inches(w(92)), Inches(h(28)))
    set_fill(top, "#008C91")
    set_line(top, "#00737B", 1.0)
    body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x(left_px)), Inches(y(top_px + 14)), Inches(w(92)), Inches(h(58)))
    set_fill(body, "#008C91")
    set_line(body, "#00737B", 1.0)
    bottom = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x(left_px)), Inches(y(top_px + 58)), Inches(w(92)), Inches(h(28)))
    set_fill(bottom, "#008C91")
    set_line(bottom, "#00737B", 1.0)
    for yy in (top_px + 28, top_px + 50):
        add_line(slide, left_px + 6, yy, left_px + 86, yy, color="#FFFFFF", width_pt=1.4)
    for i, (dx, dy, s) in enumerate([(80, 48, 1.0), (105, 50, 0.8), (66, 63, 0.8)]):
        head = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x(left_px + dx)), Inches(y(top_px + dy)), Inches(w(20 * s)), Inches(h(20 * s)))
        set_fill(head, "#0A6973")
        set_line(head, "#0A6973", 0.8)
        body2 = slide.shapes.add_shape(MSO_SHAPE.ARC, Inches(x(left_px + dx - 8 * s)), Inches(y(top_px + dy + 16 * s)), Inches(w(36 * s)), Inches(h(26 * s)))
        no_fill(body2)
        set_line(body2, "#0A6973", 4.0)


def icon_camera(slide, left_px, top_px):
    cam = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x(left_px)), Inches(y(top_px)), Inches(w(50)), Inches(h(32)))
    set_fill(cam, "#2B2B2B")
    set_line(cam, "#111111", 1)
    lens = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x(left_px + 13)), Inches(y(top_px + 6)), Inches(w(18)), Inches(h(18)))
    set_fill(lens, "#111111")
    set_line(lens, "#444444", 0.8)
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(x(left_px + 48)), Inches(y(top_px + 7)), Inches(w(26)), Inches(h(18)))
    tri.rotation = 90
    set_fill(tri, "#2B2B2B")
    set_line(tri, "#111111", 0.8)
    add_line(slide, left_px + 25, top_px + 32, left_px + 25, top_px + 78, color="#111111", width_pt=2)
    add_line(slide, left_px + 25, top_px + 50, left_px + 5, top_px + 96, color="#111111", width_pt=1.8)
    add_line(slide, left_px + 25, top_px + 50, left_px + 50, top_px + 96, color="#111111", width_pt=1.8)
    add_line(slide, left_px + 25, top_px + 50, left_px + 25, top_px + 96, color="#111111", width_pt=1.8)


def draw_rebuild() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb("#FFFFFF")

    blue = "#0B3A91"
    cyan = "#007E86"
    orange = "#C95100"
    purple = "#6B57C8"
    green = "#2D7B36"

    # Main panels.
    add_panel(slide, 10, 44, 245, 660, "1. Virtual Interview\nSetup", blue, 12)
    add_panel(slide, 303, 172, 213, 438, "2. Raw Video:\n287 participants\n× 36 questions", blue, 12)
    add_panel(slide, 585, 44, 392, 690, "3. Preprocessing Modules", blue, 12)
    add_panel(slide, 1022, 190, 155, 445, "4. Timestamp\nAlignment", blue, 11.5)
    add_panel(slide, 1206, 136, 192, 571, "5. Five Modalities", cyan, 11.5)
    add_panel(slide, 1430, 136, 188, 571, "6. NEO-FFI-3\nSelf-report Labels", orange, 11.2)
    add_panel(slide, 1654, 240, 158, 447, "7. Multimodal\nPersonality\nDataset", cyan, 11.2)

    # Panel 1.
    icon_person(slide, 18, 140, scale=1.25, shirt="#5D6C76")
    monitor = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x(88)), Inches(y(128)), Inches(w(150)), Inches(h(96)))
    set_fill(monitor, "#9EAAB5")
    set_line(monitor, "#111111", 1.5)
    add_text(slide, "00:15", x(94), y(132), w(42), h(18), size=5.5, color="#FFFFFF", bold=True)
    icon_face(slide, 151, 140, size_px=42)
    add_text(slide, "Participant in front of\n49-inch screen", x(47), y(292), w(165), h(42), style="body")
    screen = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x(31)), Inches(y(352)), Inches(w(197)), Inches(h(104)))
    set_fill(screen, "#9EAAB5")
    set_line(screen, "#111111", 1.4)
    icon_face(slide, 105, 368, size_px=58)
    add_text(slide, "3D virtual interviewer\non screen", x(46), y(465), w(150), h(42), style="body")
    icon_camera(slide, 31, 530)
    add_line(slide, 83, 552, 207, 520, color=blue, width_pt=0.9, dash=True)
    add_line(slide, 83, 552, 208, 630, color=blue, width_pt=0.9, dash=True)
    icon_person(slide, 205, 518, scale=0.55, shirt="#6E7F87")
    add_text(slide, "Wide-angle camera\ncapturing full body", x(44), y(648), w(165), h(45), style="body")

    # Panel 2 raw video.
    icon_video(slide, 335, 300, 152, 110)
    for row in range(2):
        for col in range(3 if row == 0 else 2):
            px = 323 + col * 58
            py = 440 + row * 78
            thumb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x(px)), Inches(y(py)), Inches(w(52)), Inches(h(62)))
            set_fill(thumb, "#D3D0C8")
            set_line(thumb, "#555555", 0.6)
            icon_person(slide, px + 16, py + 8, scale=0.36, shirt="#4C6C78")
    add_text(slide, "...", x(447), y(520), w(42), h(30), size=17, bold=True)

    # Panel 3 preprocessing cards.
    cards = [
        (596, 98, 372, 102, "#2E76CC", "#F8FBFF", "FFmpeg", "Audio", blue),
        (596, 224, 372, 102, "#2D8A45", "#F7FFF7", "FunASR", "Spoken Text\n+ Timestamps", green),
        (596, 352, 372, 102, purple, "#FCFAFF", "MTCNN", "Face Clips", purple),
        (596, 478, 372, 102, cyan, "#F7FFFF", "AlphaPose", "Full-body\nPose\nSkeletons", cyan),
        (596, 604, 372, 102, orange, "#FFF8F2", "Frame\nSampling", "Video\nFrames", orange),
    ]
    for idx, (lx, ty, ww, hh, border, fill, label, out_label, label_color) in enumerate(cards):
        add_card(slide, lx, ty, ww, hh, border, fill)
        cy = ty + hh / 2
        if idx == 0:
            icon_video(slide, lx + 18, ty + 22, 66, 52)
            icon_wave(slide, lx + 224, ty + 32, 70, 36)
            add_text(slide, label, x(lx + 92), y(ty + 34), w(96), h(34), color=blue, style="module_label")
            add_arrow(slide, lx + 180, cy, lx + 216, cy, color="#111111")
            add_text(slide, out_label, x(lx + 300), y(ty + 35), w(68), h(32), style="module_output")
        elif idx == 1:
            icon_mic(slide, lx + 18, ty + 21, color="#173B5D")
            icon_doc(slide, lx + 210, ty + 23)
            add_text(slide, label, x(lx + 74), y(ty + 35), w(82), h(30), color=green, style="module_label")
            add_arrow(slide, lx + 154, cy, lx + 202, cy, color="#111111")
            clock = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x(lx + 236)), Inches(y(ty + 64)), Inches(w(24)), Inches(h(24)))
            no_fill(clock)
            set_line(clock, "#111111", 1.0)
            add_text(slide, out_label, x(lx + 270), y(ty + 37), w(94), h(38), size=6.9, style="module_output")
        elif idx == 2:
            icon_face(slide, lx + 20, ty + 30, size_px=42)
            icon_face(slide, lx + 206, ty + 21, size_px=70)
            add_text(slide, label, x(lx + 74), y(ty + 38), w(84), h(28), color=purple, style="module_label")
            add_arrow(slide, lx + 154, cy, lx + 200, cy, color="#111111")
            add_text(slide, out_label, x(lx + 282), y(ty + 38), w(82), h(28), style="module_output")
        elif idx == 3:
            icon_pose(slide, lx + 16, ty + 19, scale=0.82)
            icon_pose(slide, lx + 205, ty + 18, scale=0.82, color="#254D8E")
            add_text(slide, label, x(lx + 70), y(ty + 36), w(94), h(30), color=cyan, style="module_label")
            add_arrow(slide, lx + 154, cy, lx + 200, cy, color="#111111")
            add_text(slide, out_label, x(lx + 278), y(ty + 28), w(88), h(58), size=7.0, style="module_output")
        else:
            icon_frames(slide, lx + 18, ty + 20)
            person_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x(lx + 210)), Inches(y(ty + 22)), Inches(w(62)), Inches(h(72)))
            set_fill(person_box, "#D3D0C8")
            set_line(person_box, "#555555", 0.7)
            icon_person(slide, lx + 230, ty + 32, scale=0.38)
            add_text(slide, label, x(lx + 100), y(ty + 24), w(82), h(56), color=orange, style="module_label")
            add_arrow(slide, lx + 172, cy, lx + 204, cy, color="#111111")
            add_text(slide, out_label, x(lx + 292), y(ty + 34), w(70), h(38), style="module_output")

    # Main flow arrows around preprocessing.
    add_arrow(slide, 255, 392, 297, 392, color="#111111")
    add_arrow(slide, 516, 268, 582, 268, color="#111111")
    add_arrow(slide, 516, 392, 582, 392, color="#111111")
    add_arrow(slide, 516, 522, 582, 522, color="#111111")
    add_polyline(slide, [(516, 392), (550, 392), (550, 145), (582, 145)], color="#111111", width_pt=1.2, arrow=True)
    add_polyline(slide, [(516, 392), (550, 392), (550, 650), (582, 650)], color="#111111", width_pt=1.2, arrow=True)
    add_polyline(slide, [(977, 145), (1000, 145), (1000, 270), (1018, 270)], color="#111111", width_pt=1.2, arrow=True)
    add_polyline(slide, [(977, 270), (1000, 270), (1000, 270), (1018, 270)], color="#111111", width_pt=1.2, arrow=True)
    add_polyline(slide, [(977, 397), (1000, 397), (1000, 270), (1018, 270)], color="#111111", width_pt=1.2, arrow=True)
    add_polyline(slide, [(977, 522), (1000, 522), (1000, 270), (1018, 270)], color="#111111", width_pt=1.2, arrow=True)
    add_polyline(slide, [(977, 650), (1000, 650), (1000, 270), (1018, 270)], color="#111111", width_pt=1.2, arrow=True)

    # Timestamp alignment.
    clock = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x(1056)), Inches(y(312)), Inches(w(88)), Inches(h(88)))
    no_fill(clock)
    set_line(clock, blue, 3.0)
    add_line(slide, 1100, 356, 1100, 326, color=blue, width_pt=3)
    add_line(slide, 1100, 356, 1124, 378, color=blue, width_pt=3)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x(1094)), Inches(y(350)), Inches(w(12)), Inches(h(12)))
    set_fill(dot, blue)
    set_line(dot, blue, 0.4)
    add_text(slide, "Align audio,\ntext, face,\nframes, and\npose streams\nto a unified\ntimeline", x(1048), y(438), w(103), h(135), style="body_small")
    add_arrow(slide, 1177, 410, 1202, 410, color="#111111")

    # Five modalities.
    mod_cards = [
        (1216, 200, "#6B57C8", "#FCFAFF", "Face"),
        (1216, 306, "#E17721", "#FFF8F2", "Frame"),
        (1216, 410, cyan, "#F7FFFF", "Pose"),
        (1216, 518, "#2E76CC", "#F8FBFF", "Audio"),
        (1216, 624, "#2D8A45", "#F7FFF7", "Text"),
    ]
    for i, (lx, ty, border, fill, label) in enumerate(mod_cards):
        add_card(slide, lx, ty, 172, 88, border, fill)
        if label == "Face":
            icon_face(slide, lx + 20, ty + 14, size_px=62)
        elif label == "Frame":
            icon_video(slide, lx + 18, ty + 16, 82, 56)
            icon_person(slide, lx + 50, ty + 24, scale=0.32)
        elif label == "Pose":
            icon_pose(slide, lx + 35, ty + 13, scale=0.70, color="#254D8E")
        elif label == "Audio":
            icon_wave(slide, lx + 22, ty + 25, 88, 40)
        else:
            icon_doc(slide, lx + 30, ty + 17, 45, 56)
        add_text(slide, label, x(lx + 102), y(ty + 27), w(72), h(34), color=border, style="module_label")
        add_arrow(slide, 1398, ty + 44, 1425, ty + 44, color="#111111")

    # NEO labels.
    icon_checklist(slide, 1484, 222)
    add_text(slide, "OCEAN Scores", x(1464), y(348), w(122), h(32), size=9.2, bold=True)
    traits = [
        ("O", "Openness", "#775DCC"),
        ("C", "Conscientiousness", "#D76B12"),
        ("E", "Extraversion", "#178B95"),
        ("A", "Agreeableness", "#4E9D56"),
        ("N", "Neuroticism", "#3B7FC6"),
    ]
    for idx, (letter, name, color) in enumerate(traits):
        yy = 403 + idx * 56
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x(1448)), Inches(y(yy)), Inches(w(34)), Inches(h(34)))
        set_fill(circ, color)
        set_line(circ, color, 0.6)
        add_text(slide, letter, x(1448), y(yy + 1), w(34), h(30), color="#FFFFFF", style="badge")
        add_text(slide, name, x(1488), y(yy + 2), w(128), h(30), align=PP_ALIGN.LEFT, style="trait")

    add_arrow(slide, 1618, 412, 1650, 412, color="#111111")

    # Dataset panel.
    icon_database_people(slide, 1684, 352)
    add_text(slide, "Multimodal\nvideos aligned\nacross five\nmodalities with\nNEO-FFI-3\n(OCEAN) labels", x(1668), y(475), w(130), h(152), style="body_small")

    # Legend.
    legend = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x(20)), Inches(y(780)), Inches(w(940)), Inches(h(88)))
    no_fill(legend)
    set_line(legend, "#173B5D", 1.0, dash=True)
    icon_video(slide, 48, 802, 48, 42)
    add_text(slide, "Video", x(110), y(809), w(70), h(28), align=PP_ALIGN.LEFT, style="legend")
    icon_wave(slide, 200, 806, 70, 34)
    add_text(slide, "Audio", x(278), y(809), w(70), h(28), align=PP_ALIGN.LEFT, style="legend")
    icon_doc(slide, 376, 799, 34, 46)
    add_text(slide, "Text", x(422), y(809), w(62), h(28), align=PP_ALIGN.LEFT, style="legend")
    icon_face(slide, 506, 798, size_px=48)
    add_text(slide, "Face", x(566), y(809), w(58), h(28), align=PP_ALIGN.LEFT, style="legend")
    icon_pose(slide, 642, 794, scale=0.58)
    add_text(slide, "Pose", x(700), y(809), w(70), h(28), align=PP_ALIGN.LEFT, style="legend")
    icon_frames(slide, 800, 798)
    add_text(slide, "Frames", x(862), y(809), w(82), h(28), align=PP_ALIGN.LEFT, style="legend")

    # Minimal reproducibility metadata.
    metadata = {
        "summary": "Editable PPTX reconstruction of the provided reference image using PowerPoint text, shapes, connectors, and vector-like icons.",
        "reference_image": str(REFERENCE),
        "output_pptx": str(OUT / "editable_composition.pptx"),
        "editable_layers": ["all visible text", "main containers", "module cards", "arrows", "legend", "simple icons"],
        "approximated_layers": ["human figures", "interviewer portrait", "camera perspective", "database people icon"],
        "text_style_tokens": TEXT_STYLES,
        "canvas": {"reference_px": [REF_W, REF_H], "slide_inches": [SLIDE_W, SLIDE_H]},
    }
    (OUT / "figure_program.json").write_text(json.dumps(metadata, indent=2, ensure_ascii=False), encoding="utf-8")
    (OUT / "rebuild_notes.md").write_text(
        "# Summary\n"
        "This is a PowerPoint-first editable reconstruction of the provided image.\n\n"
        "## Editable\n"
        "- Titles, labels, captions, OCEAN trait text, and legend labels are PPT text boxes.\n"
        "- Section containers, module cards, arrows, connectors, checklist, waves, document, video, pose, database, and camera icons are PPT shapes/connectors.\n\n"
        "## Approximate\n"
        "- Human figures, face portraits, screen views, and database people are editable shape approximations rather than exact traced artwork.\n"
        "- The goal is functional editability and layout fidelity, not pixel-perfect illustration tracing.\n"
        "\n## Reusable Typography Pass\n"
        "- Text is routed through role-based style tokens in `scripts/rebuild_c60_editable.py`.\n"
        "- Reuse `panel_title`, `module_label`, `module_output`, `body`, `trait`, and `legend` for future image-to-PPT rebuilds instead of hand-tuning every text box.\n"
        "- Font weight is controlled by role, with smaller bold module labels and regular body/trait labels to avoid over-heavy text.\n",
        encoding="utf-8",
    )

    pptx = OUT / "editable_composition.pptx"
    prs.save(pptx)


def export_preview() -> None:
    try:
        import win32com.client  # type: ignore
    except Exception as exc:
        (OUT / "preview_export_error.txt").write_text(str(exc), encoding="utf-8")
        return
    app = win32com.client.Dispatch("PowerPoint.Application")
    app.Visible = 1
    presentation = app.Presentations.Open(str(OUT / "editable_composition.pptx"), WithWindow=False)
    try:
        presentation.Slides(1).Export(str(OUT / "rebuild_preview.png"), "PNG", 1806, 920)
    finally:
        presentation.Close()
        app.Quit()


if __name__ == "__main__":
    draw_rebuild()
    export_preview()
