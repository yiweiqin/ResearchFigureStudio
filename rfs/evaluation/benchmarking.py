from __future__ import annotations

import hashlib
import json
import math
import re
from pathlib import Path
from typing import Any

from PIL import Image, ImageChops, ImageFilter, ImageStat
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import requests

from ..utils import ensure_dir, read_json, write_json, write_text


SUITES = {"paper-to-image", "image-to-ppt"}


def _load(path: Path) -> dict:
    if not path.exists():
        return {}
    try:
        value = read_json(path)
    except Exception:
        return {}
    return value if isinstance(value, dict) else {}


def _mean(values: list[float]) -> float:
    return round(sum(values) / max(1, len(values)), 4)


def _clamp(value: object) -> float:
    try:
        return round(max(0.0, min(1.0, float(value))), 4)
    except Exception:
        return 0.0


def _threshold_failures(metrics: dict, thresholds: dict) -> list[str]:
    maximum_metrics = {
        "hallucination_count",
        "forbidden_content_count",
        "plan_forbidden_content_count",
        "full_slide_image_count",
        "blocking_visual_issue_count",
    }
    failures = []
    for key, threshold in thresholds.items():
        value = metrics.get(key)
        if not isinstance(value, (int, float)) or not isinstance(threshold, (int, float)):
            continue
        if key in maximum_metrics and float(value) > float(threshold):
            failures.append(f"{key}={value} above maximum {threshold}")
        elif key not in maximum_metrics and float(value) < float(threshold):
            failures.append(f"{key}={value} below minimum {threshold}")
    return failures


def _normalized_text(value: object) -> str:
    return re.sub(r"[^a-z0-9]+", "", str(value or "").casefold())


def _score_planning_contract(expected: dict, specification: dict) -> dict[str, Any]:
    if not specification:
        return {"available": False}
    actual_entities = []
    for field in ("inputs", "modules", "outputs", "innovations", "must_show"):
        for index, item in enumerate(specification.get(field, []) if isinstance(specification.get(field), list) else []):
            if not isinstance(item, dict):
                continue
            label = next((str(item.get(key) or "").strip() for key in ("visible_label", "name", "text", "statement", "label", "id") if str(item.get(key) or "").strip()), "")
            actual_entities.append({"id": str(item.get("id") or f"{field}_{index}"), "label": label, "normalized": _normalized_text(label), "field": field})
    expanded_relations = set()
    for index, item in enumerate(specification.get("relations", []) if isinstance(specification.get("relations"), list) else []):
        if not isinstance(item, dict):
            continue
        source = str(item.get("source") or item.get("source_id") or "")
        target = str(item.get("target") or item.get("target_id") or "")
        if source and target:
            expanded_relations.add((source, target))
        label = str(item.get("label") or "").strip()
        if source and target and label:
            label_id = f"relation_label_{index}"
            actual_entities.append({"id": label_id, "label": label, "normalized": _normalized_text(label), "field": "relation_label"})
            expanded_relations.add((source, label_id))
            expanded_relations.add((label_id, target))
    expected_entities = [item for item in expected.get("entities", []) if isinstance(item, dict) and item.get("required", True)]
    mapping: dict[str, str] = {}
    matched = []
    used_actual_ids: set[str] = set()
    for item in expected_entities:
        primary = _normalized_text(item.get("label"))
        aliases = [_normalized_text(value) for value in item.get("aliases", []) if _normalized_text(value)]
        def match_score(actual: dict[str, str]) -> int:
            value = actual["normalized"]
            if not value or actual["id"] in used_actual_ids:
                return 0
            bonus = 2 if actual.get("field") in {"inputs", "modules", "outputs", "innovations"} else 1 if actual.get("field") == "relation_label" else 0
            if primary and primary == value:
                return 5 + bonus
            if any(alias == value for alias in aliases):
                return 4 + bonus
            if any(alias in value or value in alias for alias in aliases):
                return 3 + bonus
            if primary and (primary in value or value in primary):
                return 1 + bonus
            return 0
        ranked = sorted(((match_score(actual), actual) for actual in actual_entities), key=lambda pair: pair[0], reverse=True)
        match = ranked[0][1] if ranked and ranked[0][0] > 0 else None
        if match:
            mapping[str(item.get("id"))] = match["id"]
            used_actual_ids.add(match["id"])
            matched.append(str(item.get("id")))
    expected_relations = [item for item in expected.get("relations", []) if isinstance(item, dict) and item.get("required", True)]
    actual_relations = expanded_relations
    adjacency: dict[str, set[str]] = {}
    for source, target in actual_relations:
        adjacency.setdefault(source, set()).add(target)
    def connected(source: str, target: str, max_hops: int = 2) -> bool:
        if (source, target) in actual_relations:
            return True
        frontier = {source}
        visited = {source}
        for _ in range(max_hops):
            frontier = {neighbor for node in frontier for neighbor in adjacency.get(node, set()) if neighbor not in visited}
            if target in frontier:
                return True
            visited.update(frontier)
        return False
    relation_matches = 0
    matched_relations = []
    missing_relations = []
    for relation in expected_relations:
        source = mapping.get(str(relation.get("source")))
        target = mapping.get(str(relation.get("target")))
        if source and target and connected(source, target):
            relation_matches += 1
            matched_relations.append(f"{relation.get('source')}->{relation.get('target')}")
        else:
            missing_relations.append(f"{relation.get('source')}->{relation.get('target')}")
    actual_text = " ".join(item["label"] for item in actual_entities).casefold()
    forbidden = [label for label in expected.get("forbidden_labels", []) if str(label).casefold() in actual_text]
    return {
        "available": True,
        "entity_recall": _clamp(len(matched) / max(1, len(expected_entities))),
        "relation_recall": _clamp(relation_matches / max(1, len(expected_relations))),
        "matched_entity_ids": matched,
        "entity_mapping": mapping,
        "missing_entity_ids": [str(item.get("id")) for item in expected_entities if str(item.get("id")) not in matched],
        "matched_relations": matched_relations,
        "missing_relations": missing_relations,
        "forbidden_labels_found": forbidden,
    }


def validate_benchmark_case(case_dir: str | Path) -> dict[str, Any]:
    root = Path(case_dir).resolve()
    case = _load(root / "case.json")
    errors: list[str] = []
    warnings: list[str] = []
    suite = str(case.get("suite") or "")
    if suite not in SUITES:
        errors.append(f"case.json suite must be one of {sorted(SUITES)}")
    if not str(case.get("case_id") or "").strip():
        errors.append("case.json requires case_id")
    if not isinstance(case.get("thresholds"), dict):
        errors.append("case.json requires thresholds object")

    required_files = ["case.json"]
    if suite == "paper-to-image":
        paper_relative = str(case.get("paper") or "paper.md")
        required_files.append(str(case.get("expected_semantics") or "expected_semantics.json"))
        if not (root / paper_relative).exists():
            source_path = root / str(case.get("source") or "source.json")
            if source_path.exists():
                warnings.append(f"paper input not fetched yet: {paper_relative}")
            else:
                errors.append(f"missing required file: {paper_relative}; no source.json is available")
        expected = _load(root / str(case.get("expected_semantics") or "expected_semantics.json"))
        if not expected.get("entities"):
            errors.append("expected_semantics.json requires entities")
        if not isinstance(expected.get("relations"), list):
            errors.append("expected_semantics.json requires relations list")
        if not isinstance(expected.get("forbidden_labels", []), list):
            errors.append("expected_semantics.json forbidden_labels must be a list")
    elif suite == "image-to-ppt":
        required_files.extend([str(case.get("reference_image") or "reference.png"), str(case.get("expected_objects") or "expected_objects.json")])
        expected = _load(root / str(case.get("expected_objects") or "expected_objects.json"))
        if not isinstance(expected.get("objects"), list):
            errors.append("expected_objects.json requires objects list")
        if not isinstance(expected.get("relations", []), list):
            errors.append("expected_objects.json relations must be a list")

    for relative in required_files:
        if relative and not (root / relative).exists():
            errors.append(f"missing required file: {relative}")
    references = case.get("positive_references", []) + case.get("negative_references", [])
    for relative in references:
        if not (root / str(relative)).exists():
            warnings.append(f"missing optional reference: {relative}")

    return {
        "summary": "Benchmark case validation report.",
        "ok": not errors,
        "case_dir": str(root),
        "case_id": case.get("case_id"),
        "suite": suite,
        "errors": errors,
        "warnings": warnings,
    }


def fetch_benchmark_case(case_dir: str | Path, force: bool = False) -> dict[str, Any]:
    root = Path(case_dir).resolve()
    case = _load(root / "case.json")
    source_path = root / str(case.get("source") or "source.json")
    source = _load(source_path)
    if str(case.get("suite") or "") != "paper-to-image":
        return {"summary": "Benchmark source fetch is only required for paper-to-image cases.", "ok": True, "case_dir": str(root), "status": "not_required"}
    if not source:
        return {"summary": "Benchmark source metadata is missing.", "ok": False, "case_dir": str(root), "status": "missing_source", "error": str(source_path)}
    target = root / str(case.get("paper") or "inputs/paper.pdf")
    if target.exists() and not force:
        return {
            "summary": "Benchmark paper source already exists.",
            "ok": True,
            "case_dir": str(root),
            "status": "reused",
            "paper": str(target),
            "sha256": hashlib.sha256(target.read_bytes()).hexdigest(),
        }
    urls = source.get("paper_urls") if isinstance(source.get("paper_urls"), list) else []
    errors = []
    for url in urls:
        try:
            response = requests.get(str(url), headers={"User-Agent": "ResearchFigureStudio benchmark fetch/0.1"}, timeout=120)
            response.raise_for_status()
            content = response.content
            if not content.startswith(b"%PDF"):
                raise ValueError("downloaded response is not a PDF")
            target.parent.mkdir(parents=True, exist_ok=True)
            target.write_bytes(content)
            digest = hashlib.sha256(content).hexdigest()
            report = {
                "summary": "Benchmark paper source fetched for local use.",
                "ok": True,
                "case_dir": str(root),
                "status": "downloaded",
                "paper": str(target),
                "source_url": str(url),
                "bytes": len(content),
                "sha256": digest,
                "license_note": source.get("license_note"),
            }
            write_json(target.parent / "fetch_report.json", report)
            return report
        except Exception as exc:
            errors.append({"url": str(url), "error": str(exc)})
    return {"summary": "Unable to fetch benchmark paper source.", "ok": False, "case_dir": str(root), "status": "failed", "errors": errors}


def list_benchmark_cases(benchmarks_root: str | Path, suite: str | None = None) -> dict[str, Any]:
    root = Path(benchmarks_root).resolve()
    cases = []
    suites = [suite] if suite else sorted(SUITES)
    for suite_name in suites:
        case_root = root / suite_name / "cases"
        if not case_root.exists():
            continue
        for case_dir in sorted(path for path in case_root.iterdir() if path.is_dir()):
            validation = validate_benchmark_case(case_dir)
            cases.append({
                "case_id": validation.get("case_id") or case_dir.name,
                "suite": suite_name,
                "case_dir": str(case_dir),
                "valid": validation["ok"],
                "errors": validation["errors"],
            })
    return {"summary": "Available ResearchFigureStudio benchmark cases.", "benchmarks_root": str(root), "cases": cases}


def _select_candidate_review(run: Path) -> dict:
    report = _load(run / "candidate_review.json")
    candidates = [item for item in report.get("candidates", []) if isinstance(item, dict)]
    selected_id = report.get("selected_candidate_id")
    selected = next((item for item in candidates if item.get("candidate_id") == selected_id), None)
    if not selected and candidates:
        selected = max(candidates, key=lambda item: float(item.get("score") or 0.0))
    return selected.get("review", {}) if isinstance(selected, dict) else {}


def score_paper_to_image(case_dir: str | Path, run_dir: str | Path) -> dict[str, Any]:
    case_root, run = Path(case_dir).resolve(), Path(run_dir).resolve()
    case = _load(case_root / "case.json")
    expected = _load(case_root / str(case.get("expected_semantics") or "expected_semantics.json"))
    review = _select_candidate_review(run)
    benchmark_review = _load(run / "benchmark_review.json")
    if benchmark_review:
        review = {**review, **benchmark_review}

    scientific = review.get("scientific", {}) if isinstance(review.get("scientific"), dict) else {}
    ocr = review.get("ocr", {}) if isinstance(review.get("ocr"), dict) else {}
    aesthetic = review.get("aesthetic", {}) if isinstance(review.get("aesthetic"), dict) else {}
    clarity = review.get("clarity", {}) if isinstance(review.get("clarity"), dict) else {}
    information = review.get("information", {}) if isinstance(review.get("information"), dict) else {}
    stability = review.get("stability", {}) if isinstance(review.get("stability"), dict) else {}
    planning_contract = _score_planning_contract(expected, _load(run / "figure_specification.json"))

    expected_entities = [item for item in expected.get("entities", []) if isinstance(item, dict) and item.get("required", True)]
    expected_relations = [item for item in expected.get("relations", []) if isinstance(item, dict) and item.get("required", True)]
    missing_entities = scientific.get("missing_modules", []) or benchmark_review.get("missing_entities", [])
    missing_relations = scientific.get("missing_relations", []) or benchmark_review.get("missing_relations", [])
    invented = scientific.get("invented_items", []) or benchmark_review.get("invented_items", [])
    forbidden_found = ocr.get("forbidden_labels_found", []) or benchmark_review.get("forbidden_labels_found", [])
    missing_labels = ocr.get("missing_labels", []) or benchmark_review.get("missing_labels", [])
    misspelled = ocr.get("misspelled_labels", []) or benchmark_review.get("misspelled_labels", [])

    entity_recall = _clamp((len(expected_entities) - len(missing_entities)) / max(1, len(expected_entities)))
    relation_recall = _clamp((len(expected_relations) - len(missing_relations)) / max(1, len(expected_relations)))
    exact_label_rate = _clamp((len(expected_entities) - len(missing_labels) - len(misspelled)) / max(1, len(expected_entities)))
    scientific_score = _clamp(scientific.get("score", _mean([entity_recall, relation_recall])))
    aesthetic_dimensions = [aesthetic.get(key) for key in ("hierarchy", "balance", "whitespace", "color", "icon_consistency", "readability") if aesthetic.get(key) is not None]
    aesthetic_score = _clamp(aesthetic.get("score", _mean([_clamp(value) for value in aesthetic_dimensions])))
    clarity_score = _clamp(clarity.get("score", aesthetic.get("readability", 0.0)))
    information_score = _clamp(information.get("score", information.get("coverage", scientific_score)))
    stability_score = _clamp(stability.get("production_pass_rate", stability.get("score", 0.0)))
    hard_failures = list(review.get("hard_errors", []))
    if not scientific or not aesthetic or scientific.get("engineering_only") or aesthetic.get("engineering_only"):
        hard_failures.append("frozen benchmark judge did not provide scientific and aesthetic sections")
    if invented:
        hard_failures.append("invented scientific content")
    if forbidden_found:
        hard_failures.append("forbidden labels found")
    if relation_recall < 1.0:
        hard_failures.append("required relations missing or incorrect")
    if exact_label_rate < 1.0:
        hard_failures.append("required labels missing or misspelled")
    plan_entity_min = float(case.get("thresholds", {}).get("plan_entity_recall", 1.0))
    plan_relation_min = float(case.get("thresholds", {}).get("plan_relation_recall", 1.0))
    if planning_contract.get("available") and float(planning_contract.get("entity_recall", 0.0)) < plan_entity_min:
        hard_failures.append("paper planning missed required benchmark entities")
    if planning_contract.get("available") and float(planning_contract.get("relation_recall", 0.0)) < plan_relation_min:
        hard_failures.append("paper planning missed required benchmark relations")
    if planning_contract.get("forbidden_labels_found"):
        hard_failures.append("paper planning introduced forbidden benchmark labels")

    metrics = {
        "entity_recall": entity_recall,
        "relation_recall": relation_recall,
        "exact_label_rate": exact_label_rate,
        "scientific_score": scientific_score,
        "information_score": information_score,
        "clarity_score": clarity_score,
        "aesthetic_score": aesthetic_score,
        "stability_score": stability_score,
        "hallucination_count": len(invented),
        "forbidden_content_count": len(forbidden_found),
        "plan_entity_recall": planning_contract.get("entity_recall"),
        "plan_relation_recall": planning_contract.get("relation_recall"),
        "plan_forbidden_content_count": len(planning_contract.get("forbidden_labels_found", [])),
    }
    thresholds = case.get("thresholds", {})
    threshold_failures = _threshold_failures(metrics, thresholds)
    total_score = round(0.40 * scientific_score + 0.15 * information_score + 0.15 * clarity_score + 0.20 * aesthetic_score + 0.10 * stability_score, 4)
    return {
        "summary": "Paper-to-image benchmark score.",
        "suite": "paper-to-image",
        "case_id": case.get("case_id"),
        "case_dir": str(case_root),
        "run_dir": str(run),
        "metrics": metrics,
        "planning_contract": planning_contract,
        "total_score": total_score,
        "hard_failures": sorted(set(hard_failures)),
        "threshold_failures": threshold_failures,
        "passed": not hard_failures and not threshold_failures,
        "policy": "scientific, relation, terminology, and hallucination failures cannot be offset by aesthetics",
    }


def _image_similarity(reference: Path, preview: Path) -> dict[str, float | None]:
    if not reference.exists() or not preview.exists():
        return {"pixel_similarity": None, "edge_similarity": None, "color_similarity": None}
    with Image.open(reference) as ref_image, Image.open(preview) as out_image:
        ref = ref_image.convert("RGB").resize((512, 512), Image.Resampling.LANCZOS)
        out = out_image.convert("RGB").resize((512, 512), Image.Resampling.LANCZOS)
    difference = ImageChops.difference(ref, out)
    rms = math.sqrt(sum(value * value for value in ImageStat.Stat(difference).rms) / 3.0)
    pixel_similarity = _clamp(1.0 - rms / 255.0)
    ref_edge = ref.convert("L").filter(ImageFilter.FIND_EDGES).point(lambda value: 255 if value > 32 else 0)
    out_edge = out.convert("L").filter(ImageFilter.FIND_EDGES).point(lambda value: 255 if value > 32 else 0)
    edge_diff = ImageChops.difference(ref_edge, out_edge)
    edge_similarity = _clamp(1.0 - ImageStat.Stat(edge_diff).mean[0] / 255.0)
    ref_mean, out_mean = ImageStat.Stat(ref).mean, ImageStat.Stat(out).mean
    color_error = sum(abs(a - b) for a, b in zip(ref_mean, out_mean)) / (3.0 * 255.0)
    return {"pixel_similarity": pixel_similarity, "edge_similarity": edge_similarity, "color_similarity": _clamp(1.0 - color_error)}


def _ppt_editability(pptx_path: Path) -> dict[str, Any]:
    if not pptx_path.exists():
        return {"status": "missing", "full_slide_image_count": 0, "text_shape_count": 0, "connector_count": 0, "shape_count": 0}
    presentation = Presentation(str(pptx_path))
    full_slide_images = 0
    text_shapes = 0
    connectors = 0
    shapes = 0
    slide_width, slide_height = float(presentation.slide_width), float(presentation.slide_height)
    for slide in presentation.slides:
        for shape in slide.shapes:
            shapes += 1
            if getattr(shape, "has_text_frame", False) and str(getattr(shape, "text", "")).strip():
                text_shapes += 1
            if shape.shape_type == MSO_SHAPE_TYPE.LINE:
                connectors += 1
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                coverage = float(shape.width) * float(shape.height) / max(1.0, slide_width * slide_height)
                if coverage >= 0.90:
                    full_slide_images += 1
    return {
        "status": "pass" if full_slide_images == 0 else "blocked",
        "full_slide_image_count": full_slide_images,
        "text_shape_count": text_shapes,
        "connector_count": connectors,
        "shape_count": shapes,
    }


def _bbox_iou(left: dict, right: dict) -> float:
    lx0, ly0 = float(left.get("x", 0)), float(left.get("y", 0))
    lx1, ly1 = lx0 + float(left.get("w", 0)), ly0 + float(left.get("h", 0))
    rx0, ry0 = float(right.get("x", 0)), float(right.get("y", 0))
    rx1, ry1 = rx0 + float(right.get("w", 0)), ry0 + float(right.get("h", 0))
    intersection = max(0.0, min(lx1, rx1) - max(lx0, rx0)) * max(0.0, min(ly1, ry1) - max(ly0, ry0))
    union = max(0.000001, float(left.get("w", 0)) * float(left.get("h", 0)) + float(right.get("w", 0)) * float(right.get("h", 0)) - intersection)
    return intersection / union


def _center_error(left: dict, right: dict) -> float:
    lx = float(left.get("x", 0)) + float(left.get("w", 0)) / 2
    ly = float(left.get("y", 0)) + float(left.get("h", 0)) / 2
    rx = float(right.get("x", 0)) + float(right.get("w", 0)) / 2
    ry = float(right.get("y", 0)) + float(right.get("h", 0)) / 2
    return math.sqrt((lx - rx) ** 2 + (ly - ry) ** 2)


def _match_expected_objects(expected_objects: list[dict], geometry: dict) -> tuple[list[dict], dict[str, str]]:
    actual = []
    for collection, object_type in (("panels", "panel"), ("cards", "card"), ("slots", "asset")):
        actual.extend({**item, "benchmark_type": object_type} for item in geometry.get(collection, []) if isinstance(item, dict) and isinstance(item.get("bbox_percent"), dict))
    unmatched = list(actual)
    matches = []
    id_mapping: dict[str, str] = {}
    for expected in expected_objects:
        expected_box = expected.get("bbox_percent") if isinstance(expected.get("bbox_percent"), dict) else {}
        expected_type = str(expected.get("type") or "")
        same_id = next((item for item in unmatched if str(item.get("id")) == str(expected.get("id"))), None)
        candidates = [item for item in unmatched if expected_type in {"", str(item.get("benchmark_type"))}]
        if not candidates:
            candidates = unmatched
        chosen = same_id or max(candidates, key=lambda item: _bbox_iou(expected_box, item["bbox_percent"]), default=None)
        if not chosen:
            matches.append({"expected_id": expected.get("id"), "status": "unmatched", "iou": 0.0, "center_error": 1.0})
            continue
        unmatched.remove(chosen)
        iou = _bbox_iou(expected_box, chosen["bbox_percent"])
        center_error = _center_error(expected_box, chosen["bbox_percent"])
        id_mapping[str(expected.get("id"))] = str(chosen.get("id"))
        matches.append({
            "expected_id": expected.get("id"),
            "actual_id": chosen.get("id"),
            "expected_type": expected_type,
            "actual_type": chosen.get("benchmark_type"),
            "iou": round(iou, 4),
            "center_error": round(center_error, 4),
            "status": "matched" if iou >= 0.30 else "low_overlap",
        })
    return matches, id_mapping


def score_image_to_ppt(case_dir: str | Path, run_dir: str | Path) -> dict[str, Any]:
    case_root, run = Path(case_dir).resolve(), Path(run_dir).resolve()
    case = _load(case_root / "case.json")
    expected = _load(case_root / str(case.get("expected_objects") or "expected_objects.json"))
    reference = case_root / str(case.get("reference_image") or "reference.png")
    preview = run / str(case.get("preview") or "rebuild_preview.png")
    pptx = run / str(case.get("pptx") or "editable_composition.pptx")
    similarity = _image_similarity(reference, preview)
    editability = _ppt_editability(pptx)
    visual = _load(run / "rebuild_visual_quality_report.json")
    semantic = _load(run / "semantic_binding_report.json")
    composition = _load(run / "composition_quality_report.json")
    alignment = _load(run / "text_alignment_report.json")
    geometry = _load(run / "reference_geometry.json")

    expected_objects = [item for item in expected.get("objects", []) if isinstance(item, dict)]
    expected_relations = [item for item in expected.get("relations", []) if isinstance(item, dict)]
    object_matches, id_mapping = _match_expected_objects(expected_objects, geometry)
    object_coverage = _clamp(sum(1 for item in object_matches if item["status"] == "matched") / max(1, len(expected_objects))) if expected_objects else 1.0
    mean_object_iou = _mean([float(item["iou"]) for item in object_matches]) if object_matches else 1.0
    mean_center_error = _mean([float(item["center_error"]) for item in object_matches]) if object_matches else 0.0
    actual_relations = {
        (str(item.get("source_id") or item.get("source") or ""), str(item.get("target_id") or item.get("target") or ""))
        for item in composition.get("arrows", [])
        if isinstance(item, dict)
    }
    matched_relations = 0
    for relation in expected_relations:
        source = id_mapping.get(str(relation.get("source")))
        target = id_mapping.get(str(relation.get("target")))
        if source and target and (source, target) in actual_relations:
            matched_relations += 1
    relation_coverage = _clamp(matched_relations / max(1, len(expected_relations))) if expected_relations else 1.0
    max_text_delta = max((float(item.get("center_delta_percent") or 0.0) for item in alignment.get("items", []) if isinstance(item, dict)), default=0.0)
    text_alignment_score = _clamp(1.0 - max_text_delta / 0.10)
    blocking_visual_issues = int(visual.get("blocking_issue_count") or 0)
    semantic_binding_score = _clamp(semantic.get("mapped_entity_count", 0) / max(1, semantic.get("entity_count", 0))) if semantic else 1.0
    editable_score = 0.0 if editability["status"] == "blocked" else _clamp((min(1, editability["text_shape_count"]) + min(1, editability["connector_count"]) + 1) / 3)
    fidelity_values = [value for value in similarity.values() if isinstance(value, float)]
    render_fidelity = _mean(fidelity_values) if fidelity_values else 0.0
    metrics = {
        **similarity,
        "render_fidelity": render_fidelity,
        "object_coverage": object_coverage,
        "mean_object_iou": mean_object_iou,
        "mean_center_error": mean_center_error,
        "relation_coverage": relation_coverage,
        "text_alignment_score": text_alignment_score,
        "semantic_binding_score": semantic_binding_score,
        "editability_score": editable_score,
        "full_slide_image_count": editability["full_slide_image_count"],
        "blocking_visual_issue_count": blocking_visual_issues,
    }
    hard_failures = []
    if editability["full_slide_image_count"]:
        hard_failures.append("full-slide reference image detected")
    if blocking_visual_issues:
        hard_failures.append("blocking visual issues remain")
    thresholds = case.get("thresholds", {})
    threshold_failures = _threshold_failures(metrics, thresholds)
    total_score = round(0.40 * render_fidelity + 0.20 * object_coverage + 0.15 * relation_coverage + 0.10 * text_alignment_score + 0.15 * editable_score, 4)
    return {
        "summary": "Image-to-editable-PPT benchmark score.",
        "suite": "image-to-ppt",
        "case_id": case.get("case_id"),
        "case_dir": str(case_root),
        "run_dir": str(run),
        "metrics": metrics,
        "editability": editability,
        "object_matches": object_matches,
        "total_score": total_score,
        "hard_failures": hard_failures,
        "threshold_failures": threshold_failures,
        "passed": not hard_failures and not threshold_failures,
        "limitations": [
            "pixel, edge, and color similarity are baseline metrics; calibrated SSIM/LPIPS and region annotations remain future work",
            "mutation-based editability testing requires a dedicated render-after-edit harness",
        ],
    }


def score_benchmark_case(case_dir: str | Path, run_dir: str | Path, out: str | Path | None = None) -> dict[str, Any]:
    validation = validate_benchmark_case(case_dir)
    if not validation["ok"]:
        return {**validation, "passed": False}
    result = score_paper_to_image(case_dir, run_dir) if validation["suite"] == "paper-to-image" else score_image_to_ppt(case_dir, run_dir)
    target = ensure_dir(out or run_dir)
    write_json(target / "benchmark_result.json", result)
    lines = [
        "# Benchmark Result",
        "",
        f"- Suite: {result.get('suite')}",
        f"- Case: {result.get('case_id')}",
        f"- Passed: {result.get('passed')}",
        f"- Total score: {result.get('total_score')}",
        "",
        "## Metrics",
    ]
    lines.extend(f"- {key}: {value}" for key, value in result.get("metrics", {}).items())
    lines.extend(["", "## Hard failures"])
    lines.extend(f"- {item}" for item in result.get("hard_failures", []))
    lines.extend(["", "## Threshold failures"])
    lines.extend(f"- {item}" for item in result.get("threshold_failures", []))
    write_text(target / "benchmark_result.md", "\n".join(lines) + "\n")
    return result


def run_benchmark_case(case_dir: str | Path, out: str | Path) -> dict[str, Any]:
    case_root = Path(case_dir).resolve()
    case_preflight = _load(case_root / "case.json")
    if str(case_preflight.get("suite") or "") == "paper-to-image":
        paper_path = case_root / str(case_preflight.get("paper") or "paper.md")
        if not paper_path.exists():
            fetched = fetch_benchmark_case(case_root)
            if not fetched.get("ok"):
                return {**fetched, "passed": False}
    validation = validate_benchmark_case(case_dir)
    if not validation["ok"]:
        return {**validation, "passed": False}
    root = Path(case_dir).resolve()
    target = ensure_dir(out).resolve()
    case = _load(root / "case.json")
    config = case.get("run_config", {}) if isinstance(case.get("run_config"), dict) else {}
    if validation["suite"] == "paper-to-image":
        from ..paper_to_image import run_paper_to_image

        workflow_result = run_paper_to_image(
            paper=root / str(case.get("paper") or "paper.md"),
            out=target,
            preferences_path=(root / str(case["preferences"])) if case.get("preferences") else None,
            positive_references=[str(root / str(path)) for path in case.get("positive_references", [])],
            negative_references=[str(root / str(path)) for path in case.get("negative_references", [])],
            planner_mode=str(config.get("planner_mode") or "heuristic"),
            asset_mode=str(config.get("image_asset_mode") or "placeholder"),
            candidates=int(config.get("image_candidates") or 1),
            aspect_ratio=str(config.get("aspect_ratio") or "auto"),
            language=str(config.get("language") or "English"),
            review_mode=str(config.get("review_mode") or "heuristic"),
            domain_profile=str(config.get("domain_profile") or "auto"),
            template=str(config.get("template") or "auto"),
            repair_rounds=int(config.get("repair_rounds") or 0),
            ocr_engine=str(config.get("ocr_engine") or "off"),
            ocr_lang=str(config.get("ocr_lang") or "en_ch"),
        )
    else:
        from ..editable_rebuild import rebuild_editable

        workflow_result = rebuild_editable(
            reference=root / str(case.get("reference_image") or "reference.png"),
            out=target,
            asset_mode=str(config.get("asset_mode") or "placeholder"),
            asset_policy=str(config.get("asset_policy") or "smart-api"),
            text_mode=str(config.get("text_mode") or "off"),
            layout_mode=str(config.get("layout_mode") or "heuristic"),
            control_mode=str(config.get("control_mode") or "heuristic"),
            design_plan_mode=str(config.get("design_plan_mode") or "heuristic"),
            export_preview=True,
            ocr_engine=str(config.get("ocr_engine") or "off"),
        )
    score = score_benchmark_case(root, target, target)
    result = {
        "summary": "Benchmark workflow and scoring completed.",
        "ok": bool(workflow_result.get("ok")),
        "suite": validation["suite"],
        "case_id": validation["case_id"],
        "out_dir": str(target),
        "workflow_result": workflow_result,
        "benchmark_result": score,
    }
    write_json(target / "benchmark_run.json", result)
    return result


def run_fast_benchmark_case(
    case_dir: str | Path,
    out: str | Path,
    deadline_seconds: int = 180,
    planner_mode: str = "vlm",
    planner_model: str | None = None,
    ocr_engine: str = "off",
) -> dict[str, Any]:
    case_root = Path(case_dir).resolve()
    case = _load(case_root / "case.json")
    if str(case.get("suite") or "") != "paper-to-image":
        return {"summary": "Fast framework benchmark only supports paper-to-image cases.", "ok": False, "case_dir": str(case_root)}
    paper_path = case_root / str(case.get("paper") or "paper.md")
    if not paper_path.exists():
        fetched = fetch_benchmark_case(case_root)
        if not fetched.get("ok"):
            return {**fetched, "passed": False}
    validation = validate_benchmark_case(case_root)
    if not validation.get("ok"):
        return {**validation, "passed": False}
    from ..paper_to_image import run_fast_framework_prompt

    preferences = case_root / str(case.get("preferences")) if case.get("preferences") else None
    run = run_fast_framework_prompt(
        paper=paper_path,
        out=out,
        deadline_seconds=deadline_seconds,
        planner_mode=planner_mode,
        planner_model=planner_model,
        ocr_engine=ocr_engine,
        preferences_path=preferences if preferences and preferences.exists() else None,
        aspect_ratio=str(case.get("run_config", {}).get("aspect_ratio") or "16:9"),
    )
    score = score_benchmark_case(case_root, out)
    result = {
        "summary": "Fast paper framework benchmark completed.",
        "ok": bool(run.get("ok")),
        "case_id": case.get("case_id"),
        "run": run,
        "benchmark": score,
        "passed_planning_thresholds": not any("plan_" in value for value in score.get("threshold_failures", [])),
    }
    write_json(Path(out) / "fast_benchmark_result.json", result)
    return result


def run_fast_benchmark_suite(
    benchmarks_root: str | Path,
    out: str | Path,
    case_ids: list[str] | None = None,
    deadline_seconds: int = 180,
    planner_mode: str = "vlm",
    planner_model: str | None = None,
    ocr_engine: str = "off",
) -> dict[str, Any]:
    root = Path(benchmarks_root).resolve()
    target = ensure_dir(out).resolve()
    selected = set(case_ids or [])
    listing = list_benchmark_cases(root, suite="paper-to-image")
    cases = [item for item in listing.get("cases", []) if item.get("valid") and (not selected or str(item.get("case_id")) in selected)]
    results = []
    for item in cases:
        case_id = str(item.get("case_id"))
        case_out = target / case_id
        try:
            result = run_fast_benchmark_case(
                item["case_dir"],
                case_out,
                deadline_seconds=deadline_seconds,
                planner_mode=planner_mode,
                planner_model=planner_model,
                ocr_engine=ocr_engine,
            )
        except Exception as exc:
            result = {"summary": "Fast benchmark case raised an exception.", "ok": False, "case_id": case_id, "error": str(exc)}
        run = result.get("run", {}) if isinstance(result.get("run"), dict) else {}
        benchmark = result.get("benchmark", {}) if isinstance(result.get("benchmark"), dict) else {}
        metrics = benchmark.get("metrics", {}) if isinstance(benchmark.get("metrics"), dict) else {}
        provider = run.get("provider", {}) if isinstance(run.get("provider"), dict) else {}
        results.append({
            "case_id": case_id,
            "ok": bool(result.get("ok")),
            "production_ready": bool(run.get("production_ready")),
            "passed_planning_thresholds": bool(result.get("passed_planning_thresholds")),
            "contract_source": run.get("contract_source"),
            "cache_hit": bool(run.get("cache_hit")),
            "document_cache_hit": bool(run.get("document_cache_hit")),
            "elapsed_seconds": run.get("elapsed_seconds"),
            "stage_timings": run.get("stage_timings", {}),
            "provider": provider,
            "plan_entity_recall": metrics.get("plan_entity_recall"),
            "plan_relation_recall": metrics.get("plan_relation_recall"),
            "plan_forbidden_content_count": metrics.get("plan_forbidden_content_count"),
            "threshold_failures": benchmark.get("threshold_failures", []),
            "error": result.get("error"),
        })

    def numbers(field: str) -> list[float]:
        values = []
        for item in results:
            value = item.get(field)
            if isinstance(value, (int, float)):
                values.append(float(value))
        return values

    def timing(name: str) -> list[float]:
        return [float(item.get("stage_timings", {}).get(name)) for item in results if isinstance(item.get("stage_timings", {}).get(name), (int, float))]

    def percentile(values: list[float], fraction: float) -> float:
        if not values:
            return 0.0
        ordered = sorted(values)
        index = min(len(ordered) - 1, max(0, math.ceil(len(ordered) * fraction) - 1))
        return round(ordered[index], 4)

    provider_attempts = sum(int(item.get("provider", {}).get("attempts") or 0) for item in results)
    provider_retries = sum(int(item.get("provider", {}).get("retries_used") or 0) for item in results)
    provider_calls = [item for item in results if item.get("provider", {}).get("attempts")]
    failure_categories: dict[str, int] = {}
    for item in provider_calls:
        for category in item.get("provider", {}).get("failure_categories", []) or []:
            failure_categories[str(category)] = failure_categories.get(str(category), 0) + 1
    total_times = numbers("elapsed_seconds")
    entity_scores = numbers("plan_entity_recall")
    relation_scores = numbers("plan_relation_recall")
    aggregate = {
        "case_count": len(results),
        "successful_run_count": sum(bool(item.get("ok")) for item in results),
        "production_ready_count": sum(bool(item.get("production_ready")) for item in results),
        "planning_threshold_pass_count": sum(bool(item.get("passed_planning_thresholds")) for item in results),
        "cache_hit_count": sum(bool(item.get("cache_hit")) for item in results),
        "cache_hit_rate": round(sum(bool(item.get("cache_hit")) for item in results) / max(1, len(results)), 4),
        "document_cache_hit_count": sum(bool(item.get("document_cache_hit")) for item in results),
        "document_cache_hit_rate": round(sum(bool(item.get("document_cache_hit")) for item in results) / max(1, len(results)), 4),
        "mean_plan_entity_recall": _mean(entity_scores),
        "mean_plan_relation_recall": _mean(relation_scores),
        "forbidden_content_total": sum(int(item.get("plan_forbidden_content_count") or 0) for item in results),
        "mean_total_seconds": _mean(total_times),
        "p95_total_seconds": percentile(total_times, 0.95),
        "mean_document_preparation_seconds": _mean(timing("document_preparation_seconds")),
        "mean_document_extraction_seconds": _mean(timing("document_extraction_seconds")),
        "mean_semantic_compilation_seconds": _mean(timing("semantic_compilation_seconds")),
        "provider_call_count": len(provider_calls),
        "provider_success_count": sum(bool(item.get("provider", {}).get("success")) for item in provider_calls),
        "provider_success_rate": round(sum(bool(item.get("provider", {}).get("success")) for item in provider_calls) / len(provider_calls), 4) if provider_calls else None,
        "provider_attempts": provider_attempts,
        "provider_retries_used": provider_retries,
        "provider_failure_categories": failure_categories,
    }
    report = {
        "summary": "Fast paper framework benchmark suite completed.",
        "ok": bool(results) and all(item.get("ok") for item in results),
        "benchmarks_root": str(root),
        "out_dir": str(target),
        "planner_mode": planner_mode,
        "planner_model": planner_model,
        "deadline_seconds": deadline_seconds,
        "ocr_engine": ocr_engine,
        "selected_case_ids": [item.get("case_id") for item in cases],
        "aggregate": aggregate,
        "cases": results,
    }
    write_json(target / "fast_suite_report.json", report)
    return report
