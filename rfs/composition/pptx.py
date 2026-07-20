from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml import parse_xml
from pptx.util import Inches, Pt

from ..utils import pct_to_inches, write_json


def _rgb(hex_color: str) -> RGBColor:
    value = hex_color.strip().lstrip("#") or "000000"
    return RGBColor(int(value[:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _apply_arrow(connector, size: str = "sm") -> None:
    ln = connector.line._get_or_add_ln()
    ln.append(parse_xml(f'<a:tailEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" type="triangle" w="{size}" len="{size}"/>'))


def _apply_line_cap(connector, cap: str = "round") -> None:
    value = {"round": "rnd", "square": "sq"}.get(str(cap).lower())
    if not value:
        return
    ln = connector.line._get_or_add_ln()
    ln.set("cap", value)


def _connector_type_for_route(route_style: str):
    if str(route_style).lower() in {"soft_curve", "dashed_loop", "dashed_spline_like"}:
        return MSO_CONNECTOR.CURVE
    return MSO_CONNECTOR.STRAIGHT


def _set_text(shape, text: str, font_size: float = 10, bold: bool = False, color: str = "#163B4D", align=PP_ALIGN.CENTER, font_family: str | None = None) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(float(font_size))
    run.font.bold = bold
    if font_family:
        run.font.name = str(font_family)
    run.font.color.rgb = _rgb(color)


def _add_round_rect(slide, x: float, y: float, w: float, h: float, fill: str, stroke: str, width_pt: float = 1.2):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    shape.line.color.rgb = _rgb(stroke)
    shape.line.width = Pt(width_pt)
    shape.shadow.inherit = False
    return shape


def _add_card_frame(slide, card: dict, width_in: float, height_in: float):
    x, y, w, h = pct_to_inches(card["bbox_percent"], width_in, height_in)
    shape_type = MSO_SHAPE.RECTANGLE if str(card.get("shape_kind") or "rounded_rect").lower() == "rect" else MSO_SHAPE.ROUNDED_RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if float(card.get("fill_transparency", 1.0)) >= 1.0:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(str(card.get("fill_color") or "#FFFFFF"))
    shape.line.color.rgb = _rgb(str(card.get("stroke_color") or "#59AFCB"))
    shape.line.width = Pt(float(card.get("stroke_width_pt") or 1.5))
    dash_style = str(card.get("dash_style") or "solid").lower()
    if dash_style in {"dash", "dashed"}:
        shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    elif dash_style in {"dot", "dotted"}:
        shape.line.dash_style = getattr(MSO_LINE_DASH_STYLE, "ROUND_DOT", MSO_LINE_DASH_STYLE.DASH)
    shape.shadow.inherit = False
    return shape


def _add_label(slide, text: str, x: float, y: float, w: float, h: float, font_size: float = 9, bold: bool = False, align=PP_ALIGN.CENTER, color: str = "#163B4D", font_family: str | None = None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    _set_text(box, text, font_size=font_size, bold=bold, color=color, align=align, font_family=font_family)
    return box


def _add_title_block(slide, program: dict, width_in: float, height_in: float) -> None:
    title_block = program.get("title_block")
    if not isinstance(title_block, dict):
        return
    title = str(title_block.get("title", "")).strip()
    subtitle = str(title_block.get("subtitle", "")).strip()
    bbox = title_block.get("bbox_percent")
    if not title or not isinstance(bbox, dict):
        return
    x, y, w, h = pct_to_inches(bbox, width_in, height_in)
    title_h = min(h * 0.58, 0.42)
    subtitle_h = max(0.18, min(h - title_h, 0.26))
    _add_label(
        slide,
        title,
        x,
        y,
        w,
        title_h,
        font_size=int(title_block.get("title_font_size", 24)),
        bold=True,
        align=PP_ALIGN.LEFT,
    )
    if subtitle:
        subtitle_box = _add_label(
            slide,
            subtitle,
            x,
            y + title_h * 0.94,
            w,
            subtitle_h,
            font_size=int(title_block.get("subtitle_font_size", 11)),
            bold=False,
            align=PP_ALIGN.LEFT,
        )
        for paragraph in subtitle_box.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = _rgb(str(title_block.get("subtitle_color", "#333333")))


def _short_caption(text: str) -> str:
    replacements = {
        "participant seated before 49-inch screen": "Participant setup",
        "3D virtual interviewer on screen": "Virtual interviewer",
        "wide-angle camera capturing full body": "Full-body camera",
        "virtual interview room setup": "Interview room",
        "interview question prompt on screen": "Question prompt",
        "raw interview video recording": "Raw video",
        "287 participant video collection": "287 participants",
        "36 interview questions": "36 questions",
        "full-body video thumbnails": "Full-body clips",
        "participant clip grid": "Clip grid",
        "FFmpeg extracts audio from video": "FFmpeg audio",
        "audio waveform stream": "Audio stream",
        "FunASR speech recognition": "FunASR ASR",
        "spoken text with timestamps": "Text + timestamps",
        "MTCNN face detection": "MTCNN face",
        "face clips": "Face clips",
        "AlphaPose full-body skeleton extraction": "AlphaPose pose",
        "full-body pose skeleton stream": "Pose stream",
        "frame sampling from video": "Frame sampling",
        "sampled video frames": "Video frames",
        "timestamp alignment clock": "Alignment clock",
        "aligned audio and text streams": "Audio + text",
        "aligned face and pose streams": "Face + pose",
        "aligned frame stream": "Frame stream",
        "face modality": "Face",
        "frame modality": "Frame",
        "pose modality": "Pose",
        "audio modality": "Audio",
        "text modality": "Text",
        "NEO-FFI-3 questionnaire": "NEO-FFI-3",
        "OCEAN score vector": "OCEAN scores",
        "Openness score": "Openness",
        "Conscientiousness score": "Conscientiousness",
        "Extraversion score": "Extraversion",
        "Agreeableness score": "Agreeableness",
        "Neuroticism score": "Neuroticism",
    }
    if text in replacements:
        return replacements[text]
    if len(text) <= 22:
        return text
    words = text.replace("/", " ").split()
    return " ".join(words[:3]) if len(words) > 3 else text[:22]


def _picture_contain_box(image_path: Path, x: float, y: float, w: float, h: float) -> tuple[float, float, float, float, float]:
    with Image.open(image_path) as img:
        iw, ih = img.size
    image_ratio = iw / max(ih, 1)
    slot_ratio = w / max(h, 0.001)
    if image_ratio > slot_ratio:
        fit_w = w
        fit_h = w / image_ratio
    else:
        fit_h = h
        fit_w = h * image_ratio
    left = x + (w - fit_w) / 2
    top = y + (h - fit_h) / 2
    fill_percent = fit_w * fit_h / max(w * h, 0.001) * 100
    return left, top, fit_w, fit_h, fill_percent


def _add_picture_contain(slide, image_path: Path, x: float, y: float, w: float, h: float):
    left, top, fit_w, fit_h, _fill_percent = _picture_contain_box(image_path, x, y, w, h)
    return slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(fit_w), height=Inches(fit_h))


def _panel_map(program: dict) -> dict[str, dict]:
    return {panel["id"]: panel for panel in program.get("panels", [])}


def _text_program_items(program: dict) -> list[dict]:
    text_program = program.get("text_program")
    if not isinstance(text_program, dict):
        return []
    items = text_program.get("items", [])
    return [item for item in items if isinstance(item, dict) and item.get("visible", True)]


def _text_item_for_target(program: dict, target_id: str, role: str) -> dict | None:
    for item in _text_program_items(program):
        if str(item.get("target_id")) == target_id and str(item.get("role")) == role:
            return item
    return None


def _align_from_text_item(item: dict):
    value = str(item.get("align") or "").lower()
    if value == "left":
        return PP_ALIGN.LEFT
    if value == "right":
        return PP_ALIGN.RIGHT
    return PP_ALIGN.CENTER


def _object_map(program: dict) -> dict[str, dict]:
    objects = {panel["id"]: panel for panel in program.get("panels", [])}
    objects.update({slot["id"]: slot for slot in program.get("slots", [])})
    return objects


def _bbox_center(bbox: dict, canvas_w: float, canvas_h: float) -> tuple[float, float]:
    x, y, w, h = pct_to_inches(bbox, canvas_w, canvas_h)
    return x + w / 2, y + h / 2


def _arrow_endpoints(source: dict, target: dict, canvas_w: float, canvas_h: float) -> tuple[float, float, float, float]:
    sbox = source["bbox_percent"]
    tbox = target["bbox_percent"]
    sx, sy = _bbox_center(sbox, canvas_w, canvas_h)
    tx, ty = _bbox_center(tbox, canvas_w, canvas_h)
    s_left, s_top, s_w, s_h = pct_to_inches(sbox, canvas_w, canvas_h)
    t_left, t_top, t_w, t_h = pct_to_inches(tbox, canvas_w, canvas_h)

    dx = tx - sx
    dy = ty - sy
    if abs(dx) >= abs(dy):
        if dx >= 0:
            return s_left + s_w + 0.03, sy, t_left - 0.03, ty
        return s_left - 0.03, sy, t_left + t_w + 0.03, ty
    if dy >= 0:
        return sx, s_top + s_h + 0.03, tx, t_top - 0.03
    return sx, s_top - 0.03, tx, t_top + t_h + 0.03


def _arrow_points_from_path(path: list, width_in: float, height_in: float) -> list[tuple[float, float]]:
    points = []
    for point in path:
        if isinstance(point, list) and len(point) >= 2:
            points.append((float(point[0]) * width_in, float(point[1]) * height_in))
    return points


def _draw_program_arrows(slide, program: dict, width_in: float, height_in: float) -> list[dict]:
    objects_by_id = _object_map(program)
    style = program.get("style", {}) if isinstance(program.get("style"), dict) else {}
    token_map = {str(item.get("token_id")): item for item in style.get("color_tokens", []) if isinstance(item, dict)}
    rendered: list[dict] = []
    for arrow in program.get("arrows", []):
        if arrow.get("type") == "custom_bus":
            continue
        source = objects_by_id.get(arrow.get("source") or arrow.get("source_id"))
        target = objects_by_id.get(arrow.get("target") or arrow.get("target_id"))
        path = arrow.get("path_percent") if isinstance(arrow.get("path_percent"), list) else []
        if len(path) >= 2 and all(isinstance(point, list) and len(point) >= 2 for point in path):
            points = _arrow_points_from_path(path, width_in, height_in)
        elif source and target:
            x1, y1, x2, y2 = _arrow_endpoints(source, target, width_in, height_in)
            points = [(x1, y1), (x2, y2)]
        else:
            continue
        token = token_map.get(str(arrow.get("style_token_id")))
        arrow_color = str(token.get("hex")) if token else str(arrow.get("stroke_color") or "#1F6F8B")
        control_kind = str(arrow.get("control_kind") or arrow.get("type", "")).lower()
        dashed = control_kind in {"dashed_loop", "dashed", "loop"}
        if str(arrow.get("line_pattern", "")).lower() in {"dash", "dashed"}:
            dashed = True
        line_width = float(arrow.get("stroke_width_pt") or style.get("arrow_weight_pt") or 1.7)
        route_style = str(arrow.get("route_style") or "")
        connector_type = _connector_type_for_route(route_style)
        halo_width = float(arrow.get("halo_width_pt") or 0.0)
        halo_color = str(arrow.get("halo_color") or "#FFFFFF")
        arrowhead_size = str(arrow.get("arrowhead_size") or "sm").lower()
        if arrowhead_size not in {"sm", "med", "lg"}:
            arrowhead_size = "sm"
        segment_count = 0
        for idx, ((x1, y1), (x2, y2)) in enumerate(zip(points[:-1], points[1:])):
            if halo_width > 0:
                halo = slide.shapes.add_connector(connector_type, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
                halo.line.color.rgb = _rgb(halo_color)
                halo.line.width = Pt(max(line_width + 1.2, halo_width))
                _apply_line_cap(halo, str(arrow.get("line_cap") or "round"))
                if dashed:
                    halo.line.dash_style = MSO_LINE_DASH_STYLE.DASH
            connector = slide.shapes.add_connector(connector_type, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
            connector.line.color.rgb = _rgb(arrow_color)
            connector.line.width = Pt(line_width)
            _apply_line_cap(connector, str(arrow.get("line_cap") or "round"))
            if dashed:
                connector.line.dash_style = MSO_LINE_DASH_STYLE.DASH
            if idx == len(points) - 2:
                _apply_arrow(connector, arrowhead_size)
            segment_count += 1
        rendered.append({
            "arrow_id": arrow.get("id"),
            "control_kind": control_kind or "straight_arrow",
            "semantic_role": arrow.get("semantic_role"),
            "route_style": route_style,
            "bundle_id": arrow.get("bundle_id"),
            "lane_index": arrow.get("lane_index"),
            "line_cap": arrow.get("line_cap", "round"),
            "line_pattern": "dash" if dashed else "solid",
            "connector_type": "curve" if connector_type == MSO_CONNECTOR.CURVE else "straight",
            "halo_width_pt": halo_width,
            "halo_color": halo_color if halo_width > 0 else None,
            "stroke_width_pt": line_width,
            "arrowhead_size": arrowhead_size,
            "routing_algorithm": arrow.get("routing_algorithm"),
            "route_generation_status": arrow.get("route_generation_status"),
            "segment_count": segment_count,
            "point_count": len(points),
            "editable_in": "pptx",
            "render_policy": "ppt_shape_not_image_asset",
            "status": "ok" if segment_count else "not_rendered",
        })
    return rendered


def _ocean_label(text: str) -> str:
    low = text.lower()
    if "openness" in low:
        return "Openness"
    if "conscientiousness" in low:
        return "Conscientiousness"
    if "extraversion" in low:
        return "Extraversion"
    if "agreeableness" in low:
        return "Agreeableness"
    if "neuroticism" in low:
        return "Neuroticism"
    return _short_caption(text)


def compile_ppt(program: dict, out_dir: str | Path) -> Path:
    out = Path(out_dir)
    canvas = program["canvas"]
    width_in = float(canvas["width_in"])
    height_in = float(canvas["height_in"])
    style = program.get("style", {})
    palette = style.get("palette") or style.get("reference_palette") or ["#2D6FB7", "#E17721", "#6B57C8", "#1B9A94", "#4B9B52", "#D44E5D"]
    panel_styles = style.get("panel_styles", {}) if isinstance(style.get("panel_styles"), dict) else {}

    prs = Presentation()
    prs.slide_width = Inches(width_in)
    prs.slide_height = Inches(height_in)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_title_block(slide, program, width_in, height_in)

    panels_by_id = _panel_map(program)
    panel_shapes = {}
    has_text_program = bool(_text_program_items(program))
    ocr_panel_title_targets = {
        str(item.get("target_id"))
        for item in _text_program_items(program)
        if str(item.get("role")) == "panel_title" and "ocr" in str(item.get("reference_binding") or item.get("fit_strategy") or "").lower()
    }
    for idx, panel in enumerate(program["panels"]):
        x, y, w, h = pct_to_inches(panel["bbox_percent"], width_in, height_in)
        local_style = panel_styles.get(panel["id"], {}) if isinstance(panel_styles.get(panel["id"], {}), dict) else {}
        fill = local_style.get("fill_color") or palette[idx % len(palette)]
        stroke = local_style.get("stroke_color") or palette[(idx + 1) % len(palette)]
        header_color = local_style.get("header_color") or stroke
        shape = _add_round_rect(slide, x, y, w, h, fill, stroke, width_pt=1.4)
        panel_shapes[panel["id"]] = shape
        header_h = min(0.34, h * 0.17)
        header = _add_round_rect(slide, x, y, w, header_h, header_color, header_color, width_pt=0.8)
        title_text_item = _text_item_for_target(program, panel["id"], "panel_title")
        _set_text(
            header,
            "" if str(panel["id"]) in ocr_panel_title_targets else panel["title"],
            font_size=float(title_text_item.get("font_size_pt")) if title_text_item else (9 if w < 2.0 else 10),
            bold=True,
            color=str(title_text_item.get("color_hex")) if title_text_item else "#FFFFFF",
            font_family=str(title_text_item.get("font_family_guess") or "") if title_text_item else None,
        )

    rendered_cards = []
    for card in sorted(program.get("cards", []), key=lambda item: int(item.get("z_index", 12))):
        if not isinstance(card, dict) or not isinstance(card.get("bbox_percent"), dict):
            continue
        _add_card_frame(slide, card, width_in, height_in)
        x, y, w, h = pct_to_inches(card["bbox_percent"], width_in, height_in)
        if card.get("title") and not has_text_program:
            _add_label(slide, str(card.get("title")), x + 0.03, y + 0.03, max(0.05, w - 0.06), min(0.22, h * 0.32), font_size=7, bold=True)
        rendered_cards.append({
            "card_id": card.get("id"),
            "semantic_role": card.get("semantic_role"),
            "shape_kind": card.get("shape_kind", "rounded_rect"),
            "bbox_percent": card.get("bbox_percent"),
            "dash_style": card.get("dash_style", "solid"),
            "fill_transparency": float(card.get("fill_transparency", 1.0)),
            "editable_in": "pptx",
            "render_policy": "ppt_shape_not_image_asset",
            "status": "ok",
        })

    rendered_arrows = _draw_program_arrows(slide, program, width_in, height_in)

    # Slot image layer and editable captions.
    composition_items = []
    caption_queue = []
    ordered_slots = sorted(program["slots"], key=lambda item: int(item.get("z_index", 20)))
    for slot in ordered_slots:
        x, y, w, h = pct_to_inches(slot["bbox_percent"], width_in, height_in)
        asset_path = out / "assets" / f"{slot['asset_id']}.png"
        tile_added = False
        if slot.get("composition_type") == "symbol_cutout":
            if asset_path.exists():
                _add_picture_contain(slide, asset_path, x, y, w, h)
                _left, _top, fit_w, fit_h, fill_percent = _picture_contain_box(asset_path, x, y, w, h)
            else:
                fill_percent = 0.0
            if bool(slot.get("show_slot_caption", False)) and not has_text_program:
                parent = panels_by_id.get(slot.get("panel_id"))
                if parent:
                    px, py, pw, ph = pct_to_inches(parent["bbox_percent"], width_in, height_in)
                    label_x = x + w + 0.035
                    label_w = max(0.28, min(0.86, px + pw - label_x - 0.03))
                    caption_queue.append((
                        _ocean_label(slot.get("display_label") or slot["paper_concept"]),
                        label_x,
                        y + h * 0.10,
                        label_w,
                        h * 0.80,
                        5 if label_w < 0.55 else 6,
                        PP_ALIGN.LEFT,
                    ))
            composition_items.append({
                "slot_id": slot["id"],
                "asset_id": slot["asset_id"],
                "slot_frame_policy": slot.get("slot_frame_policy", "frameless_slot"),
                "picture_fill_policy": slot.get("picture_fill_policy", "direct_full_slot_contain_no_tile"),
                "tile_frame_added": tile_added,
                "caption_inside_image_slot": False,
                "slot_bbox_percent": slot["bbox_percent"],
                "image_slot_area_fill_percent": round(fill_percent, 2),
                "status": "ok" if fill_percent >= 95 else "image_area_below_95",
            })
            continue
        if asset_path.exists():
            _add_picture_contain(slide, asset_path, x, y, w, h)
            _left, _top, fit_w, fit_h, fill_percent = _picture_contain_box(asset_path, x, y, w, h)
        else:
            fill_percent = 0.0
        if bool(slot.get("show_slot_caption", False)) and not has_text_program:
            caption_h = min(0.20, h * 0.18)
            caption = slot.get("display_label") or _short_caption(slot["paper_concept"])
            caption_queue.append((caption, x + w * 0.03, y + h + 0.01, w * 0.94, caption_h, 5 if w < 0.8 else 6, PP_ALIGN.CENTER))
        composition_items.append({
            "slot_id": slot["id"],
            "asset_id": slot["asset_id"],
            "slot_frame_policy": slot.get("slot_frame_policy", "frameless_slot"),
            "picture_fill_policy": slot.get("picture_fill_policy", "direct_full_slot_contain_no_tile"),
            "tile_frame_added": tile_added,
            "caption_inside_image_slot": False,
            "slot_bbox_percent": slot["bbox_percent"],
            "image_slot_area_fill_percent": round(fill_percent, 2),
            "status": "ok" if fill_percent >= 95 else "image_area_below_95",
        })

    for caption, x, y, w, h, font_size, align in caption_queue:
        _add_label(slide, caption, x, y, w, h, font_size=font_size, align=align)

    rendered_text_items = []
    for item in _text_program_items(program):
        is_ocr_text = "ocr" in str(item.get("reference_binding") or item.get("fit_strategy") or "").lower()
        if str(item.get("role")) == "panel_title" and not is_ocr_text:
            rendered_text_items.append({
                "text_id": item.get("id"),
                "role": item.get("role"),
                "target_id": item.get("target_id"),
                "rendered_as": "panel_header_text",
                "editable_in": "pptx",
            })
            continue
        bbox = item.get("bbox_percent")
        if not isinstance(bbox, dict):
            continue
        x, y, w, h = pct_to_inches(bbox, width_in, height_in)
        _add_label(
            slide,
            str(item.get("text") or ""),
            x,
            y,
            w,
            h,
            font_size=float(item.get("font_size_pt") or 6),
            bold=bool(item.get("bold")),
            align=_align_from_text_item(item),
            color=str(item.get("color_hex") or "#263747"),
            font_family=str(item.get("font_family_guess") or "") or None,
        )
        rendered_text_items.append({
            "text_id": item.get("id"),
            "role": item.get("role"),
            "target_id": item.get("target_id"),
            "source_reference_text_id": item.get("source_reference_text_id"),
            "bbox_percent": item.get("bbox_percent"),
            "font_size_pt": item.get("font_size_pt"),
            "color_hex": item.get("color_hex"),
            "font_family_guess": item.get("font_family_guess"),
            "fit_strategy": item.get("fit_strategy"),
            "ocr_confidence": item.get("ocr_confidence"),
            "editable_in": "pptx",
            "rendered_as": "ppt_textbox",
        })

    for label in program.get("labels", []):
        bbox = label.get("bbox_percent")
        if not isinstance(bbox, dict):
            continue
        x, y, w, h = pct_to_inches(bbox, width_in, height_in)
        _add_label(
            slide,
            str(label.get("text") or ""),
            x,
            y,
            w,
            h,
            font_size=float(label.get("font_size_pt") or 9),
            bold=bool(label.get("bold", True)),
            align=PP_ALIGN.LEFT if str(label.get("align")).lower() == "left" else PP_ALIGN.CENTER,
            color=str(label.get("color_hex") or "#263747"),
        )

    # Shared resource bus as editable connectors.
    shared = panels_by_id.get("shared_resource_library")
    if shared:
        sx, sy, sw, sh = pct_to_inches(shared["bbox_percent"], width_in, height_in)
        bus_y = sy - 0.30
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(sx + 0.1), Inches(bus_y), Inches(sx + sw - 0.1), Inches(bus_y))
        line.line.color.rgb = _rgb("#1F6F8B")
        line.line.width = Pt(1.2)
        for panel in program["panels"]:
            if panel["id"] == "shared_resource_library":
                continue
            cx, _cy = _bbox_center(panel["bbox_percent"], width_in, height_in)
            if not (sx + 0.1 <= cx <= sx + sw - 0.1):
                continue
            down = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(cx), Inches(bus_y), Inches(cx), Inches(sy))
            down.line.color.rgb = _rgb("#1F6F8B")
            down.line.width = Pt(1.0)
        if not has_text_program:
            _add_label(slide, "shared resources", sx + 0.2, bus_y - 0.22, 1.5, 0.18, font_size=7, align=PP_ALIGN.LEFT)

    if program.get("show_visible_title"):
        title = program.get("paper_brief", {}).get("title_guess") or "Research System Figure"
        _add_label(slide, title[:90], width_in - 5.2, height_in - 0.50, 4.9, 0.38, font_size=7, bold=True, align=PP_ALIGN.RIGHT)

    pptx_path = out / "editable_composition.pptx"
    prs.save(pptx_path)
    write_json(out / "composition_quality_report.json", {
        "summary": "PPT composition quality report checking frameless slot insertion, no extra white tiles, and image area fill inside each reference slot.",
        "policy": {
            "slot_frame_policy": "frameless_slot",
            "picture_fill_policy": "direct_full_slot_contain_no_tile",
            "min_image_slot_area_fill_percent": 95,
            "caption_inside_image_slot_allowed": False,
        },
        "slots": composition_items,
        "cards": rendered_cards,
        "arrows": rendered_arrows,
        "text": rendered_text_items,
    })
    return pptx_path
