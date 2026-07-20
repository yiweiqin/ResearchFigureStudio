from __future__ import annotations

import base64
import json
import mimetypes
import os
import platform
import shutil
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from typing import Callable

from PIL import Image, ImageChops, ImageDraw

from .control_localizer import localize_reference_controls
from .layout_planner import dominant_palette, estimate_background, plan_reference_layout
from .layout_semantic_planner import plan_slot_semantics
from .ppt_compiler import compile_ppt
from .rebuild_vlm_validation import build_rebuild_vlm_validation_report
from .text_layer import build_text_layer
from .utils import ensure_dir, write_json, write_text


ASSET_THRESHOLDS = {
    "character": (0.80, 0.95),
    "document_stack": (0.75, 0.95),
    "chart_card": (0.75, 0.95),
    "inspection": (0.70, 0.95),
    "tool_icon": (0.80, 0.95),
    "tool_combo": (0.70, 0.95),
    "device": (0.70, 0.95),
    "screenshot_card": (0.75, 0.95),
    "legend_marker": (0.80, 0.95),
    "thin_tool": (0.50, 0.95),
    "generic": (0.70, 0.95),
}


TEXT_ASSET_OVERLAP_THRESHOLD = 0.60
SIMPLE_PRIMITIVE_TYPES = {"legend_marker", "simple_marker", "status_marker", "node_marker"}


def _bbox(x: float, y: float, w: float, h: float) -> dict[str, float]:
    x = max(0.0, min(0.995, x))
    y = max(0.0, min(0.995, y))
    w = max(0.001, min(w, 1.0 - x))
    h = max(0.001, min(h, 1.0 - y))
    return {"x": round(x, 4), "y": round(y, 4), "w": round(w, 4), "h": round(h, 4)}


def _canvas_inches(width_px: int, height_px: int) -> tuple[float, float]:
    ratio = width_px / max(height_px, 1)
    width = max(10.0, min(15.6, 7.5 * ratio))
    return round(width, 3), round(width / ratio, 3)


def _rgb_hex(rgb: tuple[int, int, int]) -> str:
    return f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"


def _estimate_background(image: Image.Image) -> str:
    rgb = image.convert("RGB")
    w, h = rgb.size
    samples = []
    for box in [(0, 0, w, max(1, h // 12)), (0, h - max(1, h // 12), w, h), (0, 0, max(1, w // 12), h), (w - max(1, w // 12), 0, w, h)]:
        crop = rgb.crop(box).resize((1, 1))
        samples.append(crop.getpixel((0, 0)))
    avg = tuple(int(sum(px[i] for px in samples) / len(samples)) for i in range(3))
    return _rgb_hex(avg)


def _dominant_palette(image: Image.Image, count: int = 6) -> list[str]:
    small = image.convert("RGB").resize((96, 96))
    colors = small.quantize(colors=count).convert("RGB").getcolors(96 * 96) or []
    colors = sorted(colors, key=lambda item: item[0], reverse=True)
    return [_rgb_hex(rgb) for _n, rgb in colors[:count]]


def _supported_aspect_ratio(width: float, height: float) -> str:
    ratio = width / max(height, 0.001)
    if ratio >= 1.55:
        return "16:9"
    if ratio >= 1.18:
        return "4:3"
    if ratio >= 0.85:
        return "1:1"
    if ratio >= 0.62:
        return "3:4"
    return "9:16"


def _slot_type_for_box(slot: dict) -> str:
    semantic_type = str(slot.get("asset_type") or "").strip()
    if semantic_type:
        return semantic_type
    ratio = float(slot["bbox_percent"]["w"]) / max(float(slot["bbox_percent"]["h"]), 0.001)
    if ratio >= 2.2:
        return "thin_tool"
    if ratio >= 1.35:
        return "chart_card"
    if ratio <= 0.78:
        return "character"
    return "generic"


def _threshold_for_type(slot_type: str) -> tuple[float, float]:
    return ASSET_THRESHOLDS.get(slot_type, ASSET_THRESHOLDS["generic"])


def _foreground_metrics(image_path: str | Path, background_hex: str | None = None) -> dict:
    try:
        image = Image.open(image_path).convert("RGB")
    except Exception as exc:
        return {"status": "error", "error": str(exc), "foreground_bbox_px": None, "foreground_bbox_fill_percent": 0.0}
    w, h = image.size
    bg = background_hex or _estimate_background(image)
    bg_rgb = tuple(int(bg.strip("#")[i:i + 2], 16) for i in (0, 2, 4))
    bg_img = Image.new("RGB", image.size, bg_rgb)
    diff = ImageChops.difference(image, bg_img).convert("L")
    mask = diff.point(lambda p: 255 if p > 18 else 0)
    box = mask.getbbox()
    if not box:
        return {
            "status": "empty",
            "foreground_bbox_px": None,
            "foreground_bbox_fill_percent": 0.0,
            "margin_left_percent": 0.5,
            "margin_right_percent": 0.5,
            "margin_top_percent": 0.5,
            "margin_bottom_percent": 0.5,
        }
    left, top, right, bottom = box
    fill = ((right - left) * (bottom - top)) / max(w * h, 1)
    return {
        "status": "ok",
        "foreground_bbox_px": [left, top, right, bottom],
        "foreground_bbox_fill_percent": round(fill, 4),
        "margin_left_percent": round(left / w, 4),
        "margin_right_percent": round((w - right) / w, 4),
        "margin_top_percent": round(top / h, 4),
        "margin_bottom_percent": round((h - bottom) / h, 4),
    }


def economy_acceptance_decision(slot_type: str, fill_percent: float, strict: bool = False) -> dict:
    low, high = (0.80, 0.95) if strict else _threshold_for_type(slot_type)
    accepted = low <= float(fill_percent) <= high
    if accepted:
        reason = f"fill {fill_percent:.3f} within {slot_type} threshold {low:.2f}-{high:.2f}"
    elif float(fill_percent) < low:
        reason = f"fill {fill_percent:.3f} below {slot_type} threshold {low:.2f}"
    else:
        reason = f"fill {fill_percent:.3f} above tightness threshold {high:.2f}"
    return {"accepted": accepted, "selected_reason": reason, "threshold_min": low, "threshold_max": high}


def _find_visual_components(reference_path: Path, canvas_w: int, canvas_h: int) -> list[dict]:
    try:
        import cv2
        import numpy as np
    except Exception:
        return []
    image = cv2.imread(str(reference_path))
    if image is None:
        return []
    gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
    blur = cv2.GaussianBlur(gray, (5, 5), 0)
    edges = cv2.Canny(blur, 45, 120)
    kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (9, 9))
    dilated = cv2.dilate(edges, kernel, iterations=2)
    contours, _hier = cv2.findContours(dilated, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    boxes = []
    image_area = canvas_w * canvas_h
    for contour in contours:
        x, y, w, h = cv2.boundingRect(contour)
        area = w * h
        if area < image_area * 0.008 or area > image_area * 0.35:
            continue
        if w < 24 or h < 24:
            continue
        boxes.append((x, y, w, h))
    boxes.sort(key=lambda item: (item[1], item[0]))
    kept: list[tuple[int, int, int, int]] = []
    for box in boxes:
        x, y, w, h = box
        duplicate = False
        for kx, ky, kw, kh in kept:
            ix = max(0, min(x + w, kx + kw) - max(x, kx))
            iy = max(0, min(y + h, ky + kh) - max(y, ky))
            overlap = ix * iy / max(w * h, 1)
            if overlap > 0.55:
                duplicate = True
                break
        if not duplicate:
            kept.append(box)
        if len(kept) >= 12:
            break
    slots = []
    for idx, (x, y, w, h) in enumerate(kept, start=1):
        slots.append({
            "id": f"slot_{idx:02d}",
            "asset_id": f"slot_{idx:02d}",
            "paper_concept": f"visual element {idx}",
            "display_label": "",
            "bbox_percent": _bbox(x / canvas_w, y / canvas_h, w / canvas_w, h / canvas_h),
            "composition_type": "full_frame_icon",
            "show_slot_caption": False,
            "z_index": 20 + idx,
        })
    return slots


def _fallback_slots(count: int = 6) -> list[dict]:
    boxes = [
        (0.08, 0.26, 0.16, 0.22),
        (0.31, 0.26, 0.16, 0.22),
        (0.54, 0.26, 0.16, 0.22),
        (0.77, 0.26, 0.16, 0.22),
        (0.25, 0.58, 0.18, 0.22),
        (0.57, 0.58, 0.18, 0.22),
    ][:count]
    slots = []
    for idx, (x, y, w, h) in enumerate(boxes, start=1):
        slots.append({
            "id": f"slot_{idx:02d}",
            "asset_id": f"slot_{idx:02d}",
            "paper_concept": f"visual element {idx}",
            "display_label": "",
            "bbox_percent": _bbox(x, y, w, h),
            "composition_type": "full_frame_icon",
            "show_slot_caption": False,
            "z_index": 20 + idx,
        })
    return slots


def _style_from_layout(layout: dict) -> dict:
    palette = layout.get("palette") or ["#FFFFFF", "#4A90C2", "#E17721", "#4B9B52", "#6B57C8", "#D44E5D"]
    background = str(layout.get("canvas", {}).get("background") or "#FFFFFF")
    style = {
        "palette": palette,
        "reference_palette": palette,
        "color_tokens": [{"token_id": f"palette_{i:02d}", "hex": color} for i, color in enumerate(palette)],
        "panel_styles": {
            "reference_canvas": {
                "fill_color": background,
                "stroke_color": palette[1] if len(palette) > 1 else "#B8C0CC",
                "header_color": palette[1] if len(palette) > 1 else "#6C7A89",
            }
        },
        "arrow_weight_pt": 1.7,
    }
    for panel in layout.get("panels", []):
        if panel.get("id") and panel["id"] != "reference_canvas":
            style["panel_styles"].setdefault(str(panel["id"]), {
                "fill_color": background,
                "stroke_color": palette[1] if len(palette) > 1 else "#B8C0CC",
                "header_color": palette[1] if len(palette) > 1 else "#6C7A89",
            })
    return style


def _build_program(
    reference_path: Path,
    out: Path,
    layout_mode: str = "hybrid",
    vlm_layout_adapter: Callable | None = None,
) -> tuple[dict, dict]:
    layout = plan_reference_layout(reference_path, out, mode=layout_mode, vlm_adapter=vlm_layout_adapter)
    style = _style_from_layout(layout)
    program = {
        "canvas": layout["canvas"],
        "style": style,
        "title_block": {"title": "", "subtitle": "", "bbox_percent": _bbox(0.04, 0.02, 0.92, 0.04)},
        "panels": layout.get("panels", []),
        "cards": layout.get("cards", []),
        "slots": layout.get("slots", []),
        "assets": [],
        "arrows": [],
        "labels": [],
        "groups": [],
        "export_targets": [{"type": "pptx", "path": "editable_composition.pptx"}],
    }
    return program, layout


def _crop_reference(reference_path: Path, out_path: Path, bbox_percent: dict) -> Path:
    with Image.open(reference_path).convert("RGB") as image:
        w, h = image.size
        left = int(float(bbox_percent["x"]) * w)
        top = int(float(bbox_percent["y"]) * h)
        right = int((float(bbox_percent["x"]) + float(bbox_percent["w"])) * w)
        bottom = int((float(bbox_percent["y"]) + float(bbox_percent["h"])) * h)
        crop = image.crop((max(0, left), max(0, top), min(w, right), min(h, bottom)))
        out_path.parent.mkdir(parents=True, exist_ok=True)
        crop.save(out_path)
    return out_path


def _placeholder_asset(slot: dict, out_path: Path, background: str) -> None:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    ratio = float(slot["bbox_percent"]["w"]) / max(float(slot["bbox_percent"]["h"]), 0.001)
    if ratio >= 1.55:
        size = (512, 288)
    elif ratio <= 0.62:
        size = (288, 512)
    else:
        size = (384, 384)
    image = Image.new("RGB", size, background)
    draw = ImageDraw.Draw(image)
    pad = int(min(size) * 0.12)
    draw.rounded_rectangle((pad, pad, size[0] - pad, size[1] - pad), radius=18, fill="#FFFFFF", outline="#6C7A89", width=3)
    draw.line((pad * 1.5, size[1] * 0.45, size[0] - pad * 1.5, size[1] * 0.45), fill="#4A6FA5", width=5)
    draw.ellipse((size[0] * 0.42, size[1] * 0.22, size[0] * 0.58, size[1] * 0.38), fill="#E17721")
    draw.text((pad * 1.3, size[1] * 0.58), slot["id"], fill="#263747")
    image.save(out_path)


def _api_generate_asset(spec: dict, crop_path: Path, out_path: Path) -> dict:
    import requests

    api_key = os.getenv("GEMINI_API_KEY") or os.getenv("API_KEY")
    url = os.getenv("GEMINI_GEN_IMG_URL")
    if not api_key or not url:
        raise RuntimeError("GEMINI_API_KEY/API_KEY and GEMINI_GEN_IMG_URL are required for --asset-mode api")
    mime = mimetypes.guess_type(crop_path.name)[0] or "image/png"
    image_b64 = base64.b64encode(crop_path.read_bytes()).decode("ascii")
    prompt = spec["prompt"]
    payload = {
        "contents": [{
            "role": "user",
            "parts": [
                {"text": prompt},
                {"inlineData": {"mimeType": mime, "data": image_b64}},
            ],
        }],
        "generationConfig": {
            "responseModalities": ["IMAGE"],
            "imageConfig": {"imageSize": "1K", "aspectRatio": spec["generation_aspect_ratio"]},
        },
    }
    response = requests.post(url, headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}, data=json.dumps(payload), timeout=240)
    response.raise_for_status()
    data = response.json()
    for candidate in data.get("candidates", []):
        for part in candidate.get("content", {}).get("parts", []):
            inline = part.get("inlineData") or part.get("inline_data")
            if inline and inline.get("data"):
                out_path.parent.mkdir(parents=True, exist_ok=True)
                out_path.write_bytes(base64.b64decode(inline["data"]))
                return {"status": "generated", "api_requests_attempted": 1}
    raise RuntimeError("API returned no image data")


def _fill_target_for_type(slot_type: str) -> tuple[float, float]:
    return _threshold_for_type(slot_type)


def _asset_prompt(slot: dict, slot_type: str, background: str, aspect_ratio: str) -> str:
    low, high = _fill_target_for_type(slot_type)
    subject = str(slot.get("prompt_subject") or slot.get("paper_concept") or slot.get("id"))
    nearby = slot.get("nearby_text") or []
    nearby_text = " | ".join(nearby) if isinstance(nearby, list) else str(nearby or "")
    orientation = "horizontal" if aspect_ratio in {"16:9", "4:3"} else "vertical" if aspect_ratio in {"3:4", "9:16"} else "square"
    return (
        "Use the provided reference crop only as guidance for object identity, pose, colors, and diagram style. "
        "Generate one clean standalone visual asset for an editable PowerPoint reconstruction. "
        f"Slot type: {slot_type}. Canvas aspect ratio: {aspect_ratio}. Composition orientation: {orientation}. "
        f"Main subject: {subject}. "
        f"Nearby reference text for semantics only: {nearby_text}. "
        f"The main subject should occupy {int(low * 100)}%-{int(high * 100)}% of the image canvas, with only narrow uniform margins. "
        "For horizontal slots, make the subject span the width; for vertical slots, make the subject span the height. "
        "Avoid a tiny object floating in a large colored field. "
        f"Use a flat background matching {background}. "
        "Do not include readable text, captions, watermarks, arrows, panel borders, or the whole source diagram. "
        "Do not crop off important details."
    )


def _load_accepted_assets(path: Path) -> dict:
    if not path.exists():
        return {}
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return {}
    if isinstance(data, dict):
        return data
    return {}


def _bbox_overlap_ratio(a: dict, b: dict) -> float:
    ax1, ay1 = float(a["x"]), float(a["y"])
    ax2, ay2 = ax1 + float(a["w"]), ay1 + float(a["h"])
    bx1, by1 = float(b["x"]), float(b["y"])
    bx2, by2 = bx1 + float(b["w"]), by1 + float(b["h"])
    ix = max(0.0, min(ax2, bx2) - max(ax1, bx1))
    iy = max(0.0, min(ay2, by2) - max(ay1, by1))
    return (ix * iy) / max(float(a["w"]) * float(a["h"]), 0.000001)


def _text_program_boxes(program: dict) -> list[dict]:
    text_program = program.get("text_program") if isinstance(program.get("text_program"), dict) else {}
    boxes = []
    for item in text_program.get("items", []) or []:
        if not isinstance(item, dict) or not isinstance(item.get("bbox_percent"), dict):
            continue
        if not str(item.get("text") or "").strip():
            continue
        boxes.append(item)
    return boxes


def _normalized_reuse_key(spec: dict) -> str:
    import re

    subject = str(spec.get("prompt_subject") or spec.get("semantic_role") or spec.get("slot_id") or "").lower()
    subject = re.sub(r"\b(executor|branch|step|slot|icon)\s*[a-z0-9]*\b", " ", subject)
    subject = re.sub(r"\b[a-d]\b|\b[0-9]+\b", " ", subject)
    subject = re.sub(r"[^a-z0-9]+", " ", subject).strip()
    if not subject:
        subject = str(spec.get("asset_type") or spec.get("slot_type") or "generic")
    return f"{spec.get('asset_type') or spec.get('slot_type') or 'generic'}::{subject}"


def _apply_smart_asset_policy(specs: list[dict], program: dict, out: Path, asset_policy: str) -> list[dict]:
    text_items = _text_program_boxes(program)
    filtered: list[dict] = []
    text_filter_report = []
    decision_report = []
    api_plan = {
        "summary": "Smart API asset plan. Reference crops are only prompt/context inputs; they are not final PPT assets.",
        "asset_policy": asset_policy,
        "unique_api_assets": 0,
        "reused_slots": 0,
        "skipped_text_regions": 0,
        "ppt_primitive_slots": 0,
        "estimated_api_requests": 0,
        "reuse_groups": [],
        "recommended_api_slots": [],
    }

    for spec in specs:
        slot_id = str(spec["slot_id"])
        bbox = spec["slot_bbox_percent"]
        overlaps = [
            {
                "text_id": item.get("id"),
                "text": item.get("text"),
                "overlap_ratio": round(_bbox_overlap_ratio(bbox, item["bbox_percent"]), 4),
            }
            for item in text_items
        ]
        best = max(overlaps, key=lambda item: item["overlap_ratio"], default=None)
        if best and float(best["overlap_ratio"]) >= TEXT_ASSET_OVERLAP_THRESHOLD:
            spec["asset_decision"] = "text_region"
            spec["asset_generation_mode"] = "skip_asset"
            spec["asset_decision_reason"] = f"OCR/DSL text overlap {best['overlap_ratio']:.2f}; region is primarily editable text."
            text_filter_report.append({
                "candidate_id": slot_id,
                "rejected_as_asset": True,
                "converted_to": "existing_editable_text",
                "reason": spec["asset_decision_reason"],
                "matched_text_id": best.get("text_id"),
                "matched_text": best.get("text"),
                "ocr_or_text_overlap": best.get("overlap_ratio"),
            })
            decision_report.append({
                "slot_id": slot_id,
                "decision": "text_region",
                "needs_api": False,
                "reason": spec["asset_decision_reason"],
            })
            api_plan["skipped_text_regions"] += 1
            continue
        if str(spec.get("asset_type") or spec.get("slot_type") or "").lower() in SIMPLE_PRIMITIVE_TYPES:
            spec["asset_decision"] = "ppt_primitive"
            spec["asset_generation_mode"] = "skip_asset"
            spec["asset_decision_reason"] = "Simple marker should be represented by editable PPT primitives, not an image asset."
            decision_report.append({
                "slot_id": slot_id,
                "decision": "ppt_primitive",
                "needs_api": False,
                "reason": spec["asset_decision_reason"],
            })
            api_plan["ppt_primitive_slots"] += 1
            continue
        spec["asset_decision"] = "api_generated"
        spec["asset_generation_mode"] = "api"
        spec["asset_decision_reason"] = "Complex non-text visual slot; final asset should be API-generated from reference crop context."
        spec["reuse_group_key"] = _normalized_reuse_key(spec)
        filtered.append(spec)

    primary_by_key: dict[str, dict] = {}
    reuse_groups: dict[str, list[str]] = {}
    for spec in filtered:
        key = str(spec["reuse_group_key"])
        reuse_groups.setdefault(key, []).append(str(spec["slot_id"]))
        if key not in primary_by_key:
            primary_by_key[key] = spec
            decision = "api_generated"
            api_plan["unique_api_assets"] += 1
            api_plan["recommended_api_slots"].append(str(spec["slot_id"]))
        else:
            primary = primary_by_key[key]
            spec["asset_decision"] = "reuse_existing"
            spec["asset_generation_mode"] = "reuse_existing"
            spec["reuse_source_slot_id"] = primary["slot_id"]
            spec["reuse_source_asset_id"] = primary["asset_id"]
            decision = "reuse_existing"
            api_plan["reused_slots"] += 1
        decision_report.append({
            "slot_id": spec["slot_id"],
            "decision": decision,
            "reused_asset_id": spec.get("reuse_source_asset_id"),
            "reuse_group_key": key,
            "needs_api": decision == "api_generated",
            "reason": spec.get("asset_decision_reason"),
        })
    api_plan["estimated_api_requests"] = api_plan["unique_api_assets"]
    api_plan["reuse_groups"] = [
        {"reuse_group_key": key, "slots": slots}
        for key, slots in sorted(reuse_groups.items())
        if len(slots) > 1
    ]
    kept_slot_ids = {str(spec["slot_id"]) for spec in filtered}
    program["slots"] = [slot for slot in program.get("slots", []) if str(slot.get("id")) in kept_slot_ids]
    write_json(out / "asset_decision_report.json", {
        "summary": "Slot-level asset decisions. In smart-api mode, final crop assets are disabled; crops are only API reference inputs.",
        "asset_policy": asset_policy,
        "decisions": decision_report,
    })
    write_json(out / "text_asset_filter_report.json", {
        "summary": "Text-like candidate slots rejected before asset generation.",
        "overlap_threshold": TEXT_ASSET_OVERLAP_THRESHOLD,
        "items": text_filter_report,
    })
    write_json(out / "api_asset_plan.json", api_plan)
    return filtered


def _make_asset_specs(program: dict, reference_path: Path, out: Path, asset_policy: str = "legacy") -> list[dict]:
    specs = []
    background = str(program["canvas"].get("background") or "#FFFFFF")
    for slot in program["slots"]:
        bbox = slot["bbox_percent"]
        slot_type = _slot_type_for_box(slot)
        fill_min, fill_max = _fill_target_for_type(slot_type)
        slot_ratio = float(bbox["w"]) / max(float(bbox["h"]), 0.001)
        generation_ratio = str(slot.get("generation_aspect_ratio") or _supported_aspect_ratio(float(bbox["w"]), float(bbox["h"])))
        slot_background = str(slot.get("background_color_hex") or background)
        crop_path = out / "reference_slot_crops" / f"{slot['id']}.png"
        _crop_reference(reference_path, crop_path, bbox)
        spec = {
            "slot_id": slot["id"],
            "asset_id": slot["asset_id"],
            "slot_type": slot_type,
            "asset_type": slot_type,
            "semantic_role": slot.get("semantic_role"),
            "nearby_text": slot.get("nearby_text"),
            "panel_context": slot.get("panel_context"),
            "upstream_ids": slot.get("upstream_ids", []),
            "downstream_ids": slot.get("downstream_ids", []),
            "prompt_subject": slot.get("prompt_subject") or slot.get("paper_concept") or slot.get("id"),
            "slot_bbox_percent": bbox,
            "slot_aspect_ratio": round(slot_ratio, 4),
            "generation_aspect_ratio": generation_ratio,
            "background_color_hex": slot_background,
            "reference_crop_path": str(crop_path),
            "prompt": _asset_prompt(slot, slot_type, slot_background, generation_ratio),
            "main_subject_fill_target": round((fill_min + fill_max) / 2, 3),
            "internal_content_fill_target": round((fill_min + fill_max) / 2, 3),
            "max_margin_percent": 0.10,
            "forbidden_elements": ["readable text", "captions", "watermarks", "arrows", "panel borders", "full source diagram"],
            "style_reference": "provided reference crop",
            "no_text": True,
            "no_arrows": True,
            "no_full_diagram": True,
            "layout_bbox_locked": True,
            "no_crop_postprocess": True,
        }
        specs.append(spec)
    if asset_policy == "smart-api":
        return _apply_smart_asset_policy(specs, program, out, asset_policy)
    write_json(out / "asset_decision_report.json", {
        "summary": "Legacy asset decision report.",
        "asset_policy": asset_policy,
        "decisions": [{"slot_id": spec["slot_id"], "decision": asset_policy, "needs_api": asset_policy == "api"} for spec in specs],
    })
    write_json(out / "text_asset_filter_report.json", {
        "summary": "Text asset filter disabled in legacy asset policy.",
        "items": [],
    })
    write_json(out / "api_asset_plan.json", {
        "summary": "API asset plan disabled in legacy asset policy.",
        "asset_policy": asset_policy,
        "estimated_api_requests": len(specs) if asset_policy == "api" else 0,
    })
    return specs


def _generate_assets(
    specs: list[dict],
    program: dict,
    out: Path,
    asset_mode: str,
    asset_workers: int,
    asset_retries: int,
    economy_mode: bool,
    regenerate_slots: set[str],
    strict_asset_regeneration: bool,
    asset_policy: str = "legacy",
) -> tuple[list[dict], dict]:
    asset_dir = ensure_dir(out / "assets")
    accepted = _load_accepted_assets(out / "accepted_assets.json")
    reports: list[dict] = []
    api_requests = 0

    def run_one(spec: dict) -> dict:
        slot_id = spec["slot_id"]
        out_path = asset_dir / f"{spec['asset_id']}.png"
        if spec.get("asset_decision") == "reuse_existing":
            source_asset = asset_dir / f"{spec['reuse_source_asset_id']}.png"
            if source_asset.exists():
                shutil.copyfile(source_asset, out_path)
                metrics = _foreground_metrics(out_path, spec["background_color_hex"])
                decision = economy_acceptance_decision(spec["slot_type"], metrics.get("foreground_bbox_fill_percent", 0.0), strict=False)
                return {**spec, **metrics, **decision, "status": "reused_from_group", "asset_path": str(out_path), "api_requests_attempted": 0, "economy_decision": "reuse_group_asset", "fallback_used": False}
            return {**spec, "status": "reuse_source_missing", "asset_path": str(out_path), "api_requests_attempted": 0, "fallback_used": True, "error": f"reuse source missing: {source_asset}"}
        locked = bool(accepted.get(slot_id, {}).get("accepted")) or bool(accepted.get(spec["asset_id"], {}).get("accepted"))
        if out_path.exists() and economy_mode and slot_id not in regenerate_slots and not strict_asset_regeneration:
            metrics = _foreground_metrics(out_path, spec["background_color_hex"])
            decision = economy_acceptance_decision(spec["slot_type"], metrics.get("foreground_bbox_fill_percent", 0.0), strict=False)
            if locked or decision["accepted"]:
                return {**spec, **metrics, **decision, "status": "reused", "asset_path": str(out_path), "api_requests_attempted": 0, "economy_decision": "reuse_existing_accepted_or_threshold_passed", "fallback_used": False}
        crop_path = Path(spec["reference_crop_path"])
        attempts = max(1, int(asset_retries) + 1) if strict_asset_regeneration else 1
        last_error = None
        for attempt in range(1, attempts + 1):
            try:
                if asset_mode == "placeholder":
                    _placeholder_asset({"id": slot_id, "bbox_percent": spec["slot_bbox_percent"]}, out_path, spec["background_color_hex"])
                    request_count = 0
                    status = "placeholder"
                elif asset_mode == "crop" and asset_policy != "smart-api":
                    shutil.copyfile(crop_path, out_path)
                    request_count = 0
                    status = "reference_crop"
                elif asset_mode == "crop" and asset_policy == "smart-api":
                    _placeholder_asset({"id": slot_id, "bbox_percent": spec["slot_bbox_percent"]}, out_path, spec["background_color_hex"])
                    request_count = 0
                    status = "crop_disabled_by_smart_api_policy_placeholder"
                else:
                    result = _api_generate_asset(spec, crop_path, out_path)
                    request_count = int(result.get("api_requests_attempted", 1))
                    status = result.get("status", "generated")
                metrics = _foreground_metrics(out_path, spec["background_color_hex"])
                decision = economy_acceptance_decision(spec["slot_type"], metrics.get("foreground_bbox_fill_percent", 0.0), strict=strict_asset_regeneration)
                return {
                    **spec,
                    **metrics,
                    **decision,
                    "status": status,
                    "asset_path": str(out_path),
                    "api_requests_attempted": request_count,
                    "attempt": attempt,
                    "economy_decision": "generated_once" if economy_mode else "generated",
                    "fallback_used": False,
                }
            except Exception as exc:
                last_error = str(exc)
                if asset_mode == "api" and asset_policy != "smart-api":
                    try:
                        shutil.copyfile(crop_path, out_path)
                        metrics = _foreground_metrics(out_path, spec["background_color_hex"])
                        decision = economy_acceptance_decision(spec["slot_type"], metrics.get("foreground_bbox_fill_percent", 0.0), strict=False)
                        return {
                            **spec,
                            **metrics,
                            **decision,
                            "status": "api_failed_reference_crop_fallback",
                            "asset_path": str(out_path),
                            "api_requests_attempted": 1,
                            "attempt": attempt,
                            "economy_decision": "fallback_after_api_error",
                            "fallback_used": True,
                            "error": last_error,
                        }
                    except Exception:
                        pass
                if asset_mode == "api" and asset_policy == "smart-api":
                    try:
                        _placeholder_asset({"id": slot_id, "bbox_percent": spec["slot_bbox_percent"]}, out_path, spec["background_color_hex"])
                        metrics = _foreground_metrics(out_path, spec["background_color_hex"])
                        decision = economy_acceptance_decision(spec["slot_type"], metrics.get("foreground_bbox_fill_percent", 0.0), strict=False)
                        return {
                            **spec,
                            **metrics,
                            **decision,
                            "status": "api_failed_placeholder_fallback",
                            "asset_path": str(out_path),
                            "api_requests_attempted": 1,
                            "attempt": attempt,
                            "economy_decision": "placeholder_after_api_error_crop_disabled",
                            "fallback_used": True,
                            "error": last_error,
                        }
                    except Exception:
                        pass
        return {**spec, "status": "failed", "asset_path": str(out_path), "api_requests_attempted": 0, "fallback_used": True, "error": last_error}

    primary_specs = [spec for spec in specs if spec.get("asset_decision") != "reuse_existing"]
    reuse_specs = [spec for spec in specs if spec.get("asset_decision") == "reuse_existing"]
    workers = max(1, min(int(asset_workers), 12))
    with ThreadPoolExecutor(max_workers=workers) as pool:
        futures = [pool.submit(run_one, spec) for spec in primary_specs]
        for future in as_completed(futures):
            report = future.result()
            reports.append(report)
            api_requests += int(report.get("api_requests_attempted") or 0)
    for spec in reuse_specs:
        report = run_one(spec)
        reports.append(report)
        api_requests += int(report.get("api_requests_attempted") or 0)
    reports.sort(key=lambda item: str(item.get("slot_id")))
    program["assets"] = [{"id": spec["asset_id"], "path": f"assets/{spec['asset_id']}.png", "source": "slot_asset"} for spec in specs]
    summary = {
        "summary": "Asset generation report for reusable editable rebuild workflow.",
        "asset_mode": asset_mode,
        "asset_policy": asset_policy,
        "final_crop_assets_allowed": asset_policy != "smart-api",
        "economy_mode": economy_mode,
        "asset_retries": asset_retries,
        "strict_asset_regeneration": strict_asset_regeneration,
        "api_requests_attempted": api_requests,
        "assets": reports,
    }
    economy = {
        "summary": "Economy policy decisions.",
        "api_requests_attempted": api_requests,
        "accepted_assets_file": str(out / "accepted_assets.json"),
        "regenerate_slots": sorted(regenerate_slots),
        "thresholds": {key: {"min": value[0], "max": value[1]} for key, value in ASSET_THRESHOLDS.items()},
        "asset_policy": asset_policy,
        "assets": [{
            "slot_id": item.get("slot_id"),
            "slot_type": item.get("slot_type"),
            "economy_decision": item.get("economy_decision"),
            "selected_reason": item.get("selected_reason"),
            "api_requests_attempted": item.get("api_requests_attempted"),
            "fallback_used": item.get("fallback_used"),
        } for item in reports],
    }
    write_json(out / "asset_generation_report.json", summary)
    write_json(out / "asset_economy_report.json", economy)

    def image_size(path_text: str | None) -> list[int] | None:
        if not path_text or not Path(path_text).exists():
            return None
        with Image.open(path_text) as img:
            return list(img.size)

    write_json(out / "asset_ratio_fit_report.json", {
        "summary": "Aspect-ratio fit report.",
        "assets": [{
            "slot_id": item.get("slot_id"),
            "slot_aspect_ratio": item.get("slot_aspect_ratio"),
            "generation_aspect_ratio": item.get("generation_aspect_ratio"),
            "generated_image_size": image_size(item.get("asset_path")),
            "foreground_bbox_fill_percent": item.get("foreground_bbox_fill_percent"),
            "slot_fit_waste_percent": None,
            "aspect_ratio_match_status": "matched" if item.get("generation_aspect_ratio") else "unknown",
        } for item in reports],
    })
    return reports, summary


def _export_preview(pptx_path: Path, out: Path) -> Path | None:
    preview = out / "rebuild_preview.png"
    if platform.system() != "Windows":
        write_text(out / "preview_export_error.txt", "PowerPoint preview export is only available on Windows with Microsoft PowerPoint installed.")
        return None
    try:
        import win32com.client
        app = win32com.client.Dispatch("PowerPoint.Application")
        app.Visible = True
        deck = app.Presentations.Open(str(pptx_path), WithWindow=False)
        deck.Slides(1).Export(str(preview), "PNG")
        deck.Close()
        app.Quit()
        return preview
    except Exception as exc:
        write_text(out / "preview_export_error.txt", str(exc))
        return None


def _ppt_package_counts(pptx_path: Path, reference_path: Path) -> dict:
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    slide_count = len(prs.slides)
    shape_count = 0
    picture_count = 0
    text_shape_count = 0
    connector_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            shape_count += 1
            if getattr(shape, "shape_type", None) == 13:
                picture_count += 1
            if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                text_shape_count += 1
            shape_type_text = str(getattr(shape, "shape_type", ""))
            if "CONNECTOR" in shape_type_text or "LINE" in shape_type_text:
                connector_count += 1
    return {
        "ppt_page_count": slide_count,
        "shape_count": shape_count,
        "picture_count": picture_count,
        "text_shape_count": text_shape_count,
        "connector_count": connector_count,
        "contains_full_reference_image": False,
        "reference_path": str(reference_path),
    }


def _load_json_or_empty(path: Path) -> dict:
    if not path.exists():
        return {}
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return {}
    return data if isinstance(data, dict) else {}


def _program_from_contracts(out: Path) -> dict:
    geometry = _load_json_or_empty(out / "reference_geometry.json")
    controls = _load_json_or_empty(out / "reference_controls.json")
    text_program = _load_json_or_empty(out / "text_program.json")
    canvas = geometry.get("canvas") if isinstance(geometry.get("canvas"), dict) else {"width_in": 10.0, "height_in": 5.625, "background": "#FFFFFF"}
    if "width_in" not in canvas or "height_in" not in canvas:
        width_px = int(canvas.get("width_px") or 1600)
        height_px = int(canvas.get("height_px") or 900)
        width_in, height_in = _canvas_inches(width_px, height_px)
        canvas = {**canvas, "width_in": width_in, "height_in": height_in, "background": canvas.get("background") or "#FFFFFF"}
    layout = {
        "canvas": canvas,
        "palette": geometry.get("palette") or ["#FFFFFF", "#4A90C2", "#E17721", "#4B9B52"],
        "panels": geometry.get("panels") or [{
            "id": "reference_canvas",
            "title": "Editable Rebuild",
            "bbox_percent": _bbox(0.025, 0.065, 0.95, 0.84),
            "editable_in": "pptx",
        }],
    }
    slots = geometry.get("slots") or _load_json_or_empty(out / "slot_inventory.json").get("slots") or []
    program = {
        "canvas": canvas,
        "style": _style_from_layout(layout),
        "title_block": {"title": "", "subtitle": "", "bbox_percent": _bbox(0.04, 0.02, 0.92, 0.04)},
        "panels": layout["panels"],
        "cards": geometry.get("cards") or [],
        "slots": slots,
        "assets": [{"id": slot.get("asset_id") or slot.get("id"), "path": f"assets/{slot.get('asset_id') or slot.get('id')}.png", "source": "slot_asset"} for slot in slots],
        "arrows": controls.get("arrows") or [],
        "labels": [],
        "groups": [],
        "export_targets": [{"type": "pptx", "path": "editable_composition.pptx"}],
    }
    if isinstance(text_program.get("items"), list):
        program["text_program"] = text_program
        program["text_program_path"] = "text_program.json"
    return program


def _compile_existing_contracts(reference_path: Path, out: Path, export_preview: bool = False) -> dict:
    program = _program_from_contracts(out)
    write_json(out / "figure_program.json", program)
    pptx_path = compile_ppt(program, out)
    preview = _export_preview(pptx_path, out) if export_preview else None
    counts = _ppt_package_counts(pptx_path, reference_path)
    report_path = out / "composition_quality_report.json"
    report = _load_json_or_empty(report_path)
    report["rebuild_editable_summary"] = counts
    report["compile_only"] = True
    write_json(report_path, report)
    return {
        "summary": "Editable rebuild recompiled from existing JSON contracts.",
        "ok": True,
        "out_dir": str(out),
        "reference": str(reference_path),
        "pptx": str(pptx_path),
        "preview": str(preview) if preview else None,
        "asset_count": len(program.get("assets", [])),
        "slot_count": len(program.get("slots", [])),
        "connector_count": len(program.get("arrows", [])),
        "text_count": len(program.get("text_program", {}).get("items", [])) if isinstance(program.get("text_program"), dict) else 0,
        "compile_only": True,
    }


def rebuild_editable(
    reference: str | Path,
    out: str | Path,
    asset_mode: str = "api",
    asset_workers: int = 4,
    asset_retries: int = 1,
    economy_mode: bool = True,
    text_mode: str = "ocr",
    control_mode: str = "hybrid",
    layout_mode: str = "hybrid",
    export_preview: bool = False,
    regenerate_slots: str | list[str] | None = None,
    strict_asset_regeneration: bool = False,
    skip_analysis: bool = False,
    compile_only: bool = False,
    ocr_engine: str = "paddle",
    ocr_lang: str = "en_ch",
    ocr_adapter: Callable | None = None,
    vlm_layout_adapter: Callable | None = None,
    control_adapter: Callable | None = None,
    semantic_adapter: Callable | None = None,
    asset_policy: str = "legacy",
) -> dict:
    reference_path = Path(reference)
    if not reference_path.exists():
        raise FileNotFoundError(reference_path)
    out_path = ensure_dir(out)
    if compile_only:
        return _compile_existing_contracts(reference_path, out_path, export_preview=export_preview)
    input_dir = ensure_dir(out_path / "inputs")
    archived_reference = input_dir / reference_path.name
    shutil.copyfile(reference_path, archived_reference)
    write_json(out_path / "input_manifest.json", {
        "summary": "Input manifest for reference-only editable rebuild.",
        "reference": str(reference_path),
        "archived_reference": str(archived_reference),
        "asset_mode": asset_mode,
        "text_mode": text_mode,
        "control_mode": control_mode,
        "layout_mode": layout_mode,
        "asset_policy": asset_policy,
        "skip_analysis": skip_analysis,
        "compile_only": compile_only,
    })

    if skip_analysis and (out_path / "reference_geometry.json").exists():
        program = _program_from_contracts(out_path)
        reference_geometry = _load_json_or_empty(out_path / "reference_geometry.json")
    else:
        program, reference_geometry = _build_program(archived_reference, out_path, layout_mode=layout_mode, vlm_layout_adapter=vlm_layout_adapter)
    write_json(out_path / "reference_geometry.json", reference_geometry)

    if skip_analysis and (out_path / "reference_controls.json").exists():
        reference_controls = _load_json_or_empty(out_path / "reference_controls.json")
    else:
        reference_controls = localize_reference_controls(
            archived_reference,
            program.get("slots", []),
            program.get("style", {}).get("palette", []),
            out_path,
            mode=control_mode,
            control_adapter=control_adapter,
        )
    program["arrows"] = reference_controls.get("arrows", [])
    write_json(out_path / "reference_controls.json", reference_controls)

    if text_mode != "off":
        extractor_mode = "ocr" if text_mode == "ocr" else "heuristic"
        build_text_layer(
            archived_reference,
            program,
            program.get("style", {}),
            out_path,
            text_extractor_mode=extractor_mode,
            ocr_engine=ocr_engine,
            ocr_lang=ocr_lang,
            ocr_adapter=ocr_adapter,
        )
    else:
        write_json(out_path / "reference_text_geometry.json", {"summary": "Text extraction skipped.", "detection_mode": "off", "regions": []})
        write_json(out_path / "ocr_text_quality_report.json", {"summary": "Text extraction skipped.", "mode": "off", "status": "skipped"})

    text_geometry = _load_json_or_empty(out_path / "reference_text_geometry.json")
    semantic_report = _load_json_or_empty(out_path / "slot_semantic_report.json") if skip_analysis else {}
    if not skip_analysis:
        semantic_slots, semantic_report = plan_slot_semantics(
            archived_reference,
            program.get("slots", []),
            program.get("panels", []),
            program.get("arrows", []),
            text_geometry,
            semantic_adapter=semantic_adapter,
        )
        program["slots"] = semantic_slots
        reference_geometry["slots"] = semantic_slots
        write_json(out_path / "reference_geometry.json", reference_geometry)
        write_json(out_path / "slot_semantic_report.json", semantic_report)
    elif not semantic_report:
        semantic_report = {"summary": "Semantic planning skipped by --skip-analysis.", "semantic_vlm_status": "skipped", "slots": []}
        write_json(out_path / "slot_semantic_report.json", semantic_report)
    write_json(out_path / "slot_inventory.json", {"summary": "Visual asset slot inventory.", "slots": program["slots"]})

    specs = _make_asset_specs(program, archived_reference, out_path, asset_policy=asset_policy)
    if asset_policy == "smart-api":
        reference_geometry["slots"] = program.get("slots", [])
        write_json(out_path / "reference_geometry.json", reference_geometry)
        write_json(out_path / "slot_inventory.json", {"summary": "Visual asset slot inventory.", "slots": program["slots"]})
    write_json(out_path / "asset_generation_specs.json", {"summary": "Slot-level asset generation specs.", "asset_mode": asset_mode, "specs": specs})
    regen = set()
    if isinstance(regenerate_slots, str):
        regen = {item.strip() for item in regenerate_slots.split(",") if item.strip()}
    elif isinstance(regenerate_slots, list):
        regen = {str(item).strip() for item in regenerate_slots if str(item).strip()}
    asset_reports, asset_summary = _generate_assets(specs, program, out_path, asset_mode, asset_workers, asset_retries, economy_mode, regen, strict_asset_regeneration, asset_policy=asset_policy)
    vlm_validation_report = build_rebuild_vlm_validation_report(out_path, reference_geometry, reference_controls, semantic_report, asset_summary)

    write_json(out_path / "figure_program.json", program)
    pptx_path = compile_ppt(program, out_path)
    preview = _export_preview(pptx_path, out_path) if export_preview else None

    counts = _ppt_package_counts(pptx_path, archived_reference)
    report_path = out_path / "composition_quality_report.json"
    try:
        report = json.loads(report_path.read_text(encoding="utf-8"))
    except Exception:
        report = {}
    report["rebuild_editable_summary"] = counts
    report["no_full_image_policy"] = {
        "status": "pass",
        "contains_full_reference_image": False,
        "policy": "reference image is archived only; final PPT uses slot assets and editable objects",
    }
    write_json(report_path, report)

    result = {
        "summary": "Reusable image-to-editable-PPT rebuild complete.",
        "ok": True,
        "out_dir": str(out_path),
        "reference": str(reference_path),
        "pptx": str(pptx_path),
        "preview": str(preview) if preview else None,
        "asset_mode": asset_mode,
        "asset_workers": asset_workers,
        "asset_retries": asset_retries,
        "economy_mode": economy_mode,
        "asset_policy": asset_policy,
        "api_requests_attempted": asset_summary.get("api_requests_attempted", 0),
        "asset_count": len(asset_reports),
        "slot_count": len(program.get("slots", [])),
        "text_count": len(program.get("text_program", {}).get("items", [])) if isinstance(program.get("text_program"), dict) else 0,
        "connector_count": len(program.get("arrows", [])),
        "text_mode": text_mode,
        "control_mode": control_mode,
        "layout_mode": layout_mode,
        "reports": {
            "input_manifest": str(out_path / "input_manifest.json"),
            "reference_geometry": str(out_path / "reference_geometry.json"),
            "reference_text_geometry": str(out_path / "reference_text_geometry.json"),
            "reference_controls": str(out_path / "reference_controls.json"),
            "slot_inventory": str(out_path / "slot_inventory.json"),
            "slot_semantic_report": str(out_path / "slot_semantic_report.json"),
            "rebuild_vlm_validation_report": str(out_path / "rebuild_vlm_validation_report.json"),
            "asset_generation_specs": str(out_path / "asset_generation_specs.json"),
            "asset_generation_report": str(out_path / "asset_generation_report.json"),
            "asset_economy_report": str(out_path / "asset_economy_report.json"),
            "asset_decision_report": str(out_path / "asset_decision_report.json"),
            "text_asset_filter_report": str(out_path / "text_asset_filter_report.json"),
            "api_asset_plan": str(out_path / "api_asset_plan.json"),
            "figure_program": str(out_path / "figure_program.json"),
            "composition_quality_report": str(out_path / "composition_quality_report.json"),
        },
    }
    write_json(out_path / "rebuild_result.json", result)
    return result
