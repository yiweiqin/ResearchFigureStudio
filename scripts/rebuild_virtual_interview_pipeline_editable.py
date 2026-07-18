from __future__ import annotations

import argparse
import base64
import json
import mimetypes
import os
import time
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path

from PIL import Image
import requests
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
REFERENCE = Path(r"C:\Users\zhang\AppData\Local\Temp\codex-clipboard-ff76bc61-99eb-43af-a58c-c3179570fd8e.png")
OUT = ROOT / "output" / "virtual_interview_pipeline_editable_rebuild"
ASSET_DIR = OUT / "assets"
REFERENCE_CROP_DIR = OUT / "reference_slot_crops"
AI_CANDIDATE_DIR = OUT / "asset_candidates"

REF_W = 1774
REF_H = 887
SLIDE_W = 15.60
SLIDE_H = SLIDE_W * REF_H / REF_W
FONT = "Arial"


ASSET_SPECS = [
    {
        "id": "setup_top",
        "bbox": (22, 120, 238, 280),
        "prompt": "participant sitting at a desk facing a large monitor with a virtual female interviewer on screen",
        "background": "#FFFFFF",
    },
    {
        "id": "setup_mid",
        "bbox": (31, 342, 238, 428),
        "prompt": "large flat monitor showing a 3D female virtual interviewer portrait, clean lab illustration",
        "background": "#FFFFFF",
    },
    {
        "id": "setup_cam",
        "bbox": (22, 505, 235, 625),
        "prompt": "wide-angle camera on tripod capturing a standing participant, blue dotted field-of-view lines",
        "background": "#FFFFFF",
    },
    {
        "id": "raw_people",
        "bbox": (318, 430, 490, 530),
        "prompt": "small video thumbnails showing a standing participant in front of a neutral background",
        "background": "#FFFFFF",
    },
    {
        "id": "face",
        "bbox": (1205, 199, 1280, 267),
        "prompt": "front-facing male face portrait inside a face-detection crop frame",
        "background": "#F8FBFF",
    },
    {
        "id": "person_frame",
        "bbox": (1200, 305, 1280, 372),
        "prompt": "film-strip frame containing a full-body standing participant",
        "background": "#F8FBFF",
    },
    {
        "id": "pose",
        "bbox": (1212, 402, 1278, 470),
        "prompt": "full-body human pose skeleton with colored joints and limbs",
        "background": "#F8FBFF",
    },
    {
        "id": "dataset",
        "bbox": (1640, 345, 1745, 445),
        "prompt": "teal database cylinder with grouped user silhouettes in front",
        "background": "#FFFFFF",
    },
]


def x(px: float) -> float:
    return px / REF_W * SLIDE_W


def y(px: float) -> float:
    return px / REF_H * SLIDE_H


def w(px: float) -> float:
    return px / REF_W * SLIDE_W


def h(px: float) -> float:
    return px / REF_H * SLIDE_H


def rgb(hex_color: str) -> RGBColor:
    value = hex_color.strip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def crop_asset(slot_id: str, bbox: tuple[int, int, int, int]) -> Path:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    out = ASSET_DIR / f"{slot_id}.png"
    with Image.open(REFERENCE).convert("RGB") as image:
        image.crop(bbox).save(out)
    return out


def crop_reference_slot(slot_id: str, bbox: tuple[int, int, int, int]) -> Path:
    REFERENCE_CROP_DIR.mkdir(parents=True, exist_ok=True)
    out = REFERENCE_CROP_DIR / f"{slot_id}.png"
    with Image.open(REFERENCE).convert("RGB") as image:
        image.crop(bbox).save(out)
    return out


def supported_aspect_ratio(width: int, height: int) -> str:
    ratio = width / max(height, 1)
    if ratio >= 1.55:
        return "16:9"
    if ratio >= 1.18:
        return "4:3"
    if ratio >= 0.85:
        return "1:1"
    if ratio >= 0.62:
        return "3:4"
    return "9:16"


def generate_ai_asset(spec: dict, crop_path: Path, out_path: Path) -> dict:
    api_key = os.getenv("GEMINI_API_KEY") or os.getenv("API_KEY")
    url = os.getenv("GEMINI_GEN_IMG_URL")
    if not api_key or not url:
        raise RuntimeError("GEMINI_API_KEY/API_KEY and GEMINI_GEN_IMG_URL are required for API asset generation")

    left, top, right, bottom = spec["bbox"]
    aspect_ratio = supported_aspect_ratio(right - left, bottom - top)
    mime = mimetypes.guess_type(crop_path.name)[0] or "image/png"
    image_b64 = base64.b64encode(crop_path.read_bytes()).decode("ascii")
    prompt = (
        "Use the provided crop as a reference for object identity, pose, and scientific-diagram style. "
        "Create a clean standalone raster asset for an editable PowerPoint research pipeline figure. "
        f"Main subject: {spec['prompt']}. "
        f"Canvas aspect ratio: {aspect_ratio}. "
        "Keep the subject large and centered, with only narrow background margins. "
        "Preserve the reference's clean academic illustration style, subtle shadows, dark outline, and muted colors. "
        f"Use a flat background matching {spec['background']}. "
        "Do not include any readable text, labels, numbers, captions, watermarks, panel borders, or arrows. "
        "Do not generate a full diagram; generate only this local visual asset."
    )
    payload = {
        "contents": [{
            "role": "user",
            "parts": [
                {"text": prompt},
                {"inlineData": {"mimeType": mime, "data": image_b64}},
            ],
        }],
        "generationConfig": {
            "responseModalities": ["IMAGE"],
            "imageConfig": {"imageSize": "1K", "aspectRatio": aspect_ratio},
        },
    }
    out_path.parent.mkdir(parents=True, exist_ok=True)
    response = requests.post(
        url,
        headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
        data=json.dumps(payload),
        timeout=240,
    )
    response.raise_for_status()
    data = response.json()
    for candidate in data.get("candidates", []):
        for part in candidate.get("content", {}).get("parts", []):
            inline = part.get("inlineData") or part.get("inline_data")
            if inline and inline.get("data"):
                out_path.write_bytes(base64.b64decode(inline["data"]))
                return {
                    "slot_id": spec["id"],
                    "status": "generated",
                    "asset_path": str(out_path),
                    "reference_crop_path": str(crop_path),
                    "aspect_ratio": aspect_ratio,
                    "prompt": prompt,
                }
    raise RuntimeError(f"No image returned for {spec['id']}")


def prepare_visual_assets(asset_mode: str = "api", workers: int = 4, force_generate: bool = False) -> tuple[dict[str, Path], dict]:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    assets: dict[str, Path] = {}
    reports: list[dict] = []
    crops = {spec["id"]: crop_reference_slot(spec["id"], spec["bbox"]) for spec in ASSET_SPECS}

    if asset_mode == "crop":
        for spec in ASSET_SPECS:
            out_path = ASSET_DIR / f"{spec['id']}.png"
            with Image.open(crops[spec["id"]]).convert("RGB") as image:
                image.save(out_path)
            assets[spec["id"]] = out_path
            reports.append({"slot_id": spec["id"], "status": "reference_crop", "asset_path": str(out_path), "api_requests_attempted": 0})
        return assets, {"summary": "Visual asset preparation report.", "asset_mode": asset_mode, "api_requests_attempted": 0, "assets": reports}

    def run_one(spec: dict) -> dict:
        out_path = ASSET_DIR / f"{spec['id']}.png"
        if out_path.exists() and not force_generate:
            return {"slot_id": spec["id"], "status": "reused_existing_ai_asset", "asset_path": str(out_path), "api_requests_attempted": 0}
        try:
            report = generate_ai_asset(spec, crops[spec["id"]], out_path)
            return {**report, "api_requests_attempted": 1}
        except Exception as exc:
            with Image.open(crops[spec["id"]]).convert("RGB") as image:
                image.save(out_path)
            return {
                "slot_id": spec["id"],
                "status": "fallback_reference_crop_after_api_failure",
                "asset_path": str(out_path),
                "reference_crop_path": str(crops[spec["id"]]),
                "api_requests_attempted": 1,
                "error": str(exc),
            }

    with ThreadPoolExecutor(max_workers=workers) as executor:
        futures = {executor.submit(run_one, spec): spec for spec in ASSET_SPECS}
        for future in as_completed(futures):
            reports.append(future.result())
    reports.sort(key=lambda item: item["slot_id"])
    for item in reports:
        assets[item["slot_id"]] = Path(item["asset_path"])
    return assets, {
        "summary": "Visual asset preparation report.",
        "asset_mode": asset_mode,
        "api_requests_attempted": sum(int(item.get("api_requests_attempted") or 0) for item in reports),
        "assets": reports,
    }


def add_text(slide, text: str, left, top, width, height, *, size=10, color="#111111", bold=False, align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_textbox(Inches(x(left)), Inches(y(top)), Inches(w(width)), Inches(h(height)))
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return shape


def add_box(slide, left, top, width, height, *, line="#00328C", fill="#FFFFFF", radius=True, lw=1.2, dash=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x(left)), Inches(y(top)), Inches(w(width)), Inches(h(height)))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(lw)
    if dash:
        shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return shape


def add_picture_contain(slide, image_path: Path, left, top, width, height):
    with Image.open(image_path) as image:
        ratio = image.width / max(image.height, 1)
    box_ratio = width / max(height, 1)
    if ratio > box_ratio:
        fit_w = width
        fit_h = width / ratio
    else:
        fit_h = height
        fit_w = height * ratio
    return slide.shapes.add_picture(
        str(image_path),
        Inches(x(left + (width - fit_w) / 2)),
        Inches(y(top + (height - fit_h) / 2)),
        width=Inches(w(fit_w)),
        height=Inches(h(fit_h)),
    )


def add_arrow(slide, x1, y1, x2, y2, *, color="#000000", width=2.0):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x(x1)), Inches(y(y1)), Inches(x(x2)), Inches(y(y2)))
    conn.line.color.rgb = rgb(color)
    conn.line.width = Pt(width)
    conn.line.end_arrowhead = True
    return conn


def add_line(slide, x1, y1, x2, y2, *, color="#000000", width=1.5):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x(x1)), Inches(y(y1)), Inches(x(x2)), Inches(y(y2)))
    conn.line.color.rgb = rgb(color)
    conn.line.width = Pt(width)
    return conn


def add_panel(slide, left, top, width, height, title, color="#00328C"):
    add_box(slide, left, top, width, height, line=color, fill="#FFFFFF", radius=True, lw=1.5)
    add_text(slide, title, left + 8, top + 10, width - 16, 54, size=12.8, color=color, bold=True)


def add_film_icon(slide, left, top, width, height, *, fill="#3A6FA5"):
    box = add_box(slide, left, top, width, height, line="#1B2F4A", fill=fill, radius=True, lw=1.2)
    for i in range(5):
        add_box(slide, left + 8 + i * (width - 28) / 4, top + 8, 10, 8, line="#1B2F4A", fill="#BFD3E6", radius=False, lw=0.6)
        add_box(slide, left + 8 + i * (width - 28) / 4, top + height - 16, 10, 8, line="#1B2F4A", fill="#BFD3E6", radius=False, lw=0.6)
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(x(left + width * 0.42)), Inches(y(top + height * 0.30)), Inches(w(width * 0.22)), Inches(h(height * 0.40)))
    tri.rotation = 90
    tri.fill.solid()
    tri.fill.fore_color.rgb = rgb("#EAF2FF")
    tri.line.color.rgb = rgb("#1B2F4A")
    return box


def add_audio_wave(slide, cx, cy, width, height, *, color="#003E9C"):
    bars = 15
    for i in range(bars):
        frac = i / (bars - 1)
        amp = (0.25 + 0.75 * abs(0.5 - abs(frac - 0.5)) * 2) * height
        x0 = cx - width / 2 + frac * width
        add_line(slide, x0, cy - amp / 2, x0, cy + amp / 2, color=color, width=1.1)


def add_doc_icon(slide, left, top, width, height):
    add_box(slide, left, top, width, height, line="#1B1B1B", fill="#FFFFFF", radius=False, lw=1.1)
    for i in range(4):
        add_line(slide, left + 10, top + 15 + i * 10, left + width - 8, top + 15 + i * 10, color="#222222", width=0.8)


def add_clock(slide, cx, cy, r):
    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x(cx - r)), Inches(y(cy - r)), Inches(w(r * 2)), Inches(h(r * 2)))
    oval.fill.solid()
    oval.fill.fore_color.rgb = rgb("#F7FBFF")
    oval.line.color.rgb = rgb("#003A89")
    oval.line.width = Pt(3)
    add_line(slide, cx, cy, cx, cy - r * 0.58, color="#003A89", width=3)
    add_line(slide, cx, cy, cx + r * 0.42, cy + r * 0.35, color="#003A89", width=3)


def build_deck(asset_mode: str = "api", asset_workers: int = 4, force_generate_assets: bool = False) -> tuple[Path, dict]:
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb("#FFFFFF")

    assets, asset_report = prepare_visual_assets(asset_mode, asset_workers, force_generate_assets)

    blue = "#00338D"
    teal = "#007A83"
    orange = "#C65A00"
    purple = "#6E55B8"
    green = "#2D7D3A"

    add_panel(slide, 8, 43, 240, 646, "1. Virtual Interview\nSetup", blue)
    add_picture_contain(slide, assets["setup_top"], 28, 118, 205, 162)
    add_text(slide, "Participant in front of\n49-inch screen", 42, 286, 170, 48, size=8.8, bold=True)
    add_picture_contain(slide, assets["setup_mid"], 30, 342, 210, 93)
    add_text(slide, "3D virtual interviewer\non screen", 43, 451, 165, 44, size=8.8, bold=True)
    add_picture_contain(slide, assets["setup_cam"], 22, 500, 205, 128)
    add_text(slide, "Wide-angle camera\ncapturing full body", 44, 634, 170, 38, size=8.4, bold=True)

    add_panel(slide, 295, 168, 208, 429, "2. Raw Video:\n287 participants\n× 36 questions", blue)
    add_film_icon(slide, 327, 292, 149, 112, fill="#2E6FA7")
    for row, yy in enumerate([430, 506]):
        for col, xx in enumerate([316, 369, 424]):
            if row == 1 and col == 2:
                continue
            add_picture_contain(slide, assets["raw_people"], xx, yy, 49, 61)
    add_text(slide, "…", 446, 500, 40, 50, size=20, bold=True)

    add_panel(slide, 570, 43, 382, 672, "3. Preprocessing Modules", blue)
    module_specs = [
        (580, 96, 360, 100, "#0E60B8", "FFmpeg", "Audio"),
        (580, 219, 360, 100, "#2D7D3A", "FunASR", "Spoken Text\n+ Timestamps"),
        (580, 343, 360, 100, "#6E55B8", "MTCNN", "Face Clips"),
        (580, 467, 360, 100, "#00879C", "AlphaPose", "Full-body\nPose\nSkeletons"),
        (580, 590, 360, 100, "#C65A00", "Frame\nSampling", "Video\nFrames"),
    ]
    for left, top, width, height, color, method, out_label in module_specs:
        add_box(slide, left, top, width, height, line=color, fill="#F9FCFF", radius=True, lw=1.2)
        add_text(slide, method, left + 78, top + 30, 95, 42, size=10.4, color=color, bold=True)
        add_arrow(slide, left + 172, top + 50, left + 215, top + 50, color="#000000", width=1.5)
        out_left = left + 296 if method == "FFmpeg" else left + 260
        out_width = 55 if method == "FFmpeg" else 92
        add_text(slide, out_label, out_left, top + 28, out_width, 45, size=9.2 if method == "FFmpeg" else 9.5, bold=True)
    add_film_icon(slide, 596, 116, 63, 53, fill="#294C72")
    add_audio_wave(slide, 835, 146, 72, 42, color="#003E9C")
    mic = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x(606)), Inches(y(239)), Inches(w(21)), Inches(h(48)))
    mic.fill.solid(); mic.fill.fore_color.rgb = rgb("#1B4F82"); mic.line.color.rgb = rgb("#1B1B1B")
    add_line(slide, 616, 286, 616, 303, color="#1B1B1B", width=2)
    add_line(slide, 600, 303, 632, 303, color="#1B1B1B", width=2)
    add_doc_icon(slide, 780, 240, 45, 58)
    add_picture_contain(slide, assets["face"], 599, 369, 53, 51)
    add_picture_contain(slide, assets["face"], 782, 359, 59, 61)
    add_picture_contain(slide, assets["pose"], 596, 481, 55, 70)
    add_picture_contain(slide, assets["pose"], 788, 481, 55, 70)
    for i in range(3):
        add_box(slide, 594 + i * 8, 613 - i * 5, 47, 52, line="#333333", fill="#FFFFFF", radius=False, lw=0.7)
    add_picture_contain(slide, assets["person_frame"], 782, 604, 65, 70)

    add_panel(slide, 996, 184, 149, 436, "4. Timestamp\nAlignment", blue)
    add_clock(slide, 1071, 350, 45)
    add_text(slide, "Align audio,\ntext, face,\nframes, and\npose streams\nto a unified\ntimeline", 1020, 433, 104, 150, size=9.2, bold=True)

    add_panel(slide, 1176, 132, 186, 560, "5. Five Modalities", teal)
    modality_rows = [
        (1186, 186, purple, "Face", assets["face"]),
        (1186, 293, orange, "Frame", assets["person_frame"]),
        (1186, 395, teal, "Pose", assets["pose"]),
        (1186, 497, "#2076D0", "Audio", None),
        (1186, 600, green, "Text", None),
    ]
    for left, top, color, label, asset in modality_rows:
        add_box(slide, left, top, 168, 87, line=color, fill="#F8FBFF", radius=True, lw=1.1)
        if asset:
            add_picture_contain(slide, asset, left + 18, top + 12, 70, 58)
        elif label == "Audio":
            add_audio_wave(slide, left + 48, top + 45, 78, 42, color="#003E9C")
        else:
            add_doc_icon(slide, left + 29, top + 15, 43, 57)
        add_text(slide, label, left + 105, top + 27, 55, 36, size=10.5, color=color, bold=True)

    add_panel(slide, 1394, 132, 184, 560, "6. NEO-FFI-3\nSelf-report Labels", orange)
    clip = add_box(slide, 1445, 219, 75, 99, line="#0B3264", fill="#F7FBFF", radius=True, lw=1.2)
    add_box(slide, 1470, 217, 25, 10, line="#0B3264", fill="#0B3264", radius=True, lw=1.0)
    for i in range(4):
        add_text(slide, "✓", 1456, 245 + i * 19, 16, 14, size=10, color="#0B3264", bold=True)
        add_line(slide, 1476, 253 + i * 19, 1513, 253 + i * 19, color="#222222", width=1.2)
    add_text(slide, "OCEAN Scores", 1424, 345, 115, 24, size=10, bold=True)
    ocean = [("O", "Openness", purple), ("C", "Conscientiousness", "#D66300"), ("E", "Extraversion", "#00919B"), ("A", "Agreeableness", "#4B934A"), ("N", "Neuroticism", "#3C82C8")]
    for i, (letter, name, color) in enumerate(ocean):
        cy = 403 + i * 55
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x(1410)), Inches(y(cy - 17)), Inches(w(30)), Inches(h(35)))
        oval.fill.solid(); oval.fill.fore_color.rgb = rgb(color); oval.line.color.rgb = rgb(color)
        add_text(slide, letter, 1413, cy - 12, 24, 24, size=11, color="#FFFFFF", bold=True)
        add_text(slide, name, 1448, cy - 12, 120, 26, size=8.8, bold=True, align=PP_ALIGN.LEFT)

    add_panel(slide, 1612, 235, 155, 397, "7. Multimodal\nPersonality\nDataset", teal)
    add_picture_contain(slide, assets["dataset"], 1635, 344, 105, 95)
    add_text(slide, "Multimodal\nvideos aligned\nacross five\nmodalities with\nNEO-FFI-3\n(OCEAN) labels", 1630, 470, 120, 132, size=8.7, bold=True)

    # Macro arrows.
    for args in [
        (248, 384, 294, 384), (503, 384, 570, 384), (940, 268, 996, 268), (940, 389, 996, 389),
        (940, 545, 996, 545), (1145, 386, 1176, 386), (1362, 233, 1394, 233), (1362, 335, 1394, 335),
        (1362, 441, 1394, 441), (1362, 543, 1394, 543), (1362, 645, 1394, 645), (1578, 386, 1612, 386)
    ]:
        add_arrow(slide, *args, color="#000000", width=1.8)
    # Bracket-like incoming fan from raw video to preprocessing rows.
    add_line(slide, 536, 140, 536, 685, color="#000000", width=1.3)
    for yy in [140, 260, 384, 512, 685]:
        add_arrow(slide, 536, yy, 570, yy, color="#000000", width=1.3)

    # Legend.
    add_box(slide, 19, 762, 918, 84, line="#18345C", fill="#FFFFFF", radius=True, lw=1.0, dash=True)
    legend = [
        ("Video", "film"), ("Audio", "audio"), ("Text", "doc"), ("Face", "face"), ("Pose", "pose"), ("Frames", "frames")
    ]
    lx = 44
    for label, kind in legend:
        if kind == "film":
            add_film_icon(slide, lx, 784, 48, 43, fill="#294C72")
        elif kind == "audio":
            add_audio_wave(slide, lx + 23, 805, 48, 34, color="#003E9C")
        elif kind == "doc":
            add_doc_icon(slide, lx, 781, 32, 47)
        elif kind == "face":
            add_picture_contain(slide, assets["face"], lx, 780, 48, 49)
        elif kind == "pose":
            add_picture_contain(slide, assets["pose"], lx, 777, 48, 55)
        else:
            for i in range(3):
                add_box(slide, lx + i * 6, 781 - i * 3, 40, 47, line="#333333", fill="#FFFFFF", radius=False, lw=0.7)
            add_picture_contain(slide, assets["person_frame"], lx + 12, 786, 35, 39)
        add_text(slide, label, lx + 62, 790, 70, 32, size=9.4, bold=True, align=PP_ALIGN.LEFT)
        lx += 145

    pptx = OUT / "editable_composition.pptx"
    prs.save(pptx)
    return pptx, asset_report


def export_preview(pptx_path: Path) -> None:
    try:
        import win32com.client  # type: ignore
    except Exception:
        return
    app = win32com.client.Dispatch("PowerPoint.Application")
    presentation = app.Presentations.Open(str(pptx_path), WithWindow=False)
    try:
        presentation.Slides(1).Export(str(OUT / "rebuild_preview.png"), "PNG", REF_W, REF_H)
    finally:
        presentation.Close()
        app.Quit()


def write_reports(pptx_path: Path, asset_report: dict) -> None:
    report = {
        "summary": "Editable PPT reconstruction of the virtual interview multimodal personality dataset pipeline.",
        "reference": str(REFERENCE),
        "reference_size_px": [REF_W, REF_H],
        "slide_size_in": [SLIDE_W, SLIDE_H],
        "editable_layers": ["panel frames", "labels", "method names", "arrows", "legend", "simple icons"],
        "raster_assets": sorted(path.name for path in ASSET_DIR.glob("*.png")),
        "asset_policy": asset_report.get("asset_mode"),
        "api_requests_attempted": asset_report.get("api_requests_attempted"),
        "output_pptx": str(pptx_path),
    }
    (OUT / "figure_program.json").write_text(json.dumps(report, indent=2, ensure_ascii=False), encoding="utf-8")
    (OUT / "asset_generation_report.json").write_text(json.dumps(asset_report, indent=2, ensure_ascii=False), encoding="utf-8")
    (OUT / "rebuild_notes.md").write_text(
        "# Summary\n"
        "This rebuild uses editable PowerPoint shapes for the pipeline structure, text, arrows, labels, and legend. "
        "Pictorial interview and modality examples are inserted as slot-level raster assets generated from local reference crops when API mode is enabled.\n",
        encoding="utf-8",
    )


def main() -> None:
    parser = argparse.ArgumentParser(description="Rebuild the virtual interview pipeline as an editable PPTX.")
    parser.add_argument("--asset-mode", choices=["api", "crop"], default="api", help="Use API-generated slot assets or direct reference crops.")
    parser.add_argument("--asset-workers", type=int, default=4, help="Parallel workers for API asset generation.")
    parser.add_argument("--force-generate-assets", action="store_true", help="Regenerate API assets even if existing assets are present.")
    args = parser.parse_args()
    pptx, asset_report = build_deck(args.asset_mode, args.asset_workers, args.force_generate_assets)
    write_reports(pptx, asset_report)
    export_preview(pptx)


if __name__ == "__main__":
    main()
