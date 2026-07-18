from __future__ import annotations

import base64
import argparse
import json
import math
import mimetypes
import os
import shutil
import time
from collections import deque
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path

import requests
from PIL import Image, ImageDraw, ImageStat
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
REFERENCE_CANDIDATES = [
    Path(r"C:\Users\zhang\Documents\xwechat_files\wxid_h824xk1qpfoh22_adc2\temp\RWTemp\2026-06\860847045f7ccc2e7052eb182bab099b\8fc4985d14085a49b20f61b3d56d0bea.png"),
    ROOT / "output" / "autofigure_reference_real_v10_layout" / "inputs" / "reference.png",
    ROOT / "output" / "autofigure_reference_contract_v2" / "inputs" / "reference.png",
]
REFERENCE = next((path for path in REFERENCE_CANDIDATES if path.exists()), REFERENCE_CANDIDATES[0])
OUT = ROOT / "output" / "autofigure_architecture_ai_rebuild"
CROP_DIR = OUT / "reference_slot_crops"
ASSET_DIR = OUT / "assets"
CANDIDATE_DIR = OUT / "asset_candidates"

REF_W = 1770
REF_H = 975
SLIDE_W = 15.60
SLIDE_H = SLIDE_W * REF_H / REF_W

FONT = "Arial"
TEXT_SIZE_SCALE = 1.08
INTERNAL_CONTENT_FILL_TARGET = 0.88
INTERNAL_ACCEPTABLE_FILL_MIN = 0.80
INTERNAL_ACCEPTABLE_FILL_MAX = 0.95
INTERNAL_REGENERATE_FILL_THRESHOLD = 0.70
INTERNAL_MAX_MARGIN_PERCENT = 0.10
SLOT_FIT_WASTE_ACCEPTABLE_MAX = 0.16
GENERATION_ASPECT_RATIO_VALUES = {
    "1:1": 1.0,
    "3:4": 3 / 4,
    "4:3": 4 / 3,
    "9:16": 9 / 16,
    "16:9": 16 / 9,
}
TEXT_STYLES = {
    "title": {"font": FONT, "size": 23.0, "bold": True},
    "subtitle": {"font": FONT, "size": 10.8, "bold": False},
    "stage_title": {"font": FONT, "size": 13.0, "bold": True},
    "label": {"font": FONT, "size": 10.5, "bold": True},
    "body": {"font": FONT, "size": 8.6, "bold": False},
    "small": {"font": FONT, "size": 7.4, "bold": False},
    "callout": {"font": FONT, "size": 9.0, "bold": True},
    "legend": {"font": FONT, "size": 9.2, "bold": True},
}

SLOTS = [
    {
        "id": "input_text_stack",
        "bbox": (66, 254, 200, 400),
        "asset_bbox": (68, 255, 205, 398),
        "prompt": "a stack of off-white paper documents with simple abstract pseudo-lines and two tiny pseudo-number marks, cute academic flat illustration, no readable text",
        "fill_guidance": "Make the paper stack very large, nearly full height, with page corners close to all four safe margins. Avoid tiny centered sheets and avoid showing broad blue background around the documents.",
        "background_color_hex": "#EAF3FF",
        "internal_content_fill_target": INTERNAL_CONTENT_FILL_TARGET,
        "internal_max_margin_percent": INTERNAL_MAX_MARGIN_PERCENT,
        "layout_bbox_locked": True,
        "no_crop_postprocess": True,
    },
    {
        "id": "vlm_agent_robot",
        "bbox": (280, 258, 458, 404),
        "asset_bbox": (294, 250, 454, 407),
        "prompt": "a cute round robot assistant head with small antennae and soft blue-gray casing, pastel academic diagram style, no text",
        "background_color_hex": "#EAF3FF",
        "internal_content_fill_target": INTERNAL_CONTENT_FILL_TARGET,
        "internal_max_margin_percent": INTERNAL_MAX_MARGIN_PERCENT,
        "layout_bbox_locked": True,
        "no_crop_postprocess": True,
    },
    {
        "id": "ai_designer",
        "bbox": (642, 441, 793, 629),
        "asset_bbox": (650, 425, 790, 627),
        "prompt": "a cute AI designer character wearing a beret and holding a blueprint tablet, pastel orange academic illustration, no readable text",
        "background_color_hex": "#FFF1DF",
        "internal_content_fill_target": INTERNAL_CONTENT_FILL_TARGET,
        "internal_max_margin_percent": INTERNAL_MAX_MARGIN_PERCENT,
        "layout_bbox_locked": True,
        "no_crop_postprocess": True,
    },
    {
        "id": "ai_critic",
        "bbox": (1059, 438, 1212, 632),
        "asset_bbox": (1062, 434, 1208, 632),
        "prompt": "a cute robot critic wearing round glasses and holding a clipboard, pastel academic illustration, no readable text",
        "fill_guidance": "Use a large close-up robot critic with head, torso, glasses, and clipboard filling the portrait canvas. The robot should nearly touch the side safe margins and bottom safe margin.",
        "background_color_hex": "#FFF1DF",
        "internal_content_fill_target": INTERNAL_CONTENT_FILL_TARGET,
        "internal_max_margin_percent": INTERNAL_MAX_MARGIN_PERCENT,
        "layout_bbox_locked": True,
        "no_crop_postprocess": True,
    },
    {
        "id": "synthesis_tools",
        "bbox": (1375, 132, 1625, 247),
        "asset_bbox": (1375, 126, 1628, 250),
        "prompt": "a magic wand and cheerful painter palette with soft sparkles, pastel scientific illustration, no text",
        "background_color_hex": "#ECF7E9",
        "internal_content_fill_target": INTERNAL_CONTENT_FILL_TARGET,
        "internal_max_margin_percent": INTERNAL_MAX_MARGIN_PERCENT,
        "layout_bbox_locked": True,
        "no_crop_postprocess": True,
    },
    {
        "id": "erase_text_tool",
        "bbox": (1288, 493, 1420, 604),
        "asset_bbox": (1290, 486, 1425, 606),
        "prompt": "one large front-facing rounded pink eraser icon, cute academic diagram illustration, no readable text, no letters, no numbers, no pseudo-letters",
        "fill_guidance": "Make one oversized front-facing rounded eraser dominate the square canvas, with its body occupying roughly 85-90% of the canvas width and height. Use a broad blocky eraser shape rather than a thin diagonal object. Keep the eraser close to all safe margins without cutting it off. Do not add ABC, letters, numbers, words, pseudo-text, distant dust, particles, side objects, or broad empty background.",
        "background_color_hex": "#ECF7E9",
        "internal_content_fill_target": INTERNAL_CONTENT_FILL_TARGET,
        "internal_max_margin_percent": INTERNAL_MAX_MARGIN_PERCENT,
        "layout_bbox_locked": True,
        "no_crop_postprocess": True,
    },
    {
        "id": "ocr_verify",
        "bbox": (1570, 492, 1712, 604),
        "asset_bbox": (1566, 490, 1718, 608),
        "prompt": "a magnifying glass inspecting pseudo-numbers with a green check badge and short gray pseudo-lines, clean pastel academic illustration, no readable text",
        "fill_guidance": "Make the magnifying glass lens very large, occupying most of the canvas height, with the handle and check badge tightly grouped into one compact foreground silhouette. Avoid broad side margins and keep decorative side lines minimal.",
        "background_color_hex": "#ECF7E9",
        "internal_content_fill_target": INTERNAL_CONTENT_FILL_TARGET,
        "internal_max_margin_percent": INTERNAL_MAX_MARGIN_PERCENT,
        "layout_bbox_locked": True,
        "no_crop_postprocess": True,
    },
    {
        "id": "final_autofigure_card",
        "bbox": (1402, 682, 1590, 812),
        "asset_bbox": (1408, 686, 1588, 808),
        "prompt": "a small polished scientific figure card with abstract pie chart, smooth shapes, tiny pseudo-lines and pleasing pastel colors, no readable text",
        "fill_guidance": "Make the figure card itself large and nearly full canvas, with its rounded rectangle edges close to the safe margins. Internal charts should be big enough to count as foreground, not tiny decorations.",
        "background_color_hex": "#ECF7E9",
        "internal_content_fill_target": INTERNAL_CONTENT_FILL_TARGET,
        "internal_max_margin_percent": INTERNAL_MAX_MARGIN_PERCENT,
        "layout_bbox_locked": True,
        "no_crop_postprocess": True,
    },
]

CONTROL_SPECS = [
    {"id": "input_to_vlm", "kind": "line", "source_id": "input_text_stack", "target_id": "vlm_agent_robot", "points": [(222, 346), (278, 346)], "color": "#607998", "width_pt": 2.5, "arrow": True, "dash": False, "confidence": 0.92},
    {"id": "vlm_to_blueprint", "kind": "line", "source_id": "vlm_agent_robot", "target_id": "initial_blueprint", "points": [(365, 407), (365, 508)], "color": "#607998", "width_pt": 2.8, "arrow": True, "dash": False, "confidence": 0.91},
    {"id": "stage_i_to_stage_ii", "kind": "line", "source_id": "stage_i", "target_id": "stage_ii", "points": [(565, 493), (626, 493)], "color": "#555555", "width_pt": 4.0, "arrow": True, "dash": False, "confidence": 0.95},
    {"id": "stage_ii_to_stage_iii", "kind": "line", "source_id": "stage_ii", "target_id": "stage_iii", "points": [(1228, 493), (1270, 493)], "color": "#555555", "width_pt": 4.0, "arrow": True, "dash": False, "confidence": 0.95},
    {"id": "critique_refinement_loop", "kind": "oval", "source_id": "ai_critic", "target_id": "ai_designer", "bbox": (735, 335, 350, 382), "color": "#D99045", "width_pt": 2.0, "arrow": False, "dash": True, "flow": "clockwise", "confidence": 0.86},
    {"id": "designer_refine_arrow", "kind": "line", "source_id": "critique_refinement_loop", "target_id": "ai_designer", "points": [(814, 685), (774, 648)], "color": "#D99045", "width_pt": 2.0, "arrow": True, "dash": False, "confidence": 0.84},
    {"id": "critic_feedback_arrow", "kind": "line", "source_id": "critique_refinement_loop", "target_id": "ai_critic", "points": [(1038, 388), (1072, 430)], "color": "#D99045", "width_pt": 2.0, "arrow": True, "dash": False, "confidence": 0.84},
    {"id": "synthesis_to_raw", "kind": "line", "source_id": "synthesis_tools", "target_id": "raw_image", "points": [(1500, 280), (1500, 350)], "color": "#6F9A66", "width_pt": 3.0, "arrow": True, "dash": False, "confidence": 0.9},
    {"id": "raw_to_erase_route", "kind": "polyline", "source_id": "raw_image", "target_id": "erase_text_tool", "points": [(1428, 405), (1362, 405), (1348, 428), (1345, 462)], "color": "#6F9A66", "width_pt": 1.8, "arrow": True, "dash": True, "confidence": 0.9},
    {"id": "raw_to_ocr_route", "kind": "polyline", "source_id": "raw_image", "target_id": "ocr_verify", "points": [(1588, 405), (1652, 405), (1684, 426), (1692, 462)], "color": "#6F9A66", "width_pt": 1.8, "arrow": True, "dash": True, "confidence": 0.9},
    {"id": "erase_to_final_route", "kind": "polyline", "source_id": "erase_text_tool", "target_id": "final_autofigure_card", "points": [(1345, 624), (1350, 650), (1382, 680)], "color": "#6F9A66", "width_pt": 1.8, "arrow": True, "dash": True, "confidence": 0.88},
    {"id": "ocr_to_final_route", "kind": "polyline", "source_id": "ocr_verify", "target_id": "final_autofigure_card", "points": [(1644, 624), (1638, 650), (1608, 680)], "color": "#6F9A66", "width_pt": 1.8, "arrow": True, "dash": True, "confidence": 0.88},
    {"id": "finalization_arrow", "kind": "line", "source_id": "erase_text_tool", "target_id": "final_autofigure_card", "points": [(1306, 737), (1386, 737)], "color": "#6F9A66", "width_pt": 3.5, "arrow": True, "dash": False, "confidence": 0.86},
]


def px(v: float, total: float, slide_total: float) -> float:
    return v / total * slide_total


def x(v: float) -> float:
    return px(v, REF_W, SLIDE_W)


def y(v: float) -> float:
    return px(v, REF_H, SLIDE_H)


def w(v: float) -> float:
    return px(v, REF_W, SLIDE_W)


def h(v: float) -> float:
    return px(v, REF_H, SLIDE_H)


def rgb(value: str) -> RGBColor:
    text = value.strip().lstrip("#")
    return RGBColor(int(text[:2], 16), int(text[2:4], 16), int(text[4:6], 16))


def hex_to_tuple(value: str) -> tuple[int, int, int]:
    text = value.strip().lstrip("#")
    return int(text[:2], 16), int(text[2:4], 16), int(text[4:6], 16)


def tuple_to_hex(value: tuple[int, int, int]) -> str:
    return f"#{value[0]:02X}{value[1]:02X}{value[2]:02X}"


def set_fill(shape, color: str, transparency: int = 0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.fill.transparency = transparency


def no_fill(shape) -> None:
    shape.fill.background()


def set_line(shape, color: str, width_pt: float = 1.2, dash: bool = False) -> None:
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width_pt)
    if dash:
        shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH


def add_text(slide, text: str, left: float, top: float, width: float, height: float, *, style: str = "body", color: str = "#333333", align=PP_ALIGN.CENTER, size: float | None = None, bold: bool | None = None, valign=MSO_ANCHOR.MIDDLE):
    spec = TEXT_STYLES[style]
    font_size = float(size if size is not None else spec["size"]) * TEXT_SIZE_SCALE
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Inches(0.01)
    tf.margin_right = Inches(0.01)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    for idx, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = str(spec["font"])
        r.font.size = Pt(font_size)
        r.font.bold = bool(bold if bold is not None else spec["bold"])
        r.font.color.rgb = rgb(color)
    return box


def add_round(slide, left_px, top_px, width_px, height_px, fill, stroke, width_pt=1.2, radius=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x(left_px)), Inches(y(top_px)), Inches(w(width_px)), Inches(h(height_px)))
    set_fill(shape, fill)
    set_line(shape, stroke, width_pt)
    shape.shadow.inherit = False
    return shape


def add_arrowhead(connector, size="med") -> None:
    ln = connector.line._get_or_add_ln()
    ln.append(parse_xml(f'<a:tailEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" type="triangle" w="{size}" len="{size}"/>'))


def add_line(slide, x1, y1, x2, y2, color="#333333", width_pt=1.2, arrow=False, dash=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x(x1)), Inches(y(y1)), Inches(x(x2)), Inches(y(y2)))
    set_line(conn, color, width_pt, dash=dash)
    if arrow:
        add_arrowhead(conn)
    return conn


def add_poly(slide, points, color="#333333", width_pt=1.2, arrow=False, dash=False):
    for idx, (a, b) in enumerate(zip(points[:-1], points[1:])):
        add_line(slide, a[0], a[1], b[0], b[1], color=color, width_pt=width_pt, arrow=arrow and idx == len(points) - 2, dash=dash)


def bbox_percent_from_points(points: list[tuple[float, float]]) -> dict:
    xs = [point[0] for point in points]
    ys = [point[1] for point in points]
    left, right = min(xs), max(xs)
    top, bottom = min(ys), max(ys)
    return {"x": round(left / REF_W, 4), "y": round(top / REF_H, 4), "w": round((right - left) / REF_W, 4), "h": round((bottom - top) / REF_H, 4)}


def direction_label(dx: float, dy: float) -> str:
    if abs(dx) < 1e-6 and abs(dy) < 1e-6:
        return "none"
    if abs(dx) >= abs(dy) * 1.8:
        return "right" if dx > 0 else "left"
    if abs(dy) >= abs(dx) * 1.8:
        return "down" if dy > 0 else "up"
    vertical = "down" if dy > 0 else "up"
    horizontal = "right" if dx > 0 else "left"
    return f"{vertical}_{horizontal}"


def anchor_from_direction(dx: float, dy: float, at_source: bool) -> str:
    if abs(dx) >= abs(dy):
        if dx > 0:
            return "right_center" if at_source else "left_center"
        return "left_center" if at_source else "right_center"
    if dy > 0:
        return "bottom_center" if at_source else "top_center"
    return "top_center" if at_source else "bottom_center"


def position_direction_metadata(control: dict) -> dict:
    if control["kind"] == "oval":
        left, top, width, height = control["bbox"]
        return {
            "position": {
                "bbox_px": [left, top, width, height],
                "center_px": [round(left + width / 2, 2), round(top + height / 2, 2)],
                "bbox_percent": {"x": round(left / REF_W, 4), "y": round(top / REF_H, 4), "w": round(width / REF_W, 4), "h": round(height / REF_H, 4)},
            },
            "direction": {
                "path_type": "closed_loop",
                "flow": control.get("flow", "clockwise"),
                "arrowhead_at": control.get("arrowhead_at", "none"),
                "direction_confidence": control.get("confidence", 0.0),
            },
            "source_anchor": control.get("source_anchor", "loop_right_side"),
            "target_anchor": control.get("target_anchor", "loop_left_side"),
        }

    points = control["points"]
    start = points[0]
    end = points[-1]
    dx = end[0] - start[0]
    dy = end[1] - start[1]
    length = max((dx * dx + dy * dy) ** 0.5, 1e-6)
    angle = math.degrees(math.atan2(dy, dx))
    return {
        "position": {
            "start_px": [start[0], start[1]],
            "end_px": [end[0], end[1]],
            "points_px": [[px_, py_] for px_, py_ in points],
            "center_px": [round((start[0] + end[0]) / 2, 2), round((start[1] + end[1]) / 2, 2)],
            "bbox_percent": bbox_percent_from_points(points),
        },
        "direction": {
            "vector_px": [round(dx, 2), round(dy, 2)],
            "unit_vector": [round(dx / length, 4), round(dy / length, 4)],
            "angle_deg": round(angle, 2),
            "label": direction_label(dx, dy),
            "arrowhead_at": "end" if control.get("arrow") else "none",
            "direction_confidence": control.get("confidence", 0.0),
        },
        "source_anchor": control.get("source_anchor", anchor_from_direction(dx, dy, True)),
        "target_anchor": control.get("target_anchor", anchor_from_direction(dx, dy, False)),
    }


def normalize_control(control: dict) -> dict:
    item = dict(control)
    item.update(position_direction_metadata(control))
    if control["kind"] == "oval":
        left, top, width, height = control["bbox"]
        item["bbox_percent"] = {"x": round(left / REF_W, 4), "y": round(top / REF_H, 4), "w": round(width / REF_W, 4), "h": round(height / REF_H, 4)}
        item["path_percent"] = []
    else:
        item["bbox_percent"] = bbox_percent_from_points(control["points"])
        item["path_percent"] = [{"x": round(px_ / REF_W, 4), "y": round(py_ / REF_H, 4)} for px_, py_ in control["points"]]
    item["render_policy"] = "ppt_shape_not_image_asset"
    item["route_policy"] = "reference_locked"
    return item


def render_control(slide, control: dict):
    if control["kind"] == "oval":
        left, top, width, height = control["bbox"]
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x(left)), Inches(y(top)), Inches(w(width)), Inches(h(height)))
        no_fill(shape)
        set_line(shape, control["color"], control["width_pt"], dash=control["dash"])
        return shape
    if control["kind"] == "polyline":
        return add_poly(slide, control["points"], control["color"], control["width_pt"], arrow=control["arrow"], dash=control["dash"])
    start, end = control["points"]
    return add_line(slide, start[0], start[1], end[0], end[1], control["color"], control["width_pt"], arrow=control["arrow"], dash=control["dash"])


def render_controls(slide, control_ids: list[str]) -> None:
    lookup = {control["id"]: control for control in CONTROL_SPECS}
    for control_id in control_ids:
        render_control(slide, lookup[control_id])


def draw_dashed_line(draw: ImageDraw.ImageDraw, points: list[tuple[float, float]], color: str, width: int = 3, dash: int = 14, gap: int = 9) -> None:
    fill = hex_to_tuple(color)
    for start, end in zip(points[:-1], points[1:]):
        x1, y1 = start
        x2, y2 = end
        length = max(((x2 - x1) ** 2 + (y2 - y1) ** 2) ** 0.5, 1)
        steps = int(length // (dash + gap)) + 1
        for step in range(steps):
            a = step * (dash + gap) / length
            b = min((step * (dash + gap) + dash) / length, 1)
            if a >= 1:
                continue
            draw.line((x1 + (x2 - x1) * a, y1 + (y2 - y1) * a, x1 + (x2 - x1) * b, y1 + (y2 - y1) * b), fill=fill, width=width)


def write_reference_overlays() -> None:
    with Image.open(REFERENCE).convert("RGB") as image:
        slot_overlay = image.copy()
        slot_draw = ImageDraw.Draw(slot_overlay)
        for slot in SLOTS:
            slot_draw.rectangle(slot["bbox"], outline=hex_to_tuple("#2A7BD1"), width=4)
            slot_draw.text((slot["bbox"][0] + 4, slot["bbox"][1] + 4), slot["id"], fill=hex_to_tuple("#2A7BD1"))
        slot_overlay.save(OUT / "slot_overlay.png")

        control_overlay = image.copy()
        control_draw = ImageDraw.Draw(control_overlay)
        for control in CONTROL_SPECS:
            color = control["color"]
            if control["kind"] == "oval":
                left, top, width, height = control["bbox"]
                control_draw.ellipse((left, top, left + width, top + height), outline=hex_to_tuple(color), width=5)
                label_xy = (left, top)
            else:
                points = control["points"]
                if control["dash"]:
                    draw_dashed_line(control_draw, points, color, width=5)
                else:
                    control_draw.line(points, fill=hex_to_tuple(color), width=5)
                label_xy = points[0]
            control_draw.text((label_xy[0] + 5, label_xy[1] + 5), control["id"], fill=hex_to_tuple("#111111"))
        control_overlay.save(OUT / "reference_control_overlay.png")


def crop_reference_slots() -> list[dict]:
    CROP_DIR.mkdir(parents=True, exist_ok=True)
    with Image.open(REFERENCE) as image:
        image = image.convert("RGB")
        for slot in SLOTS:
            crop = image.crop(slot["bbox"])
            crop_path = CROP_DIR / f"{slot['id']}.png"
            crop.save(crop_path)
            slot["crop_path"] = str(crop_path)
            x0, y0, x1, y1 = slot["bbox"]
            slot["bbox_percent"] = {"x": round(x0 / REF_W, 4), "y": round(y0 / REF_H, 4), "w": round((x1 - x0) / REF_W, 4), "h": round((y1 - y0) / REF_H, 4)}
    return SLOTS


def estimate_edge_background(image: Image.Image) -> tuple[int, int, int]:
    rgb_image = image.convert("RGB")
    width, height = rgb_image.size
    strip = max(4, min(width, height) // 20)
    samples = []
    for crop_box in [
        (0, 0, width, strip),
        (0, height - strip, width, height),
        (0, 0, strip, height),
        (width - strip, 0, width, height),
    ]:
        crop = rgb_image.crop(crop_box)
        if hasattr(crop, "get_flattened_data"):
            samples.extend(crop.get_flattened_data())
        else:
            samples.extend(crop.getdata())
    samples.sort(key=lambda color: sum(color))
    mid = samples[len(samples) // 2]
    return int(mid[0]), int(mid[1]), int(mid[2])


def harmonize_asset_background(asset_path: Path, background_hex: str) -> dict:
    image = Image.open(asset_path).convert("RGB")
    target = hex_to_tuple(background_hex)
    source_bg = estimate_edge_background(image)
    mask = edge_connected_background_mask(image, source_bg, tolerance=42)
    output = image.copy()
    bg_layer = Image.new("RGB", image.size, target)
    output.paste(bg_layer, mask=mask)
    output.save(asset_path)
    stat = ImageStat.Stat(mask)
    replaced_percent = round(stat.mean[0] / 255 * 100, 2)
    return {
        "asset": str(asset_path),
        "estimated_source_background": tuple_to_hex(source_bg),
        "target_background": background_hex,
        "replaced_percent": replaced_percent,
    }


def measure_internal_content_fill(image_path: Path, background_hex: str) -> dict:
    image = Image.open(image_path).convert("RGB")
    width, height = image.size
    target_bg = hex_to_tuple(background_hex)
    edge_bg = estimate_edge_background(image)
    pixels = image.load()
    tolerance = 54
    xs: list[int] = []
    ys: list[int] = []
    for py_ in range(height):
        for px_ in range(width):
            color = pixels[px_, py_]
            target_dist = abs(color[0] - target_bg[0]) + abs(color[1] - target_bg[1]) + abs(color[2] - target_bg[2])
            edge_dist = abs(color[0] - edge_bg[0]) + abs(color[1] - edge_bg[1]) + abs(color[2] - edge_bg[2])
            if min(target_dist, edge_dist) > tolerance:
                xs.append(px_)
                ys.append(py_)
    if not xs:
        return {
            "candidate_path": str(image_path),
            "foreground_bbox_px": None,
            "foreground_bbox_fill_percent": 0.0,
            "margin_left_percent": 1.0,
            "margin_right_percent": 1.0,
            "margin_top_percent": 1.0,
            "margin_bottom_percent": 1.0,
            "estimated_edge_background": tuple_to_hex(edge_bg),
            "crop_applied": False,
            "warning": "no_foreground_detected",
        }
    left, right = min(xs), max(xs)
    top, bottom = min(ys), max(ys)
    bbox_width = right - left + 1
    bbox_height = bottom - top + 1
    margins = {
        "margin_left_percent": round(left / width, 4),
        "margin_right_percent": round((width - 1 - right) / width, 4),
        "margin_top_percent": round(top / height, 4),
        "margin_bottom_percent": round((height - 1 - bottom) / height, 4),
    }
    return {
        "candidate_path": str(image_path),
        "foreground_bbox_px": [left, top, right, bottom],
        "foreground_bbox_fill_percent": round((bbox_width * bbox_height) / (width * height), 4),
        **margins,
        "max_margin_percent": max(margins.values()),
        "estimated_edge_background": tuple_to_hex(edge_bg),
        "crop_applied": False,
    }


def score_internal_fill(metrics: dict) -> float:
    fill = float(metrics.get("foreground_bbox_fill_percent") or 0.0)
    max_margin = float(metrics.get("max_margin_percent") or 1.0)
    cutoff_penalty = 0.0
    for key in ["margin_left_percent", "margin_right_percent", "margin_top_percent", "margin_bottom_percent"]:
        value = float(metrics.get(key) or 0.0)
        if value < 0.015:
            cutoff_penalty += 0.08
    if INTERNAL_ACCEPTABLE_FILL_MIN <= fill <= INTERNAL_ACCEPTABLE_FILL_MAX:
        fill_score = 1.0 - abs(fill - INTERNAL_CONTENT_FILL_TARGET)
    elif fill < INTERNAL_ACCEPTABLE_FILL_MIN:
        fill_score = fill / INTERNAL_ACCEPTABLE_FILL_MIN
    else:
        fill_score = max(0.0, 1.0 - (fill - INTERNAL_ACCEPTABLE_FILL_MAX) * 3.0)
    margin_score = max(0.0, 1.0 - max(0.0, max_margin - INTERNAL_MAX_MARGIN_PERCENT) * 2.0)
    return round(fill_score * 0.78 + margin_score * 0.22 - cutoff_penalty, 4)


def evaluate_candidate(slot: dict, candidate_path: Path) -> dict:
    metrics = measure_internal_content_fill(candidate_path, slot["background_color_hex"])
    metrics.update(compute_slot_fit_metrics(candidate_path, slot))
    fill = float(metrics.get("foreground_bbox_fill_percent") or 0.0)
    slot_fit_waste = float(metrics.get("slot_fit_waste_percent") or 1.0)
    edge_touch_count = sum(
        1 for key in ["margin_left_percent", "margin_right_percent", "margin_top_percent", "margin_bottom_percent"]
        if float(metrics.get(key) or 0.0) < 0.01
    )
    metrics["edge_touch_count"] = edge_touch_count
    metrics["background_or_edge_artifact_risk"] = fill > 0.97 and edge_touch_count >= 2
    internal_score = score_internal_fill(metrics)
    slot_fit_score = max(0.0, 1.0 - slot_fit_waste)
    metrics["score"] = round(internal_score * 0.72 + slot_fit_score * 0.28, 4)
    metrics["internal_fill_score"] = internal_score
    metrics["slot_fit_score"] = round(slot_fit_score, 4)
    metrics["pass"] = (
        INTERNAL_ACCEPTABLE_FILL_MIN <= fill <= INTERNAL_ACCEPTABLE_FILL_MAX
        and slot_fit_waste <= SLOT_FIT_WASTE_ACCEPTABLE_MAX
        and not metrics["background_or_edge_artifact_risk"]
    )
    metrics["needs_regeneration"] = (
        fill < INTERNAL_REGENERATE_FILL_THRESHOLD
        or fill > INTERNAL_ACCEPTABLE_FILL_MAX
        or slot_fit_waste > SLOT_FIT_WASTE_ACCEPTABLE_MAX
        or metrics["background_or_edge_artifact_risk"]
    )
    return metrics


def select_best_candidate(candidate_metrics: list[dict]) -> dict | None:
    if not candidate_metrics:
        return None
    slot_fit_candidates = [
        item for item in candidate_metrics
        if float(item.get("slot_fit_waste_percent") or 1.0) <= SLOT_FIT_WASTE_ACCEPTABLE_MAX
    ]
    if slot_fit_candidates:
        artifact_free_candidates = [
            item for item in slot_fit_candidates
            if not item.get("background_or_edge_artifact_risk")
        ]
        candidates = artifact_free_candidates or slot_fit_candidates
        def slot_fit_key(item: dict) -> tuple[float, float, float, float]:
            fill = float(item.get("foreground_bbox_fill_percent") or 0.0)
            edge_touch_count = int(item.get("edge_touch_count") or 0)
            overfill_penalty = max(0.0, fill - INTERNAL_ACCEPTABLE_FILL_MAX) * 4.0
            fill_priority = min(fill, INTERNAL_ACCEPTABLE_FILL_MAX) - overfill_penalty
            return (
                1.0 if INTERNAL_ACCEPTABLE_FILL_MIN <= fill <= INTERNAL_ACCEPTABLE_FILL_MAX else 0.0,
                fill_priority,
                -float(edge_touch_count),
                float(item.get("score") or 0.0),
            )
        return max(candidates, key=slot_fit_key)
    return max(candidate_metrics, key=lambda item: item["score"])


ECONOMY_ACCEPTANCE_PROFILES = {
    "character": {"min_fill": 0.80, "max_fill": 0.95},
    "document_stack": {"min_fill": 0.75, "max_fill": 0.95},
    "tool_combo": {"min_fill": 0.75, "max_fill": 0.95},
    "chart_card": {"min_fill": 0.75, "max_fill": 0.95},
    "inspection": {"min_fill": 0.70, "max_fill": 0.95},
    "thin_tool": {"min_fill": 0.50, "max_fill": 0.95},
}


SLOT_ECONOMY_TYPES = {
    "input_text_stack": "document_stack",
    "vlm_agent_robot": "character",
    "ai_designer": "character",
    "ai_critic": "character",
    "synthesis_tools": "tool_combo",
    "erase_text_tool": "thin_tool",
    "ocr_verify": "inspection",
    "final_autofigure_card": "chart_card",
}


ASSET_TYPE_PROMPT_TEMPLATES = {
    "character": {
        "subject_framing": "large bust or upper-body character framing",
        "composition": "The character must be zoomed in so the head and upper body dominate the canvas, with the foreground bbox spanning at least 80% of both canvas width and canvas height. Do not show a small full-body figure floating in the center.",
        "allowed_auxiliary_objects": "Only one small carried prop is allowed when it identifies the role.",
    },
    "document_stack": {
        "subject_framing": "large stacked-document framing",
        "composition": "The document stack must be oversized and zoomed in, with page corners close to all four safe margins and the stack spanning at least 80% of canvas width and height.",
        "allowed_auxiliary_objects": "No auxiliary objects.",
    },
    "tool_combo": {
        "subject_framing": "compact grouped tool framing",
        "composition": "Make the wand and palette large, overlapping, and tightly grouped into one foreground cluster that spans most of the landscape canvas. Do not leave the tools as tiny separated objects.",
        "allowed_auxiliary_objects": "Only tiny sparkles attached to the main cluster are allowed.",
    },
    "chart_card": {
        "subject_framing": "large card-object framing",
        "composition": "The chart card rectangle itself should be large, with card edges close to the safe margins and internal chart marks large enough to count as foreground.",
        "allowed_auxiliary_objects": "No objects outside the chart card.",
    },
    "inspection": {
        "subject_framing": "large compact inspection-symbol framing",
        "composition": "The magnifying lens should be oversized and the handle and verification badge should be tightly grouped into one compact foreground silhouette that fills the 4:3 canvas.",
        "allowed_auxiliary_objects": "Only short side lines close to the lens are allowed.",
    },
    "thin_tool": {
        "subject_framing": "oversized single-tool framing",
        "composition": "Use a broad front-facing version of the tool when possible so it fills the canvas visually without relying on scattered particles.",
        "allowed_auxiliary_objects": "No auxiliary objects.",
    },
}


def build_asset_generation_spec(slot: dict, candidate_index: int = 1) -> dict:
    slot_type = SLOT_ECONOMY_TYPES.get(slot["id"], "character")
    profile = ECONOMY_ACCEPTANCE_PROFILES[slot_type]
    template = ASSET_TYPE_PROMPT_TEMPLATES[slot_type]
    aspect_ratio = choose_generation_aspect_ratio(slot)
    return {
        "slot_id": slot["id"],
        "object_type": slot_type,
        "target_aspect_ratio": aspect_ratio,
        "slot_aspect_ratio": slot_aspect_ratio(slot),
        "target_fill_range": [profile["min_fill"], profile["max_fill"]],
        "target_fill_percent_text": f"{int(profile['min_fill'] * 100)}-{int(profile['max_fill'] * 100)}%",
        "max_margin_percent": 0.12,
        "background_color_hex": slot["background_color_hex"],
        "subject_framing": template["subject_framing"],
        "composition": template["composition"],
        "allowed_auxiliary_objects": template["allowed_auxiliary_objects"],
        "main_request": slot["prompt"],
        "slot_guidance": slot.get("fill_guidance", ""),
        "forbidden": [
            "readable text",
            "letters",
            "numbers",
            "labels",
            "watermark",
            "large empty background",
            "wide margins",
            "distant detached decorations",
            "black or transparent corner artifacts",
        ],
        "candidate_index": candidate_index,
    }


def prompt_from_asset_generation_spec(spec: dict) -> str:
    forbidden = ", ".join(spec["forbidden"])
    return (
        "Use the provided crop only as a visual reference for object identity and style, not for its scale, framing, empty margins, or amount of background. "
        "Before generating, internally plan the icon layout from this generation specification and satisfy every framing constraint. "
        f"Asset id: {spec['slot_id']}. Object type: {spec['object_type']}. Main request: {spec['main_request']}. "
        f"Canvas aspect ratio: {spec['target_aspect_ratio']}. Target effective foreground fill: {spec['target_fill_percent_text']}. "
        f"Maximum intended uniform background margin: {int(spec['max_margin_percent'] * 100)}%. "
        f"Subject framing: {spec['subject_framing']}. Composition: {spec['composition']} "
        f"Auxiliary object rule: {spec['allowed_auxiliary_objects']} "
        f"{spec['slot_guidance']} "
        "Preserve the crop's object identity, approximate pose, pastel academic illustration style, soft shadows, rounded friendly shapes, and color family. "
        "Create a standalone PowerPoint scientific architecture icon where the main subject fills the image canvas according to the target fill range. "
        "The background should be only a narrow border around the subject, not a large field. Prefer a close-up crop-like composition over a full-scene composition. "
        "If the first draft would contain broad empty background, internally zoom in before finalizing the image. "
        "Keep all important details visible inside the canvas and avoid cutting off the subject. "
        f"Use a perfectly flat background color matching {spec['background_color_hex']} so the asset blends into its stage panel. "
        f"Do not include: {forbidden}. No extra unrelated objects."
    )


def economy_acceptance_decision(slot: dict, metrics: dict) -> dict:
    slot_type = SLOT_ECONOMY_TYPES.get(slot["id"], "character")
    profile = ECONOMY_ACCEPTANCE_PROFILES[slot_type]
    fill = float(metrics.get("foreground_bbox_fill_percent") or 0.0)
    slot_fit_waste = float(metrics.get("slot_fit_waste_percent") or 1.0)
    artifact_risk = bool(metrics.get("background_or_edge_artifact_risk"))
    accepted = (
        profile["min_fill"] <= fill <= profile["max_fill"]
        and slot_fit_waste <= SLOT_FIT_WASTE_ACCEPTABLE_MAX
        and not artifact_risk
    )
    if accepted:
        reason = f"{slot_type}_threshold_met"
    elif artifact_risk:
        reason = "reject_background_or_edge_artifact"
    elif slot_fit_waste > SLOT_FIT_WASTE_ACCEPTABLE_MAX:
        reason = "reject_slot_fit_waste"
    elif fill < profile["min_fill"]:
        reason = f"below_{slot_type}_min_fill"
    else:
        reason = f"above_{slot_type}_max_fill"
    return {
        "slot_type": slot_type,
        "min_fill": profile["min_fill"],
        "max_fill": profile["max_fill"],
        "accepted": accepted,
        "reason": reason,
    }


def run_asset_cost_experiment(retries: int, output_path: Path) -> dict:
    crop_reference_slots()
    existing_report_items: dict[str, dict] = {}
    generation_report_path = OUT / "asset_generation_report.json"
    if generation_report_path.exists():
        try:
            generation_report = json.loads(generation_report_path.read_text(encoding="utf-8"))
            existing_report_items = {item["slot_id"]: item for item in generation_report.get("assets", [])}
        except Exception:
            existing_report_items = {}
    items = []
    strict_max_requests = 0
    economy_max_requests = 0
    for slot in SLOTS:
        report_item = existing_report_items.get(slot["id"])
        if report_item and report_item.get("candidate_metrics"):
            candidate_metrics = report_item["candidate_metrics"]
        else:
            candidate_metrics = [evaluate_candidate(slot, path) for path in candidate_paths_for_slot(slot["id"])]
        best = report_item.get("selected_internal_fill") if report_item and report_item.get("selected_internal_fill") else select_best_candidate(candidate_metrics)
        if best is None:
            strict_would_generate = True
            economy_would_generate = True
            economy_decision = {"accepted": False, "reason": "no_existing_candidate", "slot_type": SLOT_ECONOMY_TYPES.get(slot["id"], "character")}
            selected = {}
        else:
            strict_would_generate = not best["pass"] or best["needs_regeneration"]
            economy_decision = economy_acceptance_decision(slot, best)
            economy_would_generate = not economy_decision["accepted"]
            selected = {
                "candidate_path": best.get("candidate_path"),
                "foreground_bbox_fill_percent": best.get("foreground_bbox_fill_percent"),
                "slot_fit_waste_percent": best.get("slot_fit_waste_percent"),
                "generation_aspect_ratio": best.get("generation_aspect_ratio"),
                "generated_image_size": best.get("generated_image_size"),
                "background_or_edge_artifact_risk": best.get("background_or_edge_artifact_risk"),
                "strict_pass": best.get("pass"),
                "strict_needs_regeneration": best.get("needs_regeneration"),
            }
        if strict_would_generate:
            strict_max_requests += retries
        if economy_would_generate:
            economy_max_requests += 1
        items.append({
            "slot_id": slot["id"],
            "candidate_count": len(candidate_metrics),
            "selected": selected,
            "current_strict_policy": {
                "would_generate_next_run": strict_would_generate,
                "max_api_requests_next_run": retries if strict_would_generate else 0,
            },
            "economy_policy": {
                **economy_decision,
                "would_generate_next_run": economy_would_generate,
                "max_api_requests_next_run": 1 if economy_would_generate else 0,
            },
        })
    report = {
        "summary": "Dry-run cost experiment. No image API calls are made; this compares next-run generation decisions using existing candidates.",
        "baseline_policy": {
            "name": "current_strict_80_95_policy",
            "max_candidates_per_failed_slot": retries,
        },
        "economy_policy": {
            "name": "type_aware_acceptance_with_one_candidate_retry",
            "profiles": ECONOMY_ACCEPTANCE_PROFILES,
            "max_candidates_per_failed_slot": 1,
            "locks_existing_accepted_assets": True,
        },
        "estimated_next_run_api_requests": {
            "current_strict_policy_max": strict_max_requests,
            "economy_policy_max": economy_max_requests,
            "estimated_requests_saved": strict_max_requests - economy_max_requests,
            "estimated_saving_percent": round((strict_max_requests - economy_max_requests) / max(strict_max_requests, 1) * 100, 2),
        },
        "items": items,
    }
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(json.dumps(report, indent=2, ensure_ascii=False), encoding="utf-8")
    return report


def parse_slot_ids(slot_ids: str) -> list[str]:
    if slot_ids.strip().lower() == "all":
        return [slot["id"] for slot in SLOTS]
    return [item.strip() for item in slot_ids.split(",") if item.strip()]


def create_preflight_contact_sheet(items: list[dict], output_path: Path) -> None:
    thumb = 180
    label_h = 64
    gap = 16
    columns = 4
    rows = max(1, math.ceil(len(items) / columns))
    sheet = Image.new("RGB", (columns * (thumb + gap) + gap, rows * (thumb + label_h + gap) + gap), "#F3F1EC")
    draw = ImageDraw.Draw(sheet)
    for idx, item in enumerate(items):
        col = idx % columns
        row = idx // columns
        x0 = gap + col * (thumb + gap)
        y0 = gap + row * (thumb + label_h + gap)
        if item.get("candidate_path") and Path(item["candidate_path"]).exists():
            image = Image.open(item["candidate_path"]).convert("RGB")
            image.thumbnail((thumb, thumb), Image.Resampling.LANCZOS)
            sheet.paste(image, (x0 + (thumb - image.width) // 2, y0 + (thumb - image.height) // 2))
        fill = item.get("metrics", {}).get("foreground_bbox_fill_percent")
        accepted = item.get("economy_decision", {}).get("accepted")
        label = f"{item['slot_id']}\nfill: {fill * 100:.1f}%" if isinstance(fill, (int, float)) else item["slot_id"]
        label += f"\n{'accepted' if accepted else 'review'}"
        draw.multiline_text((x0, y0 + thumb + 6), label, fill="#333333", spacing=3)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    sheet.save(output_path)


def run_preflight_first_pass_experiment(slot_ids: str, output_dir: Path, workers: int = 4) -> dict:
    crop_reference_slots()
    output_dir.mkdir(parents=True, exist_ok=True)
    candidate_dir = output_dir / "candidates"
    selected_ids = parse_slot_ids(slot_ids)
    selected_slots = [slot_by_id(slot_id) for slot_id in selected_ids]
    specs = [build_asset_generation_spec(slot, 1) for slot in selected_slots]
    (output_dir / "asset_generation_specs.json").write_text(json.dumps({
        "summary": "Preflight first-pass generation specs. These specs are used before any image API request.",
        "slots": specs,
    }, indent=2, ensure_ascii=False), encoding="utf-8")

    def run_one(slot: dict) -> dict:
        try:
            candidate_path = gemini_generate_from_crop(slot, 1, candidate_dir=candidate_dir)
            background_report = harmonize_asset_background(candidate_path, slot["background_color_hex"])
            metrics = evaluate_candidate(slot, candidate_path)
            economy_decision = economy_acceptance_decision(slot, metrics)
            return {
                "slot_id": slot["id"],
                "status": "ok",
                "candidate_path": str(candidate_path),
                "asset_generation_spec": build_asset_generation_spec(slot, 1),
                "metrics": metrics,
                "economy_decision": economy_decision,
                "strict_decision": {
                    "accepted": bool(metrics.get("pass")),
                    "needs_regeneration": bool(metrics.get("needs_regeneration")),
                },
                "background": background_report,
            }
        except Exception as exc:
            return {
                "slot_id": slot["id"],
                "status": "generation_failed",
                "asset_generation_spec": build_asset_generation_spec(slot, 1),
                "error": str(exc),
                "economy_decision": {"accepted": False, "reason": "generation_failed"},
                "strict_decision": {"accepted": False, "needs_regeneration": True},
            }

    items = []
    with ThreadPoolExecutor(max_workers=workers) as executor:
        futures = {executor.submit(run_one, slot): slot for slot in selected_slots}
        for future in as_completed(futures):
            items.append(future.result())
    items.sort(key=lambda item: item["slot_id"])
    accepted = [item for item in items if item.get("economy_decision", {}).get("accepted")]
    strict_accepted = [item for item in items if item.get("strict_decision", {}).get("accepted")]
    contact_sheet = output_dir / "preflight_first_pass_contact_sheet.png"
    create_preflight_contact_sheet(items, contact_sheet)
    report = {
        "summary": "One-shot preflight generation experiment. Each selected slot gets one planned API generation in an isolated experiment directory; main assets and PPT are not overwritten.",
        "selected_slots": selected_ids,
        "workers": workers,
        "main_assets_overwritten": False,
        "main_ppt_overwritten": False,
        "contact_sheet": str(contact_sheet),
        "result": {
            "generated_count": len([item for item in items if item["status"] == "ok"]),
            "failed_count": len([item for item in items if item["status"] != "ok"]),
            "economy_first_pass_accept_count": len(accepted),
            "economy_first_pass_accept_rate": round(len(accepted) / max(len(items), 1), 4),
            "strict_first_pass_accept_count": len(strict_accepted),
            "strict_first_pass_accept_rate": round(len(strict_accepted) / max(len(items), 1), 4),
            "estimated_requests_if_one_shot_then_manual_review": len(items),
            "estimated_requests_if_strict_five_retry_for_nonpass": len(items) + (len(items) - len(strict_accepted)) * 4,
        },
        "items": items,
    }
    (output_dir / "preflight_first_pass_experiment_report.json").write_text(json.dumps(report, indent=2, ensure_ascii=False), encoding="utf-8")
    return report


def edge_connected_background_mask(image: Image.Image, background: tuple[int, int, int], tolerance: int = 42) -> Image.Image:
    rgb_image = image.convert("RGB")
    width, height = rgb_image.size
    pixels = rgb_image.load()
    visited = bytearray(width * height)
    mask = Image.new("L", (width, height), 0)
    mask_pixels = mask.load()
    queue: deque[tuple[int, int]] = deque()

    def close_to_background(px_: int, py_: int) -> bool:
        color = pixels[px_, py_]
        return abs(color[0] - background[0]) + abs(color[1] - background[1]) + abs(color[2] - background[2]) <= tolerance

    def push(px_: int, py_: int) -> None:
        index = py_ * width + px_
        if visited[index] or not close_to_background(px_, py_):
            return
        visited[index] = 1
        mask_pixels[px_, py_] = 255
        queue.append((px_, py_))

    for px_ in range(width):
        push(px_, 0)
        push(px_, height - 1)
    for py_ in range(height):
        push(0, py_)
        push(width - 1, py_)

    while queue:
        px_, py_ = queue.popleft()
        if px_ > 0:
            push(px_ - 1, py_)
        if px_ + 1 < width:
            push(px_ + 1, py_)
        if py_ > 0:
            push(px_, py_ - 1)
        if py_ + 1 < height:
            push(px_, py_ + 1)
    return mask


def latest_candidate_for_slot(slot_id: str) -> Path | None:
    slot_dir = CANDIDATE_DIR / slot_id
    if not slot_dir.exists():
        return None
    candidates = sorted(slot_dir.glob("candidate_*.png"))
    return candidates[-1] if candidates else None


def candidate_paths_for_slot(slot_id: str) -> list[Path]:
    slot_dir = CANDIDATE_DIR / slot_id
    if not slot_dir.exists():
        return []
    return sorted(slot_dir.glob("candidate_*.png"))


def candidate_meta_path(candidate_path: Path) -> Path:
    return candidate_path.with_suffix(".generation.json")


def slot_aspect_ratio(slot: dict) -> float:
    left, top, right, bottom = slot["asset_bbox"]
    return round((right - left) / max(bottom - top, 1), 4)


def choose_generation_aspect_ratio(slot: dict) -> str:
    override = slot.get("generation_aspect_ratio_override")
    if override:
        return override
    ratio = slot_aspect_ratio(slot)
    if ratio >= 1.55:
        return "16:9"
    if ratio >= 1.18:
        return "4:3"
    if ratio >= 0.85:
        return "1:1"
    if ratio >= 0.62:
        return "3:4"
    return "9:16"


def aspect_ratio_prompt_instruction(aspect_ratio: str) -> str:
    if aspect_ratio in {"16:9", "4:3"}:
        return (
            f"Use a {aspect_ratio} landscape canvas matched to the target slot. "
            "Compose the subject horizontally so it naturally spans the canvas width without needing rotation, stacking, or artificial square filling. "
            "Keep the top and bottom margins narrow but preserve all important details. "
        )
    if aspect_ratio in {"3:4", "9:16"}:
        return (
            f"Use a {aspect_ratio} portrait canvas matched to the target slot. "
            "Compose the subject vertically so it fills the height while keeping left and right margins narrow. "
            "Use a bust or full-height framing as appropriate, but avoid a tiny centered character with broad side margins. "
        )
    return (
        "Use a square 1:1 canvas. "
        "Compose the subject to fill the square naturally with narrow, even margins. "
    )


def compute_slot_fit_metrics(image_path: Path, slot: dict) -> dict:
    with Image.open(image_path) as image:
        image_width, image_height = image.size
    image_ratio = image_width / max(image_height, 1)
    slot_ratio = slot_aspect_ratio(slot)
    left, top, right, bottom = slot["asset_bbox"]
    slot_width = right - left
    slot_height = bottom - top
    if image_ratio > slot_ratio:
        fit_width = slot_width
        fit_height = slot_width / max(image_ratio, 0.0001)
    else:
        fit_height = slot_height
        fit_width = slot_height * image_ratio
    used_area = fit_width * fit_height
    slot_area = max(slot_width * slot_height, 1)
    waste = max(0.0, 1.0 - used_area / slot_area)
    target_ratio = GENERATION_ASPECT_RATIO_VALUES[choose_generation_aspect_ratio(slot)]
    ratio_error = abs(math.log(max(image_ratio, 0.0001) / target_ratio))
    if waste <= SLOT_FIT_WASTE_ACCEPTABLE_MAX and ratio_error <= 0.18:
        status = "matched"
    elif waste <= SLOT_FIT_WASTE_ACCEPTABLE_MAX:
        status = "slot_fit_ok_generated_ratio_drift"
    else:
        status = "mismatch"
    generation_meta = {}
    meta_path = candidate_meta_path(image_path)
    if meta_path.exists():
        try:
            generation_meta = json.loads(meta_path.read_text(encoding="utf-8"))
        except Exception:
            generation_meta = {"generation_meta_warning": "failed_to_read_generation_meta"}
    return {
        "slot_aspect_ratio": slot_ratio,
        "generation_aspect_ratio": choose_generation_aspect_ratio(slot),
        "requested_generation_aspect_ratio": generation_meta.get("requested_generation_aspect_ratio", choose_generation_aspect_ratio(slot)),
        "actual_generation_request_aspect_ratio": generation_meta.get("actual_generation_request_aspect_ratio"),
        "generation_aspect_ratio_fallback_used": generation_meta.get("generation_aspect_ratio_fallback_used", False),
        "generated_image_size": [image_width, image_height],
        "generated_image_aspect_ratio": round(image_ratio, 4),
        "slot_fit_waste_percent": round(waste, 4),
        "aspect_ratio_match_status": status,
        **({"generation_meta_warning": generation_meta["generation_meta_warning"]} if "generation_meta_warning" in generation_meta else {}),
    }


def gemini_generate_from_crop(slot: dict, candidate_index: int = 1, candidate_dir: Path | None = None) -> Path:
    api_key = os.getenv("GEMINI_API_KEY") or os.getenv("API_KEY")
    url = os.getenv("GEMINI_GEN_IMG_URL")
    if not api_key or not url:
        raise RuntimeError("GEMINI_API_KEY/API_KEY and GEMINI_GEN_IMG_URL are required")
    crop_path = Path(slot["crop_path"])
    mime = mimetypes.guess_type(crop_path.name)[0] or "image/png"
    image_b64 = base64.b64encode(crop_path.read_bytes()).decode("ascii")
    spec = build_asset_generation_spec(slot, candidate_index)
    aspect_ratio = spec["target_aspect_ratio"]
    prompt = prompt_from_asset_generation_spec(spec)
    payload = {
        "contents": [{
            "role": "user",
            "parts": [
                {"text": prompt},
                {"inlineData": {"mimeType": mime, "data": image_b64}},
            ],
        }],
        "generationConfig": {
            "responseModalities": ["IMAGE"],
            "imageConfig": {"imageSize": "1K", "aspectRatio": aspect_ratio},
        },
    }
    candidate_root = candidate_dir or CANDIDATE_DIR
    out_path = candidate_root / slot["id"] / f"candidate_{candidate_index:02d}.png"
    out_path.parent.mkdir(parents=True, exist_ok=True)
    request_aspects = [aspect_ratio]
    if aspect_ratio != "1:1":
        request_aspects.append("1:1")
    last_error: Exception | None = None
    for request_aspect in request_aspects:
        payload["generationConfig"]["imageConfig"]["aspectRatio"] = request_aspect
        try:
            response = requests.post(url, headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}, data=json.dumps(payload), timeout=240)
            response.raise_for_status()
        except requests.HTTPError as exc:
            last_error = exc
            if request_aspect != "1:1":
                continue
            raise
        data = response.json()
        for candidate in data.get("candidates", []):
            for part in candidate.get("content", {}).get("parts", []):
                inline = part.get("inlineData") or part.get("inline_data")
                if inline and inline.get("data"):
                    out_path.write_bytes(base64.b64decode(inline["data"]))
                    candidate_meta_path(out_path).write_text(json.dumps({
                        "slot_id": slot["id"],
                        "requested_generation_aspect_ratio": aspect_ratio,
                        "actual_generation_request_aspect_ratio": request_aspect,
                        "generation_aspect_ratio_fallback_used": request_aspect != aspect_ratio,
                        "asset_generation_spec": spec,
                    }, indent=2, ensure_ascii=False), encoding="utf-8")
                    return out_path
    if last_error:
        raise last_error
    raise RuntimeError(f"No image returned for {slot['id']}")


def generate_one_asset_with_retries(slot: dict, retries: int = 1, economy_mode: bool = True) -> dict:
    final_path = ASSET_DIR / f"{slot['id']}.png"
    errors = []
    candidate_metrics: list[dict] = []
    existing_candidates = candidate_paths_for_slot(slot["id"]) if os.getenv("RFS_REUSE_ASSETS", "1") == "1" else []
    for candidate_path in existing_candidates:
        candidate_metrics.append(evaluate_candidate(slot, candidate_path))

    best_existing = select_best_candidate(candidate_metrics)
    existing_economy_decision = economy_acceptance_decision(slot, best_existing) if best_existing else {"accepted": False, "reason": "no_existing_candidate"}
    should_generate_more = (
        best_existing is None
        or (not existing_economy_decision["accepted"] if economy_mode else (not best_existing["pass"] or best_existing["needs_regeneration"]))
    )
    next_index = len(existing_candidates) + 1
    api_requests_attempted = 0
    if should_generate_more:
        for attempt in range(next_index, next_index + retries):
            try:
                api_requests_attempted += 1
                candidate_path = gemini_generate_from_crop(slot, attempt)
                candidate_metrics.append(evaluate_candidate(slot, candidate_path))
            except Exception as exc:
                errors.append(str(exc))
                time.sleep(min(2 * (attempt - next_index + 1), 6))
            current_best = select_best_candidate(candidate_metrics)
            if current_best and (economy_acceptance_decision(slot, current_best)["accepted"] if economy_mode else current_best["pass"]):
                break

    best = select_best_candidate(candidate_metrics)
    if best is not None:
        candidate_path = Path(best["candidate_path"])
        shutil.copyfile(candidate_path, final_path)
        background_report = harmonize_asset_background(final_path, slot["background_color_hex"])
        fill = float(best.get("foreground_bbox_fill_percent") or 0.0)
        if best["pass"]:
            selected_reason = "selected_candidate_with_internal_fill_80_95_percent"
        elif fill < INTERNAL_ACCEPTABLE_FILL_MIN:
            selected_reason = "selected_best_available_but_subject_still_too_small"
        else:
            selected_reason = "selected_best_available_but_possible_edge_cutoff_risk"
        economy_decision = economy_acceptance_decision(slot, best)
        if economy_mode and economy_decision["accepted"] and not best["pass"]:
            selected_reason = f"selected_by_economy_policy_{economy_decision['reason']}"
        return {
            "slot_id": slot["id"],
            "status": "ok" if best["pass"] else ("ok_by_economy_policy" if economy_decision["accepted"] else "ok_with_internal_fill_warning"),
            "asset_source": "candidate_selected",
            "candidate": str(candidate_path),
            "asset": str(final_path),
            "candidate_metrics": sorted(candidate_metrics, key=lambda item: item["score"], reverse=True),
            "selected_internal_fill": {**best, "selected_reason": selected_reason},
            "economy_decision": economy_decision,
            "economy_mode": economy_mode,
            "api_requests_attempted": api_requests_attempted,
            "background": background_report,
            **({"errors": errors} if errors else {}),
        }

    fallback = ASSET_DIR / f"{slot['id']}.png"
    # Keep the pipeline deliverable if API has a repeated transient failure, but make the fallback explicit.
    Image.open(slot["crop_path"]).convert("RGB").save(fallback)
    background_report = harmonize_asset_background(fallback, slot["background_color_hex"])
    fallback_metrics = evaluate_candidate(slot, fallback)
    return {
        "slot_id": slot["id"],
        "status": "fallback_to_reference_crop",
        "errors": errors,
        "asset": str(fallback),
        "selected_internal_fill": {**fallback_metrics, "selected_reason": "fallback_reference_crop_after_generation_failure"},
        "economy_decision": economy_acceptance_decision(slot, fallback_metrics),
        "economy_mode": economy_mode,
        "api_requests_attempted": api_requests_attempted,
        "background": background_report,
    }


def generate_assets_parallel(workers: int = 6, retries: int = 1, economy_mode: bool = True) -> dict:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    CANDIDATE_DIR.mkdir(parents=True, exist_ok=True)
    reports = []
    with ThreadPoolExecutor(max_workers=workers) as executor:
        futures = {executor.submit(generate_one_asset_with_retries, slot, retries, economy_mode): slot for slot in SLOTS}
        for future in as_completed(futures):
            slot = futures[future]
            try:
                reports.append(future.result())
            except Exception as exc:
                fallback = ASSET_DIR / f"{slot['id']}.png"
                Image.open(slot["crop_path"]).convert("RGB").save(fallback)
                background_report = harmonize_asset_background(fallback, slot["background_color_hex"])
                fallback_metrics = evaluate_candidate(slot, fallback)
                reports.append({
                    "slot_id": slot["id"],
                    "status": "fallback_to_reference_crop",
                    "errors": [str(exc)],
                    "asset": str(fallback),
                    "selected_internal_fill": {**fallback_metrics, "selected_reason": "fallback_reference_crop_after_parallel_generation_failure"},
                    "economy_decision": economy_acceptance_decision(slot, fallback_metrics),
                    "economy_mode": economy_mode,
                    "api_requests_attempted": 0,
                    "background": background_report,
                })
    return {
        "summary": "AI asset generation report for AutoFigure Architecture rebuild.",
        "workers": workers,
        "retries": retries,
        "economy_mode": economy_mode,
        "api_requests_attempted": sum(int(item.get("api_requests_attempted") or 0) for item in reports),
        "assets": sorted(reports, key=lambda item: item["slot_id"]),
    }


def add_picture_contain(slide, image_path: Path, left_px, top_px, width_px, height_px):
    with Image.open(image_path) as img:
        ratio = img.width / max(img.height, 1)
    box_ratio = width_px / max(height_px, 1)
    if ratio > box_ratio:
        fit_w = width_px
        fit_h = width_px / ratio
    else:
        fit_h = height_px
        fit_w = height_px * ratio
    return slide.shapes.add_picture(str(image_path), Inches(x(left_px + (width_px - fit_w) / 2)), Inches(y(top_px + (height_px - fit_h) / 2)), width=Inches(w(fit_w)), height=Inches(h(fit_h)))


def draw_document_blueprint(slide):
    box = add_round(slide, 74, 514, 472, 248, "#F4F8FB", "#6D8BA4", 1.3)
    box.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    add_text(slide, "Initial Blueprint (S0, A0)", x(160), y(528), w(270), h(34), style="stage_title", color="#111111", size=12.5)
    code = '{\n  "nodes": [{\n    id: "v1",\n    pos: [x,y]...\n  }],\n  "style": "academic_flat"\n}'
    add_text(slide, code, x(98), y(574), w(220), h(150), style="small", color="#111111", align=PP_ALIGN.LEFT, size=8.3)
    add_line(slide, 336, 656, 421, 656, color="#A6C8DE", width_pt=2.0)
    brace = slide.shapes.add_shape(MSO_SHAPE.RIGHT_BRACE, Inches(x(388)), Inches(y(610)), Inches(w(55)), Inches(h(112)))
    no_fill(brace)
    set_line(brace, "#9FC6DE", 1.6)
    add_round(slide, 424, 585, 112, 54, "#E9EEF1", "#7A8790", 1.0)
    add_text(slide, "Node A", x(433), y(595), w(90), h(28), style="body", color="#333333", size=10.5)
    add_round(slide, 424, 676, 112, 54, "#E9EEF1", "#7A8790", 1.0)
    add_text(slide, "Node B", x(433), y(686), w(90), h(28), style="body", color="#333333", size=10.5)


def add_panel(slide, left, top, width, height, title, border, fill, header):
    add_round(slide, left, top, width, height, fill, border, 1.8)
    header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x(left + 2)), Inches(y(top + 2)), Inches(w(width - 4)), Inches(h(62)))
    set_fill(header_shape, header, 0)
    set_line(header_shape, header, 0.1)
    add_text(slide, title, x(left + 16), y(top + 16), w(width - 32), h(36), style="stage_title", color="#333333")


def draw_ppt() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("#F3F1EC")

    add_text(slide, "AutoFigure Architecture", x(50), y(45), w(620), h(60), style="title", color="#333333", align=PP_ALIGN.LEFT)
    add_text(slide, "Decoupled Generative Paradigm: Reasoning, Refinement, and Layered Rendering", x(51), y(114), w(760), h(30), style="subtitle", color="#333333", align=PP_ALIGN.LEFT)

    add_panel(slide, 42, 170, 540, 650, "Stage I: Concept Extraction", "#8AA4BC", "#EAF3FF", "#D6E5F4")
    add_panel(slide, 615, 170, 610, 658, "Stage II: Critique-and-Refine", "#E1A45D", "#FFF1DF", "#F8D1A3")
    add_panel(slide, 1260, 48, 500, 892, "Stage III: Rendering Strategy", "#84A57F", "#ECF7E9", "#C9DFC1")

    asset = lambda name: ASSET_DIR / f"{name}.png"
    add_picture_contain(slide, asset("input_text_stack"), 70, 258, 145, 148)
    add_text(slide, "Input Text", x(92), y(421), w(105), h(26), style="label", color="#333333")
    render_controls(slide, ["input_to_vlm"])
    add_picture_contain(slide, asset("vlm_agent_robot"), 285, 252, 165, 160)
    add_text(slide, "VLM Agent", x(309), y(421), w(118), h(26), style="label", color="#333333")
    add_round(slide, 465, 258, 90, 90, "#EAF0F5", "#94A6B9", 1.2)
    add_text(slide, "Entities", x(476), y(293), w(68), h(24), style="body", color="#333333", bold=True)
    add_round(slide, 465, 355, 90, 90, "#EAF0F5", "#94A6B9", 1.2)
    add_text(slide, "Relations", x(474), y(389), w(72), h(24), style="body", color="#333333", bold=True)
    render_controls(slide, ["vlm_to_blueprint"])
    draw_document_blueprint(slide)
    add_text(slide, '"Map unstructured text to symbolic graph"', x(84), y(774), w(415), h(32), style="body", color="#111111", size=11.0)
    render_controls(slide, ["stage_i_to_stage_ii"])

    # Stage II
    add_round(slide, 700, 260, 440, 50, "#FFE7C8", "#D99045", 1.2)
    add_text(slide, "Critique: Alignment, Overlap, Balance", x(720), y(274), w(400), h(24), style="callout", color="#333333", size=10.6)
    add_round(slide, 638, 350, 104, 70, "#FFE7C8", "#D99045", 1.1)
    add_text(slide, "Refine\nLayout", x(657), y(360), w(64), h(42), style="callout", color="#333333", size=8.5)
    add_round(slide, 1092, 350, 110, 70, "#FFE7C8", "#D99045", 1.1)
    add_text(slide, "Feedback\nF(i)", x(1108), y(360), w(80), h(42), style="callout", color="#333333", size=8.5)
    render_controls(slide, ["critique_refinement_loop"])
    add_picture_contain(slide, asset("ai_designer"), 650, 433, 145, 185)
    add_text(slide, "AI Designer", x(670), y(642), w(120), h(24), style="label", color="#333333", size=8.7)
    add_picture_contain(slide, asset("ai_critic"), 1060, 430, 145, 190)
    add_text(slide, "AI Critic", x(1098), y(642), w(86), h(24), style="label", color="#333333", size=8.7)
    add_round(slide, 802, 480, 235, 90, "#FFE7C8", "#D99045", 1.3)
    add_text(slide, "Score Comparison", x(822), y(500), w(194), h(30), style="callout", color="#333333", size=10.0)
    add_text(slide, "q(cand)  >  q(best)", x(822), y(532), w(194), h(24), style="small", color="#333333", size=8.0)
    render_controls(slide, ["designer_refine_arrow", "critic_feedback_arrow"])
    add_round(slide, 690, 748, 460, 48, "#FFE7C8", "#D99045", 1.2)
    add_text(slide, "Update: Re-interpret method & improve", x(710), y(761), w(420), h(22), style="callout", color="#333333", size=10.5)
    render_controls(slide, ["stage_ii_to_stage_iii"])

    # Stage III
    add_picture_contain(slide, asset("synthesis_tools"), 1380, 128, 230, 120)
    add_text(slide, "Synthesis", x(1455), y(251), w(120), h(28), style="label", color="#333333")
    render_controls(slide, ["synthesis_to_raw"])
    add_round(slide, 1428, 346, 160, 108, "#F4F4EF", "#B5B5AC", 1.3)
    add_text(slide, "Raw Image", x(1450), y(364), w(116), h(26), style="label", color="#333333", size=10.2)
    add_text(slide, "Text Blur!", x(1460), y(404), w(96), h(24), style="callout", color="#E24D42", size=9.2)
    add_picture_contain(slide, asset("erase_text_tool"), 1290, 486, 140, 125)
    add_text(slide, "Erase Text", x(1320), y(595), w(108), h(28), style="label", color="#333333", size=9.5)
    add_text(slide, "ABC", x(1404), y(516), w(55), h(24), style="label", color="#333333", size=9.0)
    add_picture_contain(slide, asset("ocr_verify"), 1570, 490, 150, 125)
    add_text(slide, "OCR + Verify", x(1578), y(595), w(140), h(28), style="label", color="#333333", size=9.5)
    add_picture_contain(slide, asset("final_autofigure_card"), 1410, 684, 185, 128)
    add_text(slide, "Final AutoFigure", x(1418), y(815), w(178), h(28), style="label", color="#333333", size=10.0)
    add_text(slide, "A stunning, high-resolution scientific\ndiagram with crisp details and a pleasing\ncolor scheme, ready for publication.", x(1368), y(846), w(300), h(70), style="body", color="#333333", size=8.2)
    render_controls(slide, ["raw_to_erase_route", "raw_to_ocr_route", "erase_to_final_route", "ocr_to_final_route", "finalization_arrow"])

    # Bottom legend
    add_round(slide, 42, 858, 1180, 80, "#F4F2ED", "#D2CEC8", 1.2)
    add_text(slide, "Key Methodology:", x(105), y(890), w(210), h(24), style="legend", color="#333333", align=PP_ALIGN.LEFT)
    add_round(slide, 328, 875, 280, 50, "#DFECF8", "#8AA4BC", 1.2)
    add_text(slide, "Concept Extraction", x(380), y(890), w(180), h(22), style="legend", color="#333333")
    add_round(slide, 628, 875, 280, 50, "#FFE6C9", "#D99045", 1.2)
    add_text(slide, "Critique-and-Refine", x(670), y(890), w(200), h(22), style="legend", color="#333333")
    add_round(slide, 925, 875, 280, 50, "#DCEED7", "#84A57F", 1.2)
    add_text(slide, "Rendering Strategy", x(970), y(890), w(190), h(22), style="legend", color="#333333")

    target = OUT / "editable_composition.pptx"
    try:
        prs.save(target)
        return target
    except PermissionError:
        fallback = OUT / "editable_composition_controls_calibrated.pptx"
        prs.save(fallback)
        return fallback


def export_preview(pptx_path: Path) -> None:
    try:
        import win32com.client  # type: ignore
    except Exception as exc:
        (OUT / "preview_export_error.txt").write_text(str(exc), encoding="utf-8")
        return
    app = win32com.client.Dispatch("PowerPoint.Application")
    app.Visible = 1
    presentation = app.Presentations.Open(str(pptx_path), WithWindow=False)
    try:
        presentation.Slides(1).Export(str(OUT / "rebuild_preview.png"), "PNG", REF_W, REF_H)
    finally:
        presentation.Close()
        app.Quit()


def write_metadata(asset_report: dict) -> None:
    controls = [normalize_control(control) for control in CONTROL_SPECS]
    internal_fill_report = {
        "summary": "Internal effective-content fill report for small AI icon assets. No local crop, auto-crop, PowerPoint crop, or layout enlargement is applied.",
        "policy": {
            "icon_type": "functional_icon_or_small_mark",
            "target_fill_range": [INTERNAL_ACCEPTABLE_FILL_MIN, INTERNAL_ACCEPTABLE_FILL_MAX],
            "target_fill_center": INTERNAL_CONTENT_FILL_TARGET,
            "regenerate_below": INTERNAL_REGENERATE_FILL_THRESHOLD,
            "no_crop_postprocess": True,
            "layout_bbox_locked": True,
        },
        "assets": [
            {
                "slot_id": item["slot_id"],
                **item.get("selected_internal_fill", {}),
            }
            for item in asset_report.get("assets", [])
        ],
    }
    background_report = {
        "summary": "Background harmonization report for generated/reused AI assets.",
        "policy": "match_slot_local_background",
        "assets": [
            {
                "slot_id": item["slot_id"],
                **item.get("background", {}),
            }
            for item in asset_report.get("assets", [])
        ],
    }
    ratio_fit_report = {
        "summary": "Slot-aware aspect-ratio fit report for generated icon assets. This distinguishes true small-subject problems from square-canvas letterboxing in non-square slots.",
        "policy": {
            "layout_bbox_locked": True,
            "no_default_crop": True,
            "slot_fit_waste_acceptable_max": SLOT_FIT_WASTE_ACCEPTABLE_MAX,
            "aspect_ratio_selection": {
                ">=1.55": "16:9",
                "1.18-1.55": "4:3",
                "0.85-1.18": "1:1",
                "0.62-0.85": "3:4",
                "<0.62": "9:16",
            },
        },
        "assets": [
            {
                "slot_id": item["slot_id"],
                "slot_aspect_ratio": item.get("selected_internal_fill", {}).get("slot_aspect_ratio"),
                "generation_aspect_ratio": item.get("selected_internal_fill", {}).get("generation_aspect_ratio"),
                "requested_generation_aspect_ratio": item.get("selected_internal_fill", {}).get("requested_generation_aspect_ratio"),
                "actual_generation_request_aspect_ratio": item.get("selected_internal_fill", {}).get("actual_generation_request_aspect_ratio"),
                "generation_aspect_ratio_fallback_used": item.get("selected_internal_fill", {}).get("generation_aspect_ratio_fallback_used", False),
                "generated_image_size": item.get("selected_internal_fill", {}).get("generated_image_size"),
                "generated_image_aspect_ratio": item.get("selected_internal_fill", {}).get("generated_image_aspect_ratio"),
                "slot_fit_waste_percent": item.get("selected_internal_fill", {}).get("slot_fit_waste_percent"),
                "foreground_bbox_fill_percent": item.get("selected_internal_fill", {}).get("foreground_bbox_fill_percent"),
                "aspect_ratio_match_status": item.get("selected_internal_fill", {}).get("aspect_ratio_match_status"),
                "selected_candidate": item.get("candidate"),
                "status": item.get("status"),
                "economy_decision": item.get("economy_decision"),
                "api_requests_attempted": item.get("api_requests_attempted", 0),
                "generation_errors": item.get("errors", []),
            }
            for item in asset_report.get("assets", [])
        ],
    }
    economy_report = {
        "summary": "Economy asset generation report. Accepted assets are locked and skipped; failed slots use bounded one-shot generation unless strict mode is explicitly requested.",
        "economy_mode": asset_report.get("economy_mode"),
        "retries": asset_report.get("retries"),
        "api_requests_attempted": asset_report.get("api_requests_attempted", 0),
        "profiles": ECONOMY_ACCEPTANCE_PROFILES,
        "assets": [
            {
                "slot_id": item["slot_id"],
                "status": item.get("status"),
                "api_requests_attempted": item.get("api_requests_attempted", 0),
                "economy_decision": item.get("economy_decision"),
                "selected_candidate": item.get("candidate"),
                "foreground_bbox_fill_percent": item.get("selected_internal_fill", {}).get("foreground_bbox_fill_percent"),
                "slot_fit_waste_percent": item.get("selected_internal_fill", {}).get("slot_fit_waste_percent"),
                "selected_reason": item.get("selected_internal_fill", {}).get("selected_reason"),
            }
            for item in asset_report.get("assets", [])
        ],
    }
    inventory = {
        "summary": "Slot inventory for AI-generated image assets in AutoFigure Architecture editable rebuild.",
        "reference": str(REFERENCE),
        "slots": [
            {
                "slot_id": slot["id"],
                "bbox_percent": slot["bbox_percent"],
                "asset_bbox_px": slot["asset_bbox"],
                "slot_aspect_ratio": slot_aspect_ratio(slot),
                "generation_aspect_ratio": choose_generation_aspect_ratio(slot),
                "generation_aspect_ratio_override": slot.get("generation_aspect_ratio_override"),
                "reference_crop_path": slot["crop_path"],
                "asset_path": str(ASSET_DIR / f"{slot['id']}.png"),
                "background_color_hex": slot["background_color_hex"],
                "internal_content_fill_target": slot["internal_content_fill_target"],
                "internal_max_margin_percent": slot["internal_max_margin_percent"],
                "layout_bbox_locked": slot["layout_bbox_locked"],
                "no_crop_postprocess": slot["no_crop_postprocess"],
                "prompt": slot["prompt"],
                "fill_guidance": slot.get("fill_guidance"),
            }
            for slot in SLOTS
        ],
    }
    (OUT / "slot_inventory.json").write_text(json.dumps(inventory, indent=2, ensure_ascii=False), encoding="utf-8")
    (OUT / "asset_generation_report.json").write_text(json.dumps(asset_report, indent=2, ensure_ascii=False), encoding="utf-8")
    (OUT / "asset_internal_fill_report.json").write_text(json.dumps(internal_fill_report, indent=2, ensure_ascii=False), encoding="utf-8")
    (OUT / "asset_ratio_fit_report.json").write_text(json.dumps(ratio_fit_report, indent=2, ensure_ascii=False), encoding="utf-8")
    (OUT / "asset_economy_report.json").write_text(json.dumps(economy_report, indent=2, ensure_ascii=False), encoding="utf-8")
    (OUT / "asset_background_report.json").write_text(json.dumps(background_report, indent=2, ensure_ascii=False), encoding="utf-8")
    (OUT / "reference_control_candidates.json").write_text(json.dumps({
        "summary": "Reference-derived candidate arrows, connector lines, and dashed loops for the AutoFigure Architecture rebuild.",
        "extraction_mode": "manual_reference_locked_first_pass",
        "note": "This is the structured control layer that will later be replaced/refined by image-processing or VLM control extraction.",
        "controls": controls,
    }, indent=2, ensure_ascii=False), encoding="utf-8")
    (OUT / "reference_controls.json").write_text(json.dumps({
        "summary": "Bound editable PPT control layer for arrows, connector lines, and dashed loops.",
        "controls": controls,
    }, indent=2, ensure_ascii=False), encoding="utf-8")
    (OUT / "arrow_style_profile.json").write_text(json.dumps({
        "summary": "Arrow and connector style profile.",
        "reference_image_hard_constraint": True,
        "control_localizer_mode": "manual_reference_locked_first_pass",
        "style_tokens": {
            "stage_transition": {"color": "#555555", "width_pt": 4.0, "arrow_head": "triangle"},
            "stage_i_flow": {"color": "#607998", "width_pt": 2.5, "arrow_head": "triangle"},
            "critique_loop": {"color": "#D99045", "width_pt": 2.0, "dash": True},
            "rendering_loop": {"color": "#6F9A66", "width_pt": 1.8, "dash": True},
        },
    }, indent=2, ensure_ascii=False), encoding="utf-8")
    (OUT / "selected_arrow_routes.json").write_text(json.dumps({
        "summary": "Selected editable PowerPoint routes for all reference controls.",
        "route_generation_status": "reference_locked_manual_seed",
        "routes": controls,
    }, indent=2, ensure_ascii=False), encoding="utf-8")
    (OUT / "arrow_quality_report.json").write_text(json.dumps({
        "summary": "First-pass arrow quality report.",
        "status": "improved_structured_controls",
        "known_limitations": [
            "This pass stores measured/control coordinates in JSON but does not yet run automatic computer-vision arrow extraction.",
            "Dashed closed loops are still approximated with PPT oval/polyline controls.",
        ],
        "control_count": len(controls),
        "fallback_routing_used": False,
    }, indent=2, ensure_ascii=False), encoding="utf-8")
    (OUT / "arrow_direction_report.json").write_text(json.dumps({
        "summary": "Arrow position and direction report. This file prioritizes arrow placement, target direction, anchors, and arrowhead endpoint over style details.",
        "status": "position_direction_fields_added",
        "direction_contract": {
            "start_px": "tail position in reference-image pixel coordinates",
            "end_px": "arrowhead target position in reference-image pixel coordinates",
            "direction.label": "dominant direction inferred from start_px -> end_px",
            "source_anchor": "inferred source edge from direction unless explicitly overridden",
            "target_anchor": "inferred target edge from direction unless explicitly overridden",
        },
        "arrows": [
            {
                "id": control["id"],
                "kind": control["kind"],
                "source_id": control["source_id"],
                "target_id": control["target_id"],
                "source_anchor": control.get("source_anchor"),
                "target_anchor": control.get("target_anchor"),
                "position": control.get("position"),
                "direction": control.get("direction"),
            }
            for control in controls
        ],
    }, indent=2, ensure_ascii=False), encoding="utf-8")
    (OUT / "reference_style_profile.json").write_text(json.dumps({
        "summary": "Reference style profile for AutoFigure Architecture editable rebuild.",
        "canvas_background": "#F3F1EC",
        "stage_backgrounds": {"stage_i": "#EAF3FF", "stage_ii": "#FFF1DF", "stage_iii": "#ECF7E9"},
        "asset_background_policy": "match_slot_local_background",
        "font_family": FONT,
        "text_size_scale": TEXT_SIZE_SCALE,
    }, indent=2, ensure_ascii=False), encoding="utf-8")
    (OUT / "layout_plan.json").write_text(json.dumps({
        "summary": "Reference-guided layout plan for the AutoFigure Architecture editable rebuild.",
        "canvas": {"width_px": REF_W, "height_px": REF_H, "slide_width_in": SLIDE_W, "slide_height_in": SLIDE_H},
        "panels": [
            {"id": "stage_i", "bbox_px": [42, 170, 540, 650], "fill": "#EAF3FF"},
            {"id": "stage_ii", "bbox_px": [615, 170, 610, 658], "fill": "#FFF1DF"},
            {"id": "stage_iii", "bbox_px": [1260, 48, 500, 892], "fill": "#ECF7E9"},
        ],
        "slots": inventory["slots"],
        "controls": controls,
    }, indent=2, ensure_ascii=False), encoding="utf-8")
    (OUT / "reference_text_geometry.json").write_text(json.dumps({
        "summary": "Reference text geometry seed for editable PPT text.",
        "detection_mode": "manual_reference_seed",
        "font_size_estimation": "reference bbox height converted to PPT points, then calibrated with TEXT_SIZE_SCALE",
        "text_size_scale": TEXT_SIZE_SCALE,
    }, indent=2, ensure_ascii=False), encoding="utf-8")
    (OUT / "text_program.json").write_text(json.dumps({
        "summary": "Editable text program for AutoFigure Architecture rebuild.",
        "font_family_guess": FONT,
        "fit_strategy": "bbox_bound_with_global_scale",
        "text_style_tokens": TEXT_STYLES,
        "text_size_scale": TEXT_SIZE_SCALE,
    }, indent=2, ensure_ascii=False), encoding="utf-8")
    (OUT / "text_alignment_report.json").write_text(json.dumps({
        "summary": "Text alignment report for the AutoFigure Architecture editable rebuild.",
        "status": "scaled_up",
        "scale_applied": TEXT_SIZE_SCALE,
        "reason": "Prior preview text was visually small; this pass applies a uniform calibration factor while preserving editable PPT text boxes.",
        "next_recommended_step": "render-and-measure calibration against OCR/reference text masks.",
    }, indent=2, ensure_ascii=False), encoding="utf-8")
    (OUT / "figure_program.json").write_text(json.dumps({
        "summary": "PowerPoint-first reconstruction with AI-generated local illustration assets and editable PPT structure/text.",
        "reference_size_px": [REF_W, REF_H],
        "slide_size_in": [SLIDE_W, SLIDE_H],
        "editable_layers": ["text", "stage panels", "cards", "arrows", "dashed loops", "methodology legend"],
        "raster_asset_layers": [slot["id"] for slot in SLOTS],
        "controls": controls,
        "asset_background_policy": "match_slot_local_background",
        "asset_internal_fill_policy": {
            "target_fill_range": [INTERNAL_ACCEPTABLE_FILL_MIN, INTERNAL_ACCEPTABLE_FILL_MAX],
            "target_fill_center": INTERNAL_CONTENT_FILL_TARGET,
            "no_crop_postprocess": True,
            "layout_bbox_locked": True,
        },
        "asset_aspect_ratio_policy": {
            "mode": "slot_aware_generation_canvas",
            "slot_fit_waste_acceptable_max": SLOT_FIT_WASTE_ACCEPTABLE_MAX,
            "supported_generation_aspect_ratios": list(GENERATION_ASPECT_RATIO_VALUES.keys()),
            "report": str(OUT / "asset_ratio_fit_report.json"),
        },
        "asset_economy_policy": {
            "enabled": asset_report.get("economy_mode"),
            "api_requests_attempted": asset_report.get("api_requests_attempted", 0),
            "type_aware_profiles": ECONOMY_ACCEPTANCE_PROFILES,
            "report": str(OUT / "asset_economy_report.json"),
        },
        "text_style_tokens": TEXT_STYLES,
        "text_size_scale": TEXT_SIZE_SCALE,
    }, indent=2, ensure_ascii=False), encoding="utf-8")
    (OUT / "rebuild_notes.md").write_text(
        "# Summary\n"
        "AutoFigure Architecture was rebuilt as an editable PowerPoint composition with AI-generated local illustration assets.\n\n"
        "## Editable PPT Layers\n"
        "- All titles, labels, captions, callouts, methodology legend text, panels, cards, arrows, and dashed loops are PowerPoint objects.\n\n"
        "## AI Raster Asset Layers\n"
        "- Complex illustrations such as robot characters, document stack, synthesis tools, eraser, OCR magnifier, and final figure card were generated as local image assets from reference crops.\n"
        "- Critical text is not trusted to the generated images; it is added as editable PPT text.\n\n"
        "## Generation\n"
        f"- Parallel workers: {asset_report.get('workers')}\n"
        f"- Economy mode: {asset_report.get('economy_mode')}; image API requests attempted: {asset_report.get('api_requests_attempted', 0)}\n"
        "- If an API asset failed, the report marks it explicitly and uses the reference crop as a temporary fallback.\n\n"
        "## Control Layer\n"
        "- Arrows, connector lines, and dashed loops are now recorded in reference_controls.json and selected_arrow_routes.json before PPT rendering.\n"
        "- reference_control_overlay.png visualizes the current reference-locked control layer for manual QA.\n\n"
        "## Background and Text Calibration\n"
        "- Each generated or reused AI asset is harmonized to the local slot background color recorded in slot_inventory.json.\n"
        f"- Editable PPT text uses TEXT_SIZE_SCALE={TEXT_SIZE_SCALE} as a first-pass visual calibration factor.\n\n"
        "## Internal Icon Fill\n"
        f"- Small icon assets target {INTERNAL_ACCEPTABLE_FILL_MIN:.0%}-{INTERNAL_ACCEPTABLE_FILL_MAX:.0%} internal effective-content fill.\n"
        "- Candidate selection is based on asset_internal_fill_report.json.\n"
        "- Generation canvas aspect ratio is selected from each locked asset slot; see asset_ratio_fit_report.json.\n"
        "- No local crop, automatic crop, PowerPoint crop, or PPT layout enlargement is used.\n",
        encoding="utf-8",
    )


def slot_by_id(slot_id: str) -> dict:
    lookup = {slot["id"]: slot for slot in SLOTS}
    if slot_id not in lookup:
        raise ValueError(f"Unknown trim slot: {slot_id}")
    return lookup[slot_id]


def selected_candidate_from_internal_report(slot_id: str) -> Path | None:
    report_path = OUT / "asset_internal_fill_report.json"
    if not report_path.exists():
        return None
    try:
        data = json.loads(report_path.read_text(encoding="utf-8"))
    except Exception:
        return None
    for item in data.get("assets", []):
        if item.get("slot_id") == slot_id and item.get("candidate_path"):
            path = Path(item["candidate_path"])
            if path.exists():
                return path
    return None


def source_icon_for_trim_experiment(slot: dict) -> Path:
    return selected_candidate_from_internal_report(slot["id"]) or latest_candidate_for_slot(slot["id"]) or (ASSET_DIR / f"{slot['id']}.png")


def trim_icon_candidate_for_experiment(slot: dict, source_path: Path, out_path: Path, padding_percent: float) -> dict:
    before = measure_internal_content_fill(source_path, slot["background_color_hex"])
    bbox = before.get("foreground_bbox_px")
    if not bbox:
        shutil.copyfile(source_path, out_path)
        after = measure_internal_content_fill(out_path, slot["background_color_hex"])
        return {
            "slot_id": slot["id"],
            "source_path": str(source_path),
            "output_path": str(out_path),
            "crop_applied": False,
            "trim_rejected": True,
            "risk_flags": ["invalid_foreground_bbox"],
            "before": before,
            "after": after,
            "fill_gain": round((after.get("foreground_bbox_fill_percent") or 0) - (before.get("foreground_bbox_fill_percent") or 0), 4),
            "recommendation": "reject",
        }

    image = Image.open(source_path).convert("RGB")
    width, height = image.size
    left, top, right, bottom = bbox
    bbox_w = right - left + 1
    bbox_h = bottom - top + 1
    pad = int(max(bbox_w, bbox_h) * padding_percent)
    crop_left = max(0, left - pad)
    crop_top = max(0, top - pad)
    crop_right = min(width - 1, right + pad)
    crop_bottom = min(height - 1, bottom + pad)
    crop_w = crop_right - crop_left + 1
    crop_h = crop_bottom - crop_top + 1
    risk_flags: list[str] = []
    if crop_left / width < 0.02 or crop_top / height < 0.02 or (width - 1 - crop_right) / width < 0.02 or (height - 1 - crop_bottom) / height < 0.02:
        risk_flags.append("edge_cutoff_risk")
    if crop_w <= 8 or crop_h <= 8:
        risk_flags.append("invalid_crop_box")

    if "invalid_crop_box" in risk_flags:
        shutil.copyfile(source_path, out_path)
        crop_applied = False
    else:
        crop = image.crop((crop_left, crop_top, crop_right + 1, crop_bottom + 1))
        scale = min(width / crop.width, height / crop.height)
        fit_w = max(1, int(crop.width * scale))
        fit_h = max(1, int(crop.height * scale))
        resized = crop.resize((fit_w, fit_h), Image.Resampling.LANCZOS)
        canvas = Image.new("RGB", (width, height), hex_to_tuple(slot["background_color_hex"]))
        canvas.paste(resized, ((width - fit_w) // 2, (height - fit_h) // 2))
        canvas.save(out_path)
        crop_applied = True

    after = measure_internal_content_fill(out_path, slot["background_color_hex"])
    after_fill = float(after.get("foreground_bbox_fill_percent") or 0.0)
    before_fill = float(before.get("foreground_bbox_fill_percent") or 0.0)
    if after_fill > INTERNAL_ACCEPTABLE_FILL_MAX:
        risk_flags.append("too_tight_risk")
    fill_gain = round(after_fill - before_fill, 4)
    if crop_applied and INTERNAL_ACCEPTABLE_FILL_MIN <= after_fill <= INTERNAL_ACCEPTABLE_FILL_MAX and fill_gain >= 0.10 and not risk_flags:
        recommendation = "accept"
    elif crop_applied and fill_gain > 0 and "invalid_crop_box" not in risk_flags:
        recommendation = "needs_visual_review"
    else:
        recommendation = "reject"
    return {
        "slot_id": slot["id"],
        "source_path": str(source_path),
        "output_path": str(out_path),
        "crop_applied": crop_applied,
        "trim_rejected": recommendation == "reject",
        "padding_percent": padding_percent,
        "crop_box_px": [crop_left, crop_top, crop_right, crop_bottom],
        "risk_flags": risk_flags,
        "before": before,
        "after": after,
        "before_fill": before_fill,
        "after_fill": after_fill,
        "fill_gain": fill_gain,
        "recommendation": recommendation,
    }


def create_trim_contact_sheet(rows: list[dict], output_path: Path) -> None:
    thumb = 220
    label_h = 54
    gap = 18
    width = thumb * 2 + gap * 3
    height = (thumb + label_h + gap) * max(len(rows), 1) + gap
    sheet = Image.new("RGB", (width, height), "#F3F1EC")
    draw = ImageDraw.Draw(sheet)
    for idx, row in enumerate(rows):
        y0 = gap + idx * (thumb + label_h + gap)
        before = Image.open(row["before_path"]).convert("RGB")
        after = Image.open(row["after_path"]).convert("RGB")
        before.thumbnail((thumb, thumb), Image.Resampling.LANCZOS)
        after.thumbnail((thumb, thumb), Image.Resampling.LANCZOS)
        x_before = gap + (thumb - before.width) // 2
        x_after = gap * 2 + thumb + (thumb - after.width) // 2
        sheet.paste(before, (x_before, y0 + (thumb - before.height) // 2))
        sheet.paste(after, (x_after, y0 + (thumb - after.height) // 2))
        label = f"{row['slot_id']} | {row['before_fill']:.3f} -> {row['after_fill']:.3f} | {row['recommendation']}"
        draw.text((gap, y0 + thumb + 8), label, fill="#333333")
        draw.text((gap, y0 - 2), "before", fill="#666666")
        draw.text((gap * 2 + thumb, y0 - 2), "after", fill="#666666")
    output_path.parent.mkdir(parents=True, exist_ok=True)
    sheet.save(output_path)


def run_icon_trim_experiment(trim_slots: str, trim_padding_percent: float, trim_output_dir: Path) -> None:
    trim_output_dir.mkdir(parents=True, exist_ok=True)
    before_dir = trim_output_dir / "before"
    after_dir = trim_output_dir / "after"
    before_dir.mkdir(parents=True, exist_ok=True)
    after_dir.mkdir(parents=True, exist_ok=True)
    rows = []
    report_items = []
    for slot_id in [item.strip() for item in trim_slots.split(",") if item.strip()]:
        slot = slot_by_id(slot_id)
        source_path = source_icon_for_trim_experiment(slot)
        if not source_path.exists():
            report_items.append({"slot_id": slot_id, "error": f"source image not found: {source_path}", "recommendation": "reject"})
            continue
        before_path = before_dir / f"{slot_id}.png"
        after_path = after_dir / f"{slot_id}.png"
        shutil.copyfile(source_path, before_path)
        item = trim_icon_candidate_for_experiment(slot, before_path, after_path, trim_padding_percent)
        report_items.append(item)
        rows.append({
            "slot_id": slot_id,
            "before_path": str(before_path),
            "after_path": str(after_path),
            "before_fill": item["before_fill"],
            "after_fill": item["after_fill"],
            "recommendation": item["recommendation"],
        })
    contact_sheet = trim_output_dir / "contact_sheet_before_after.png"
    create_trim_contact_sheet(rows, contact_sheet)
    report = {
        "summary": "Controlled icon trim experiment. Main assets and main PPT are not overwritten.",
        "trim_slots": [item.strip() for item in trim_slots.split(",") if item.strip()],
        "trim_padding_percent": trim_padding_percent,
        "output_dir": str(trim_output_dir),
        "main_assets_overwritten": False,
        "main_ppt_overwritten": False,
        "contact_sheet": str(contact_sheet),
        "items": report_items,
    }
    (trim_output_dir / "icon_trim_experiment_report.json").write_text(json.dumps(report, indent=2, ensure_ascii=False), encoding="utf-8")


def main() -> None:
    parser = argparse.ArgumentParser(description="Rebuild AutoFigure Architecture as an editable PPTX.")
    parser.add_argument("--icon-trim-experiment", action="store_true", help="Run a controlled icon trim A/B experiment without changing main assets or PPT.")
    parser.add_argument("--trim-slots", default="input_text_stack,ai_designer", help="Comma-separated slot ids for icon trim experiment.")
    parser.add_argument("--trim-padding-percent", type=float, default=0.08, help="Padding added around detected foreground bbox before experimental trim.")
    parser.add_argument("--trim-output-dir", default=str(OUT / "icon_trim_experiment"), help="Output directory for trim experiment artifacts.")
    parser.add_argument("--asset-cost-experiment", action="store_true", help="Dry-run the asset generation cost-saving policy without calling the image API.")
    parser.add_argument("--asset-cost-experiment-output", default=str(OUT / "asset_cost_experiment_report.json"), help="Output JSON path for the cost-saving dry-run report.")
    parser.add_argument("--asset-cost-baseline-retries", type=int, default=5, help="Baseline retries per failed slot used for cost comparison.")
    parser.add_argument("--asset-preflight-first-pass-experiment", action="store_true", help="Generate one planned candidate per selected slot in an isolated experiment directory.")
    parser.add_argument("--preflight-slots", default="all", help="Comma-separated slot ids, or 'all', for the first-pass preflight experiment.")
    parser.add_argument("--preflight-output-dir", default=str(OUT / "preflight_first_pass_experiment"), help="Output directory for first-pass preflight experiment artifacts.")
    parser.add_argument("--preflight-workers", type=int, default=4, help="Parallel workers for the first-pass preflight experiment.")
    parser.add_argument("--asset-retries", type=int, default=1, help="Maximum image-generation attempts per slot that is not accepted by the active asset policy.")
    parser.add_argument("--asset-workers", type=int, default=6, help="Parallel workers for main asset generation.")
    parser.add_argument("--strict-asset-regeneration", action="store_true", help="Use strict 80-95 percent fill acceptance instead of economy type-aware acceptance.")
    args = parser.parse_args()
    OUT.mkdir(parents=True, exist_ok=True)
    if args.icon_trim_experiment:
        run_icon_trim_experiment(args.trim_slots, args.trim_padding_percent, Path(args.trim_output_dir))
        return
    if args.asset_cost_experiment:
        run_asset_cost_experiment(args.asset_cost_baseline_retries, Path(args.asset_cost_experiment_output))
        return
    if args.asset_preflight_first_pass_experiment:
        run_preflight_first_pass_experiment(args.preflight_slots, Path(args.preflight_output_dir), args.preflight_workers)
        return
    crop_reference_slots()
    write_reference_overlays()
    asset_report = generate_assets_parallel(
        workers=args.asset_workers,
        retries=args.asset_retries,
        economy_mode=not args.strict_asset_regeneration,
    )
    write_metadata(asset_report)
    pptx_path = draw_ppt()
    export_preview(pptx_path)


if __name__ == "__main__":
    main()
